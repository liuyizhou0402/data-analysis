# -*- coding: utf-8 -*-
"""
STAT3888 — PowerPoint deck framework (16:9, .pptx).
Reusable helpers to build full 60-minute lectures with a consistent USYD theme.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))

# palette
INK    = RGBColor(0x1B,0x24,0x30)
MUTED  = RGBColor(0x5B,0x6A,0x7D)
ACCENT = RGBColor(0xE6,0x46,0x26)
BLUE   = RGBColor(0x2A,0x6F,0x97)
TEAL   = RGBColor(0x2A,0x9D,0x8F)
GOLD   = RGBColor(0xB9,0x84,0x16)
PANEL  = RGBColor(0xF4,0xF6,0xF9)
PANEL2 = RGBColor(0xEE,0xF3,0xF9)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
LINE   = RGBColor(0xDC,0xE3,0xEC)
DARK   = RGBColor(0x0F,0x17,0x20)
CODEBG = RGBColor(0x12,0x1A,0x24)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Calibri"
MONO = "Consolas"

class Deck:
    def __init__(self, number, title):
        self.number = number
        self.prs = Presentation()
        self.prs.slide_width = EMU_W
        self.prs.slide_height = EMU_H
        self.blank = self.prs.slide_layouts[6]

    # -------- low-level helpers
    def _slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _rect(self, sl, x,y,w,h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
        sp = sl.shapes.add_shape(shape, x,y,w,h)
        sp.shadow.inherit = False
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line; sp.line.width = Pt(line_w)
        return sp

    def _text(self, sl, x,y,w,h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              space_after=6, line_spacing=1.06, wrap=True):
        """runs: list of paragraphs; each paragraph is list of (text,size,color,bold,italic,font)."""
        tb = sl.shapes.add_textbox(x,y,w,h); tf = tb.text_frame
        tf.word_wrap = wrap; tf.vertical_anchor = anchor
        tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
        for i,para in enumerate(runs):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
            p.line_spacing = line_spacing
            if isinstance(para, tuple): para=[para]
            for seg in para:
                text,size,color,bold,italic = seg[0],seg[1],seg[2],seg[3],seg[4]
                fnt = seg[5] if len(seg)>5 else FONT
                r = p.add_run(); r.text = text
                r.font.size = Pt(size); r.font.color.rgb = color
                r.font.bold = bold; r.font.italic = italic; r.font.name = fnt
        return tb

    def _bg(self, sl, color=WHITE):
        self._rect(sl, 0,0, EMU_W, EMU_H, fill=color)

    def _accentbar(self, sl):
        self._rect(sl, 0,0, Inches(0.16), EMU_H, fill=ACCENT)

    def _footer(self, sl, idx):
        self._rect(sl, Inches(0.9), Inches(7.02), Inches(11.6), Pt(1), fill=LINE)
        self._text(sl, Inches(0.9), Inches(7.08), Inches(8), Inches(0.3),
                   [[("STAT3888 · Statistical Machine Learning", 10, MUTED, False, False)]])
        self._text(sl, Inches(10.9), Inches(7.08), Inches(1.55), Inches(0.3),
                   [[(f"L{self.number:02d} · {idx}", 10, MUTED, False, False)]], align=PP_ALIGN.RIGHT)

    def _kicker(self, sl, text, x=Inches(0.9), y=Inches(0.5)):
        self._text(sl, x, y, Inches(10), Inches(0.35),
                   [[(text.upper(), 12.5, ACCENT, True, False)]])

    def _h2(self, sl, title, x=Inches(0.9), y=Inches(0.82)):
        self._text(sl, x, y, Inches(11.6), Inches(0.9),
                   [[(title, 30, INK, True, False)]], line_spacing=1.0)

    def _content_top(self):
        return Inches(1.75)

    # -------- public slide types
    def title(self, title, subtitle, presenter="School of Mathematics & Statistics · The University of Sydney"):
        sl = self._slide(); self._bg(sl, INK)
        self._rect(sl, 0,0, EMU_W, EMU_H, fill=INK)
        # accent block
        self._rect(sl, 0, Inches(0), Inches(0.35), EMU_H, fill=ACCENT)
        self._rect(sl, Inches(0.9), Inches(1.05), Inches(2.4), Inches(0.5),
                   fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        self._text(sl, Inches(0.9), Inches(1.14), Inches(2.4), Inches(0.34),
                   [[(f"LECTURE {self.number:02d}", 13, WHITE, True, False)]], align=PP_ALIGN.CENTER)
        self._text(sl, Inches(0.9), Inches(2.15), Inches(11.4), Inches(2.4),
                   [[(title, 46, WHITE, True, False)]], line_spacing=1.02)
        self._text(sl, Inches(0.92), Inches(4.55), Inches(10.6), Inches(1.2),
                   [[(subtitle, 21, RGBColor(0xC8,0xD4,0xE2), False, False)]], line_spacing=1.15)
        self._rect(sl, Inches(0.9), Inches(6.35), Inches(11.5), Pt(1.2), fill=RGBColor(0x3a,0x46,0x55))
        self._text(sl, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.4),
                   [[(presenter, 13, RGBColor(0x9f,0xad,0xbd), False, False)]])
        return sl

    def section(self, label, title):
        sl = self._slide(); self._bg(sl, BLUE)
        self._rect(sl, 0,0, EMU_W, EMU_H, fill=BLUE)
        self._rect(sl, 0, Inches(3.0), EMU_W, Inches(0.05), fill=RGBColor(0x6f,0x9c,0xbc))
        self._text(sl, Inches(0.9), Inches(2.1), Inches(11), Inches(0.5),
                   [[(label.upper(), 15, RGBColor(0xBF,0xD8,0xE8), True, False)]])
        self._text(sl, Inches(0.9), Inches(3.2), Inches(11.5), Inches(2),
                   [[(title, 40, WHITE, True, False)]], line_spacing=1.02)
        return sl

    def bullets(self, kicker, title, items, idx, sub=None):
        """items: list of (text, level, color?) or strings. level 0/1."""
        sl = self._slide(); self._bg(sl); self._accentbar(sl)
        if kicker: self._kicker(sl, kicker)
        self._h2(sl, title)
        y = self._content_top()
        if sub:
            self._text(sl, Inches(0.9), Inches(1.55), Inches(11.4), Inches(0.5),
                       [[(sub, 17, MUTED, False, True)]]); y = Inches(2.15)
        tb = sl.shapes.add_textbox(Inches(0.95), y, Inches(11.5), Inches(4.9))
        tf = tb.text_frame; tf.word_wrap=True
        for i,it in enumerate(items):
            if isinstance(it,str): text,lvl,col = it,0,INK
            else: text,lvl = it[0],it[1]; col = it[2] if len(it)>2 else INK
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.line_spacing = 1.1
            p.space_after = Pt(11 if lvl==0 else 5); p.space_before=Pt(2)
            bullet = "▸  " if lvl==0 else "–  "
            r=p.add_run(); r.text=bullet; r.font.color.rgb = ACCENT if lvl==0 else MUTED
            r.font.size=Pt(20 if lvl==0 else 16); r.font.bold=True; r.font.name=FONT
            # allow simple **bold** markup
            for seg,b in _md(text):
                r=p.add_run(); r.text=seg
                r.font.size=Pt(20 if lvl==0 else 16.5); r.font.color.rgb=col
                r.font.bold=b; r.font.name=FONT
            p.level = lvl
        self._footer(sl, idx); return sl

    def two_col(self, kicker, title, left, right, idx, left_head=None, right_head=None,
                left_tone=BLUE, right_tone=ACCENT):
        sl=self._slide(); self._bg(sl); self._accentbar(sl)
        if kicker:self._kicker(sl,kicker)
        self._h2(sl,title)
        for x, head, body, tone in [(Inches(0.9),left_head,left,left_tone),
                                     (Inches(6.95),right_head,right,right_tone)]:
            self._rect(sl, x, Inches(1.85), Inches(5.45), Inches(4.7),
                       fill=PANEL, line=LINE, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            yy=Inches(2.1)
            if head:
                self._text(sl, x+Inches(0.3), yy, Inches(4.9), Inches(0.5),
                           [[(head,19,tone,True,False)]]); yy=Inches(2.75)
            tb=sl.shapes.add_textbox(x+Inches(0.3), yy, Inches(4.85), Inches(3.6)); tf=tb.text_frame; tf.word_wrap=True
            for i,it in enumerate(body):
                text = it if isinstance(it,str) else it[0]
                p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=1.1; p.space_after=Pt(8)
                r=p.add_run(); r.text="• "; r.font.color.rgb=tone; r.font.bold=True; r.font.size=Pt(16); r.font.name=FONT
                for seg,b in _md(text):
                    r=p.add_run(); r.text=seg; r.font.size=Pt(16); r.font.color.rgb=INK; r.font.bold=b; r.font.name=FONT
        self._footer(sl,idx); return sl

    def figure(self, kicker, title, img, idx, caption=None, bullets=None, bullet_head=None):
        sl=self._slide(); self._bg(sl); self._accentbar(sl)
        if kicker:self._kicker(sl,kicker)
        self._h2(sl,title)
        if bullets:
            self._img_fit(sl, img, Inches(0.9), Inches(1.8), Inches(7.5), Inches(4.7))
            if caption:
                self._text(sl, Inches(0.9), Inches(6.55), Inches(7.4), Inches(0.5),
                           [[(caption,12.5,MUTED,False,True)]])
            x=Inches(8.7)
            self._rect(sl, x, Inches(1.8), Inches(3.75), Inches(4.75),
                       fill=PANEL2, line=LINE, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            yy=Inches(2.05)
            if bullet_head:
                self._text(sl, x+Inches(0.28), yy, Inches(3.3), Inches(0.5),
                           [[(bullet_head,17,BLUE,True,False)]]); yy=Inches(2.65)
            tb=sl.shapes.add_textbox(x+Inches(0.28), yy, Inches(3.25), Inches(3.7)); tf=tb.text_frame; tf.word_wrap=True
            for i,it in enumerate(bullets):
                text = it if isinstance(it,str) else it[0]
                p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=1.08; p.space_after=Pt(8)
                r=p.add_run(); r.text="• "; r.font.color.rgb=BLUE; r.font.bold=True; r.font.size=Pt(15); r.font.name=FONT
                for seg,b in _md(text):
                    r=p.add_run(); r.text=seg; r.font.size=Pt(15); r.font.color.rgb=INK; r.font.bold=b; r.font.name=FONT
        else:
            self._img_fit(sl, img, Inches(1.7), Inches(1.75), Inches(10.0), Inches(4.7), center=True)
            if caption:
                self._text(sl, Inches(1.5), Inches(6.55), Inches(10.3), Inches(0.5),
                           [[(caption,13,MUTED,False,True)]], align=PP_ALIGN.CENTER)
        self._footer(sl,idx); return sl

    def formula(self, kicker, title, formula_text, explain, idx, note=None):
        sl=self._slide(); self._bg(sl); self._accentbar(sl)
        if kicker:self._kicker(sl,kicker)
        self._h2(sl,title)
        self._rect(sl, Inches(1.6), Inches(2.15), Inches(10.1), Inches(1.55),
                   fill=DARK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        self._text(sl, Inches(1.7), Inches(2.15), Inches(9.9), Inches(1.55),
                   [[(formula_text, 30, RGBColor(0xEA,0xF2,0xFB), False, True, "Cambria Math")]],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tb=sl.shapes.add_textbox(Inches(1.6), Inches(4.05), Inches(10.1), Inches(2.4)); tf=tb.text_frame; tf.word_wrap=True
        for i,it in enumerate(explain):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=1.12; p.space_after=Pt(9)
            r=p.add_run(); r.text="• "; r.font.color.rgb=ACCENT; r.font.bold=True; r.font.size=Pt(18); r.font.name=FONT
            for seg,b in _md(it):
                r=p.add_run(); r.text=seg; r.font.size=Pt(18); r.font.color.rgb=INK; r.font.bold=b; r.font.name=FONT
        if note:
            self._callout(sl, note, Inches(1.6), Inches(6.0), Inches(10.1))
        self._footer(sl,idx); return sl

    def code(self, kicker, title, code_lines, idx, lang="r", caption=None):
        sl=self._slide(); self._bg(sl); self._accentbar(sl)
        if kicker:self._kicker(sl,kicker)
        self._h2(sl,title)
        h = Inches(0.34)*len(code_lines)+Inches(0.5)
        h = min(h, Inches(4.6))
        self._rect(sl, Inches(0.9), Inches(1.85), Inches(11.5), h,
                   fill=CODEBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb=sl.shapes.add_textbox(Inches(1.15), Inches(2.02), Inches(11.0), h-Inches(0.3)); tf=tb.text_frame; tf.word_wrap=False
        for i,ln in enumerate(code_lines):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=1.15; p.space_after=Pt(1)
            # comments in muted green
            if ln.strip().startswith("#"):
                r=p.add_run(); r.text=ln; r.font.color.rgb=RGBColor(0x7c,0x93,0x8a)
            else:
                r=p.add_run(); r.text=ln if ln else " "; r.font.color.rgb=RGBColor(0xE6,0xED,0xF3)
            r.font.size=Pt(15.5); r.font.name=MONO
        if caption:
            self._text(sl, Inches(0.9), Inches(1.9)+h, Inches(11.4), Inches(0.6),
                       [[(caption,14,MUTED,False,True)]])
        self._footer(sl,idx); return sl

    def table(self, kicker, title, headers, rows, idx, note=None, col_widths=None):
        sl=self._slide(); self._bg(sl); self._accentbar(sl)
        if kicker:self._kicker(sl,kicker)
        self._h2(sl,title)
        nrows=len(rows)+1; ncols=len(headers)
        tw=Inches(11.5); th=min(Inches(0.55)*nrows, Inches(4.6))
        gtb=sl.shapes.add_table(nrows,ncols,Inches(0.9),Inches(1.85),tw,th).table
        if col_widths:
            for j,w in enumerate(col_widths): gtb.columns[j].width=Inches(w)
        for j,htext in enumerate(headers):
            c=gtb.cell(0,j); c.text=""; p=c.text_frames[0].paragraphs[0] if False else c.text_frame.paragraphs[0]
            r=p.add_run(); r.text=htext; r.font.bold=True; r.font.size=Pt(15); r.font.color.rgb=WHITE; r.font.name=FONT
            c.fill.solid(); c.fill.fore_color.rgb=BLUE; c.vertical_anchor=MSO_ANCHOR.MIDDLE
        for i,row in enumerate(rows,1):
            for j,val in enumerate(row):
                c=gtb.cell(i,j); c.text=""; p=c.text_frame.paragraphs[0]
                for seg,b in _md(val):
                    r=p.add_run(); r.text=seg; r.font.size=Pt(13.5); r.font.color.rgb=INK; r.font.bold=b; r.font.name=FONT
                c.fill.solid(); c.fill.fore_color.rgb = WHITE if i%2 else PANEL
                c.vertical_anchor=MSO_ANCHOR.MIDDLE
        # strip default table style banding borders is fine
        if note:
            self._callout(sl, note, Inches(0.9), Inches(1.95)+th, Inches(11.5))
        self._footer(sl,idx); return sl

    def exercise(self, title, prompt_items, idx, hint=None, ordered=True):
        sl=self._slide(); self._bg(sl, RGBColor(0xFB,0xF3,0xF0)); self._accentbar(sl)
        self._kicker(sl,"In-class exercise")
        self._h2(sl,title)
        self._rect(sl, Inches(0.9), Inches(1.8), Inches(11.5), Inches(3.7) if hint else Inches(4.5),
                   fill=WHITE, line=LINE, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb=sl.shapes.add_textbox(Inches(1.25), Inches(2.05), Inches(10.9), Inches(3.4)); tf=tb.text_frame; tf.word_wrap=True
        for i,it in enumerate(prompt_items):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=1.15; p.space_after=Pt(10)
            lead = f"{i+1}.  " if ordered else "•  "
            r=p.add_run(); r.text=lead; r.font.color.rgb=ACCENT; r.font.bold=True; r.font.size=Pt(19); r.font.name=FONT
            for seg,b in _md(it):
                r=p.add_run(); r.text=seg; r.font.size=Pt(19); r.font.color.rgb=INK; r.font.bold=b; r.font.name=FONT
        if hint:
            self._rect(sl, Inches(0.9), Inches(5.75), Inches(11.5), Inches(1.0),
                       fill=RGBColor(0xFF,0xF6,0xEE), line=GOLD, line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            self._text(sl, Inches(1.2), Inches(5.95), Inches(10.9), Inches(0.8),
                       [[("Hint.  ",16,GOLD,True,False),(hint,16,INK,False,False)]], line_spacing=1.1)
        self._footer(sl,idx); return sl

    def summary(self, points, idx, nextup=None):
        sl=self._slide(); self._bg(sl); self._accentbar(sl)
        self._kicker(sl,"Summary")
        self._h2(sl,"Key takeaways")
        tb=sl.shapes.add_textbox(Inches(0.95), Inches(1.9), Inches(11.4), Inches(4.2)); tf=tb.text_frame; tf.word_wrap=True
        for i,pt in enumerate(points):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=1.12; p.space_after=Pt(13)
            r=p.add_run(); r.text="✓  "; r.font.color.rgb=TEAL; r.font.bold=True; r.font.size=Pt(21); r.font.name=FONT
            for seg,b in _md(pt):
                r=p.add_run(); r.text=seg; r.font.size=Pt(20); r.font.color.rgb=INK; r.font.bold=b; r.font.name=FONT
        if nextup:
            self._rect(sl, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.85),
                       fill=PANEL2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            self._text(sl, Inches(1.2), Inches(6.2), Inches(10.9), Inches(0.5),
                       [[("Next lecture →  ",16,ACCENT,True,False),(nextup,16,INK,False,False)]])
        self._footer(sl,idx); return sl

    # -------- utilities
    def _callout(self, sl, text, x, y, w, tone="tip"):
        col = TEAL if tone=="tip" else (ACCENT if tone=="warn" else BLUE)
        bg  = RGBColor(0xEE,0xF6,0xF4) if tone=="tip" else (RGBColor(0xFD,0xF0,0xEC) if tone=="warn" else PANEL2)
        self._rect(sl, x, y, w, Inches(0.85), fill=bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        self._rect(sl, x, y, Inches(0.08), Inches(0.85), fill=col)
        tb=sl.shapes.add_textbox(x+Inches(0.28), y+Inches(0.08), w-Inches(0.5), Inches(0.7)); tf=tb.text_frame
        tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        p=tf.paragraphs[0]; p.line_spacing=1.08
        for seg,b in _md(text):
            r=p.add_run(); r.text=seg; r.font.size=Pt(15); r.font.color.rgb=INK; r.font.bold=b; r.font.name=FONT

    def _img_fit(self, sl, path, x, y, maxw, maxh, center=False):
        from PIL import Image
        try:
            iw,ih = Image.open(path).size
        except Exception:
            iw,ih = 1000,700
        ar = iw/ih; boxar = maxw/maxh
        if ar > boxar:
            w = maxw; h = int(maxw/ar)
        else:
            h = maxh; w = int(maxh*ar)
        px = x + (maxw-w)//2 if center else x
        py = y + (maxh-h)//2
        sl.shapes.add_picture(path, px, py, width=w, height=h)

    def save(self, path):
        self.prs.save(path); return path


def _md(text):
    """Tiny **bold** parser -> list of (segment, is_bold)."""
    out=[]; parts=text.split("**");
    for i,seg in enumerate(parts):
        if seg=="": continue
        out.append((seg, i%2==1))
    return out or [(text,False)]
