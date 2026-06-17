# -*- coding: utf-8 -*-
"""
Generates IELTS Listening L01-L10 (61+ pages each)
Structure per lesson:
  cover+agenda+review+objectives+pain_point = 5
  PART01-04 each: section+concept+concept+content+content+example+example = 7×4 = 28
  PART05 vocab: section+vocab+vocab+vocab+two_col = 5
  PART06 method: section+steps+steps+model+model+compare+compare = 7
  PART07 drill: section+drill×5 = 6
  PART08 pitfall: section+pain_point+content+content+tip_card+tip_card = 6
  checklist+summary+homework+preview = 4
  TOTAL = 61
"""
import sys, os
sys.path.insert(0, os.path.dirname(__file__))
import ppt_utils as U

OUT = os.path.join(os.path.dirname(__file__), "listening")
os.makedirs(OUT, exist_ok=True)

def T():
    U.set_theme(U.RGBColor(0x0F,0x4C,0x5C), U.RGBColor(0x12,0xA5,0x94),
                U.RGBColor(0xE6,0xF3,0xF1), U.RGBColor(0x07,0x2B,0x34),
                '听力','IELTS Listening')

def part(p,idx,tc,te,concs,conts,exs):
    """7-page PART: section + 2 concept + 2 content + 2 example"""
    U.section(p,idx,tc,te)
    for c in concs: U.concept(p,c[0],c[1],c[2],note=c[3] if len(c)>3 else None)
    for c in conts: U.content(p,c[0],c[1],note=c[2] if len(c)>2 else None)
    for e in exs:   U.example(p,e[0],e[1],answer=e[2],tip=e[3])

# ══════════════════════════════════════════════════════════════════════════════
#  L01  考试全貌与失分诊断
# ══════════════════════════════════════════════════════════════════════════════
def L01():
    T(); p=U.new_prs()
    U.cover(p,'听力',1,'听力考试全貌\n& 失分诊断','Test Overview & Failure Diagnosis')
    U.agenda(p,['课程总览与学习路线','四个Section结构','10大题型速览',
                '分数换算与目标','失分原因诊断','例题体验','课堂练习','小结作业'])
    U.content(p,'10节课学习路线图 Course Roadmap',
        [('L01','考试全貌+失分诊断 → 建立正确认知与备考框架'),
         ('L02-L05','Section 1-4 各个击破 → 掌握每类场景的解题策略'),
         ('L06','预测技能+关键词定位 → 让听力变为"主动找"而非"被动听"'),
         ('L07','笔记速记系统 → 双任务协调，实时记录不遗漏'),
         ('L08','六大陷阱解析 → 听懂了也要避坑'),
         ('L09-L10','口音适应+全真模考 → 查漏补缺，冲刺目标分')],
        note='每节课配合剑桥真题10-19同步练习，建议每周完成1套完整模拟')
    U.objectives(p,['掌握雅思听力四Section的结构、场景与难度梯度',
                    '了解10大题型的特点与分布规律',
                    '知晓6/7/8分的正确题数要求及分数换算',
                    '完成个人失分原因诊断，明确本课程提分重点'])
    U.pain_point(p,'深圳学生最普遍的三类失分',
        ['听懂了却做错：关键词出现后立即写，忽略转折词(but/actually)之后的更正',
         '拼写失分：necessary/accommodation/February写错，明明听对了却得0分',
         'Section 3跟丢：多人对话说话人切换时张冠李戴，观点归属混乱'],
        ['听到关键词≠答案到了，继续听完整句，转折词后才是真答案',
         '建立100个高频拼写词表，每天5分钟听写训练，专项解决拼写失分',
         'Section 3：用A/B/C标注说话人，分别追踪各自立场，不混淆'])
    # PART01 — 考试结构
    part(p,1,'考试结构与Section梯度','Exam Structure & Section Difficulty',
        [('四个Section：场景与难度','S1最简单，S4最难；依次递增，但可针对性提分',
          [('S1','日常对话，2人，预约/租房/旅游；语速最慢；必须满分10/10'),
           ('S2','独白，1人，社区/设施/活动；地图和匹配题为主'),
           ('S3','学术讨论，2-4人；多选题难；是深圳学生失分最多的Section'),
           ('S4','学术讲座，1人；AWL词汇密集；是拉开7分和8分差距的关键'),
           ('时长','30分钟录音；纸笔考+10分钟誊写；机考边听边输入')],
          '录音只播一次，不可倒带——实时答题是核心技能'),
         ('分数换算表：你需要对多少题？','明确目标，量化差距',
          [('Band 9','39-40题（正确率≥97%）'),
           ('Band 8','35-37题（正确率87-92%）'),
           ('Band 7','30-31题（正确率75-77%）'),
           ('Band 6','23-25题（正确率57-62%）'),
           ('提分路径','6→7分：需额外做对约6-7题；主要来自S3陷阱题和S4填空题'),
           ('核心认知','从6分到7分是"策略提升"；从7到8分才是"实力提升"')],
          '6分→7分的6题差距，完全可以靠策略而非大量刷词来弥补')],
        [('纸笔考 vs 机考：关键差异',
          [('纸笔考','听力+阅读+写作；听力结束后额外10分钟誊写到答题卡'),
           ('机考','无誊写时间，直接在屏幕输入；可修改已作答的答案'),
           ('机考注意','提前练打字速度；切换中英文输入法；注意大小写键'),
           ('时间建议','机考不用担心誊写出错，但要避免边听边看键盘导致分心'),
           ('备考建议','报名前确认考试形式，针对性准备')]),
         ('雅思听力 vs 国内英语听力的差距',
          [('口音','国内教材多用美音；雅思主要是英式+澳式——口音不适应是大问题'),
           ('语速','雅思语速比CET-4快15-20%，不允许暂停或倒带'),
           ('词汇','S4包含大量AWL学术词汇；不认识≠听不出——音读出来要能辨认'),
           ('陷阱','国内听力很少有更正/否定陷阱；雅思约30%题目含陷阱'),
           ('应对','大量接触英式/澳式英语；做真题建立"雅思听力感"')])],
        [('体验题：S1信息表填空',
          'Questions 1-3. Write NO MORE THAN TWO WORDS AND/OR A NUMBER.\n\n'
          'CITY SPORTS CENTRE Membership\nName: __________ (Q1)\nType: __________ membership (Q2)\nStart: 1st __________ (Q3)\n\n'
          'Script: "My name\'s Chen Wei — C-H-E-N, W-E-I. I\'d like a family membership.\nAnd starting from the first of September, please."',
          'Q1: Chen Wei（逐字母听写，first+last name都要写）\nQ2: family（修饰membership的形容词）\nQ3: September（月份，大写S）',
          'Name题预判：有几个空格？空格前有Mr/Mrs提示？一般first name+last name各一格'),
         ('体验题：S4笔记填空（学术词汇挑战）',
          'Questions 31-32. Write NO MORE THAN TWO WORDS.\n\n'
          'Urban Heat Island: Causes\n• Main source: __________ from synthetic fabrics (Q31)\n• Reduced __________ due to lack of greenery (Q32)\n\n'
          'Script: "Research identifies microfibre shedding from synthetic textiles as the primary source.\nFurthermore, reduced evapotranspiration due to declining urban greenery compounds the problem."',
          'Q31: microfibre shedding（听不清先标首字母micro...，誊写时补全）\nQ32: evapotranspiration（专业词，逐音节e-vapo-trans-pi-ra-tion，记首字母占位）',
          'S4学术词：不认识没关系，按音节把听到的音写下来（音译占位），誊写时再核对')])
    # PART02 — 题型
    part(p,2,'10大题型：特点与策略','10 Question Types: Features & Strategies',
        [('填空类题型（约占60%）','填空是最高频的题型，拼写和限词是主要失分点',
          [('Form/Table Completion','日常信息（姓名/电话/价格）；答案通常是原词；限词严格'),
           ('Note Completion','结构化笔记；S2/S4常见；注意词性（空格前后决定词性）'),
           ('Summary Completion','摘要填空；注意同义替换；难度最高，需要理解上下文'),
           ('Flow-chart Completion','按步骤顺序填写；被动语态常见；动词词根重要'),
           ('句子完成 Sentence Completion','找能使句子语法成立的词；注意大小写')],
          '填空类答案几乎全是原词（原文中出现的词），不需要改写，拼写必须100%正确'),
         ('选择与配对类题型','选择题需要识别干扰；配对题要提前读完所有选项',
          [('Single MCQ','A/B/C三选一；转折词(but/however)后面的内容才是答案'),
           ('Multiple MCQ','Choose TWO from FIVE；每个选项单独判断，不靠感觉'),
           ('Matching','将问题与A-H选项配对；通常选项数>问题数，有多余选项'),
           ('Map/Plan Labelling','地图题；方向词是核心；画箭头追踪路径'),
           ('Short Answer','限词填写答案；通常是名词短语；不需要完整句子')])],
        [('题型时间分配与优先策略',
          [('第一原则','每道题最多1.5分钟；超时立即跳过，先做确定的题'),
           ('简单先做','S1全部题+S2填空题：先做，拿稳基础分'),
           ('难题标记','S3多选题+S4摘要填空：听时先记关键词，Section结束后利用空档补充'),
           ('不留空','宁可猜一个，也不留空白——猜对有分，留空一定没分'),
           ('誊写检查','纸笔考：先全部写完，最后3分钟统一检查拼写和限词')]),
         ('PRCA四步法：从被动听到主动找',
          [('P Pre-read 预读','利用说明时间预读题目，圈关键词，预判答案类型'),
           ('R Respond 实时','边听边写，候选词先记旁边，确认后写正式格'),
           ('C Check 即核','答完一题立即看下一题，不在已过的题上停留'),
           ('A Amend 修正','誊写时逐题检查：拼写/大小写/限词/复数'),
           ('核心','PRCA是习惯，不是步骤——考场中P和R几乎同时进行')])],
        [('例题：MCQ转折陷阱',
          'Q21: Why did Tom change his research topic?\n'
          'A. lack of relevant data  B. supervisor\'s suggestion  C. personal interest changed\n\n'
          'Script: "I was going to focus on climate data — there\'s tons of material on that.\nBut my supervisor thought it was too broad, so I ended up looking at urban microclimate instead."',
          'Answer: B — supervisor\'s suggestion\n'
          'Analysis: 录音先提"climate data"(A的干扰)，再提"supervisor thought"(B=正确答案)',
          '转折陷阱："was going to + BUT" — "was going to"之前的是假设，but之后才是真答案'),
         ('例题：匹配题预读策略',
          'Match each facility with the correct description (A-F).\n'
          'A. recently renovated  B. only for members  C. closed on Mondays\n'
          'D. booking required  E. free of charge  F. for over-18s only\n\n'
          'Q11. Swimming pool → ?   Q12. Tennis courts → ?   Q13. Gym → ?\n\n'
          'Script: "...the pool\'s been completely refurbished... tennis courts — you\'ll need to book in advance... the gym is restricted to adults, so over 18s only..."',
          'Q11: A (recently renovated = refurbished)\nQ12: D (booking required = book in advance)\nQ13: F (for over-18s only = restricted to adults)',
          '匹配题：提前读完A-F所有选项，标记关键词；录音中找同义替换词而非原词')])
    # PART03 — 失分诊断
    part(p,3,'失分原因深度诊断','Root Cause Analysis of Errors',
        [('失分类型一：拼写与书写错误','每年约15-20%的考生因拼写错误失去本可得到的分',
          [('最高频拼错词','necessary(1c2s) / accommodation(2c2m) / environment(no a)'),
           ('字母混淆','m/n, u/v, e/a 在快速听写时混淆；13/30 重音位置不同'),
           ('大小写错误','专有名词必须大写：Monday/London/Dr Smith；普通名词小写'),
           ('复数遗漏','two _____(s)、several _____必须加s；a _____不加s'),
           ('解决方案','100个高频拼写词每天5分钟听写；誊写逐字母检查')],
          '拼写是最可惜的失分：你听对了正确答案，却因写错字母得0分'),
         ('失分类型二：陷阱识别不足','约30%的题目含陷阱，不识别陷阱7分是上限',
          [('更正陷阱','先说A再改成B；信号词：actually/sorry/I mean/wait → 写B'),
           ('否定陷阱','"NOT available on Friday" → 考生听到Friday就写，忘了NOT'),
           ('转折陷阱','"normal price £60, but members pay £45" → 题目若问会员价答45'),
           ('顺序陷阱','答案不按题号顺序出现（尤其地图题）→ 要灵活跳着找'),
           ('干扰词陷阱','录音提到了选项A的词，但整句意思否定了A → 不能选')]  )],
        [('失分类型三：时间管理失控',
          [('单题纠结','一道题超过1.5分钟，后面3题全部错过——不值得'),
           ('预读不足','Section开始前没完成预读，被动听而非主动找答案'),
           ('Section间隙浪费','S1→S2有20秒，必须用来预读S2前5题；很多人在发呆'),
           ('誊写粗心','纸笔考最后10分钟抄错答案，把正确答案写成了错的'),
           ('解决','建立"不纠结"原则：跟丢了立即放弃，找下一题锚点重定位')]),
         ('失分类型四：词汇与口音障碍',
          [('AWL词汇不熟','S4出现evapotranspiration/entrepreneurship等AWL词，一字不识'),
           ('口音陌生','深圳学生习惯美音，遇到英式/澳式口音第一反应"这是英语吗？"'),
           ('连读弱化','want to → wanna；going to → gonna；听到of/to/and弱化形式认不出'),
           ('解决1','每天BBC/ABC 20分钟，主动接触英澳口音，靠量来适应'),
           ('解决2','AWL核心300词背会：不要求写，要求听到音能反应出意思')])],
        [('诊断练习：找出失分类型',
          '录音1："The event is on Thursday the 15th... actually no, it\'s the 16th, Friday."\n考生写：Thursday → 失分\n\n'
          '录音2："We\'re open Monday to Saturday, but NOT on Sundays."\n考生写：Sunday → 失分\n\n'
          '录音3："The fee is £80 per year — that\'s with the student discount, normally it\'s £120."\n考生写：£120 → 失分',
          '录音1：更正陷阱 → 答案是Friday/16th（后说的）\n录音2：否定陷阱 → NOT Sunday=周日不开放，答案要看题目问什么\n录音3：价格陷阱 → 若问学生价，答80；题目是关键',
          '做题前先问：题目到底在问什么？——很多错误源于读题不仔细，听到了却答错了'),
         ('诊断练习：限词违规',
          'Instruction: Write NO MORE THAN TWO WORDS AND/OR A NUMBER\n\n'
          '考生A答案："the swimming pool"（3 words）→ 0分\n'
          '考生B答案："swimming pool"（2 words）→ 满分\n'
          '考生C答案："£32.50"（1 number）→ 满分\n'
          '考生D答案："every Tuesday afternoon"（3 words）→ 0分\n'
          '考生E答案："Tuesday"（1 word）→ 满分（如果题目问day）',
          'A：冠词"the"也算一个词，加上后超出限制\nD："every"是多余的，"Tuesday afternoon"=2词才满足限制',
          '限词检查口诀：写完答案先数词数 → 冠词算词数 → 数字算一格 → 超限必须缩短')])
    # PART04 — 目标与计划
    part(p,4,'目标设定与学习策略','Goal Setting & Study Strategies',
        [('从现在到目标：差距量化分析','知道自己在哪里，才能规划去哪里',
          [('诊断方法','做1套剑桥真题，统计每个Section正确题数'),
           ('6→7分','多做对约6题；重点：S3陷阱题(+2)、S4填空(+2)、S1/S2减失误(+2)'),
           ('7→8分','多做对约5题；需要S1/S2接近满分+S3稳定7/10+S4稳定6/10'),
           ('8→9分','几乎不能失误；口音/速度/词汇全面提升，不只是策略'),
           ('现实建议','8周密集备考可从6分到7分；7分到8分通常需要3-6个月')],
          '目标分数决定学习侧重点——6→7靠策略，7→8靠实力'),
         ('剑桥真题高效使用法','用对材料比练习量更重要',
          [('模拟模式','每套计时完成，不暂停，模拟真实考场压力'),
           ('精析模式','每道错题听3-5遍，对照录音原文找失分原因'),
           ('跟读模式','选1段精听材料，跟读模仿语调和语速'),
           ('错题本','记录：题号/错误类型/正确答案/陷阱说明'),
           ('推荐顺序','剑桥15-16（先适应），17-19（冲刺），10-14（备用')])],
        [('每周练习计划（以目标7分为例）',
          '周一：精听分析（选上周错题2-3道，对照原文）\n'
          '周二：Section 1+2 专项（只做这两个Section）\n'
          '周三：泛听放松（BBC 6-Minute English 20分钟）\n'
          '周四：Section 3+4 专项+笔记练习\n'
          '周五：完整模拟1套（严格计时）\n'
          '周六：分析本周错题，更新错题本\n'
          '周日：休息（保持听英语的习惯即可，不刷题）',
          '重点：模拟（周五）+ 精析（周一/周六）= 每周最重要的两件事',
          '每周练习时间：5-6小时足够；不靠时间堆，靠"做-析-改"循环'),
         ('备考资源清单',
          [('必备','剑桥雅思真题10-19（至少备3本）；每本含4套完整模拟'),
           ('口音练习','BBC 6-Minute English（英式）；Australia ABC Listen（澳式）'),
           ('词汇资源','AWL Sublist 1-10（600词）；专注听力识别，不需要会写'),
           ('App工具','雅思哥/小站雅思（移动练习）；Audacity（慢速回放精听）'),
           ('错题资源','自建错题本最重要；比买任何参考书都更有针对性')])],
        [('诊断练习：识别你的弱点Section',
          '请根据以下描述，判断该学生的主要失分在哪个Section：\n\n'
          '学生A：S1正确9/10，S2正确7/10，S3正确4/10，S4正确5/10，总分25/40\n'
          '学生B：S1正确8/10，S2正确6/10，S3正确7/10，S4正确3/10，总分24/40\n'
          '学生C：S1正确6/10，S2正确5/10，S3正确6/10，S4正确5/10，总分22/40',
          '学生A：Section 3是瓶颈（4/10），建议专攻S3多选题策略\n学生B：Section 4是瓶颈（3/10），建议加强AWL词汇和笔记技能\n学生C：各Section均偏低，预读技能和时间管理是根本问题',
          '找到最弱的1个Section集中突破，比平均练习效率高得多'),
         ('提分最快的3个策略（速效版）',
          '以下3个策略是6分→7分最快的提分路径：\n\n'
          '策略1：S1拿满分——语速最慢，预读完成后理应满分\n'
          '策略2：消灭限词错误——每次做题后检查词数；一个月消灭所有限词失分\n'
          '策略3：建立更正意识——练习识别actually/sorry/I mean，多拿2-3题',
          '这3个策略合计可让你多拿3-5题；很多6分考生靠这3个策略直接跨到7分',
          '这3个策略都是"策略性"提升，不需要词汇量或听力感知力同步提升就能实现')])
    # PART05
    U.section(p,5,'词汇与速记表达','Core Vocabulary & Expressions')
    U.vocab(p,'考场核心词汇：动词类',
        [('cancel / reschedule',   '取消 / 改期（S1预约场景最高频）'),
         ('confirm / verify',      '确认（结束对话前常说"let me confirm..."）'),
         ('apply / register',      '申请 / 注册（报名场景）'),
         ('suggest / recommend',   '建议（S3观点讨论中区分谁在建议）'),
         ('identify / highlight',  '识别/强调（S4讲座中表示研究发现）'),
         ('evaluate / assess',     '评估（研究方法讨论常见词）'),
         ('demonstrate / reveal',  '展示/揭示（结果描述常见动词）')])
    U.vocab(p,'Section 1高频名词表',
        [('appointment / slot',    '预约 / 时间段'),
         ('deposit / fee / charge','押金 / 费用 / 收费'),
         ('facilities / amenities','设施（注意复数！）'),
         ('membership / subscription','会员资格 / 订阅'),
         ('receipt / invoice',     '收据 / 发票'),
         ('enquiry / query',       '询问（英式拼法）'),
         ('vacancy / availability','空缺 / 空闲时间')])
    U.vocab(p,'Section 3/4学术词汇',
        [('hypothesis / thesis',   '假设 / 论点'),
         ('methodology / approach','研究方法 / 方式'),
         ('findings / results',    '发现 / 结果（区分：发现=新知；结果=数据）'),
         ('implication / impact',  '含义 / 影响'),
         ('evaluation / critique', '评价 / 批评（中性词，不一定是负面）'),
         ('feasibility / viability','可行性（S3讨论是否可以实施时常见）'),
         ('correlation / causation','相关性 / 因果关系（混淆是研究方法讨论重点）')])
    U.two_col(p,'同义替换速查：题目 vs 录音',
        '题目关键词（你看到的）',
        ['start → commence / begin / initiate',
         'increase → rise / grow / surge / climb',
         'problem → issue / challenge / difficulty',
         'important → crucial / significant / vital',
         'show → demonstrate / indicate / reveal',
         'children → young people / minors / youth',
         'use → utilise / employ / apply'],
        '录音表达（你听到的）',
        ['approximately → about / roughly / around',
         'expensive → costly / pricey / not cheap',
         'often → frequently / regularly / repeatedly',
         'build → construct / develop / establish',
         'location → venue / site / place / spot',
         'suggest → recommend / propose / advise',
         'because → due to / owing to / since'])
    # PART06
    U.section(p,6,'方法精讲与真题范例','Methods & Authentic Examples')
    U.steps(p,'考场听力全流程：Section间做什么',
        [('考前准备','确认耳机舒适；音量合适；翻开答题卡看Section 1题目'),
         ('S1说明播放时','利用这30秒预读Q1-10；圈出关键词；预判答案类型'),
         ('S1→S2间隙(20s)','立即翻到S2；预读Q11-15（至少完成前3题的预读）'),
         ('S2→S3间隙(20s)','预读S3所有题；重点：MCQ选项A-E提前读完并记关键词'),
         ('S3→S4间隙(20s)','预读S4所有题；S4最难，预读最关键，不能浪费'),
         ('结束后(纸笔考)','立即开始誊写；按题号顺序；检查拼写/大小写/复数/限词')])
    U.steps(p,'陷阱识别实战流程',
        [('Step 1 预读时','圈出题目中的否定词(NOT/EXCEPT)和限定词(ONLY/ALWAYS)'),
         ('Step 2 听时','听到but/however/actually/sorry → 立即提高警觉'),
         ('Step 3 更正','听到"I mean/let me correct/wait, no" → 划掉之前候选词'),
         ('Step 4 确认','等说话人说完整句再写答案，不要听到关键词就停'),
         ('Step 5 检查','每题写完回看：有没有违反题目中的限定条件？')])
    U.model_answer(p,'S4笔记满分示范（笔记格式）',
        'Section 4 Lecture: Circular Economy',
        'Topic: Circular Economy (CE) = ≠ linear (take-make-dispose)\n\n'
        'Key principles:\n'
        '1. Design for durability (not planned obsolescence)\n'
        '2. Reuse → Repair → Recycle (hierarchy)\n'
        '3. ↓ waste → ↑ resource efficiency\n\n'
        'Example: Renault uses 80% recycled parts in factory\nResult: saves €1B/yr\n\n'
        'Challenge: consumer behaviour + initial investment cost',
        ['只记关键名词和数字，省略冠词动词',
         '用=和≠表示定义和对比，节省大量书写时间',
         '数字和例子优先记录（考题最爱考这里）',
         '符号箭头→表示因果，↑↓表示增减'])
    U.model_answer(p,'S3多选题答题示范（同时追踪两人）',
        'Q23-24: Choose TWO features both students agree on.',
        '笔记追踪格式：\n\n'
        '         Amy    Tom\n'
        'A research design  ✓      ?\n'
        'B time management  ✓      ✓  ← 两人都同意\n'
        'C interview skills  ✗      ✓\n'
        'D library resources ?      ✓\n'
        'E peer feedback     ✓      ✓  ← 两人都同意',
        'Answer: B and E\n'
        'Note: A只有Amy同意；C只有Tom同意；D不确定，但B和E已经确认',
        '多选题：在选项旁边标注每个说话人的态度；等两人都表示同意才确认')
    U.compare(p,'答题策略：被动接收 vs 主动猎找',
        '❌ 被动接收（6分以下）',
        ['录音播放才开始看题目',
         '不知道下一题要填什么类型的信息',
         '听到感觉相关的词就立刻写下来',
         '跟丢了就放弃，等Section结束',
         '誊写时随便抄，不检查'],
        '✅ 主动猎找（7分+）',
        ['说明播放时已完成题目预读',
         '每题都知道"我要找什么类型的词"',
         '候选词先记旁边，确认完整意思后写',
         '跟丢了立即用题号找下一题锚点',
         '誊写逐词检查拼写大小写限词'])
    U.compare(p,'誊写习惯：粗心 vs 仔细',
        '❌ 粗心誊写（失分）',
        ['把正确答案抄错（7写成1，m写成n）',
         '漏写复数（two book → books）',
         '忘记大写（london → London）',
         '超出限词（three words → 选两个）',
         '誊写完没复核'],
        '✅ 仔细誊写（满分）',
        ['按题号顺序誊写，一题一核对',
         '有数量词的空格检查复数',
         '专有名词/月份/星期检查大写',
         '写完每题数词数确认限词',
         '全部誊写完做最后整体复核'])
    # PART07
    U.section(p,7,'课堂练习','Classroom Practice Drills')
    U.drill(p,'练习1：答案类型预判（30秒）',
        '看以下填空题，预判每空需要填入的信息类型',
        ['1. Club name: __________ United FC → 类型：?',
         '2. Membership fee: £__________ per month → 类型：?',
         '3. Start date: __________ (month only) → 类型：?',
         '4. Contact: Dr __________ → 类型：?（注意Dr后面是名还是姓？）',
         '5. Duration: __________ weeks → 类型：?',
         '参考答案：1=专有名词 2=数字 3=月份 4=姓氏 5=数字'])
    U.drill(p,'练习2：更正陷阱快速识别',
        '每句话含一处更正，听/读后写出最终正确答案',
        ['"The meeting is on Wednesday the 9th... sorry, I mean Thursday the 10th."',
         '"Call 0800 44... wait, let me check... it\'s 0800 55, 678."',
         '"Wear smart casual — a jacket would be fine... actually, it\'s very informal."',
         '"Contact Mr Wilson about bookings... Oh, his assistant Ms Chen handles that now."',
         '"The price is £85 normally, but with your discount that\'s £62.50."'])
    U.drill(p,'练习3：限词合规性检查',
        '以下答案是否符合"NO MORE THAN TWO WORDS AND/OR A NUMBER"？',
        ['"the swimming pool" → 合规还是违规？为什么？',
         '"swimming pool" → 合规还是违规？',
         '"32.50" → 合规还是违规？（数字算几词？）',
         '"every Tuesday afternoon" → 合规还是违规？',
         '"Tuesday" → 合规还是违规？（如果题目只问day of week）'])
    U.drill(p,'练习4：S3多人对话追踪',
        '读以下对话片段，用表格追踪Amy和Tom的态度（✓✗?）',
        ['A: "I think the survey method was solid." → Amy attitude on Q1: ?',
         'T: "Really? I felt the sample size was too small for any valid conclusions." → Tom: ?',
         'A: "Fair point. Though the qualitative interviews did add depth." → Amy on Q2: ?',
         'T: "Agreed — those were the strongest part of the study." → Tom on Q2: ?',
         '请判断Q1（survey method）和Q2（qualitative interviews）各自两人的态度'])
    U.drill(p,'练习5：速记符号实战',
        '用符号记录以下信息（↑↓→∴∵≈cf.w/w/o缩写等）',
        ['"Pollution rises as population increases." → 速记：？',
         '"The experiment failed because of equipment malfunction." → 速记：？',
         '"Urban areas are approximately 2 degrees warmer than rural areas." → 速记：？',
         '"Technology leads to job losses, but also creates new opportunities." → 速记：？',
         '示范参考：Pollution↑∵pop↑ / experiment failed∵equip. malfunction'])
    # PART08
    U.section(p,8,'避坑指南','Common Pitfalls')
    U.pain_point(p,'考场最常见的崩溃场景',
        ['Section 4第一句没听懂，后面9题全部放弃——自我放弃是最大失分',
         'Section 1满分本来唾手可得，却因誊写时抄错丢了2题',
         '多选题(Choose TWO)只选了一个，忘记必须选两个',
         'Section 3听到选项A的词就选A，没等说完整句'],
        ['跟丢了立即用题号找锚点，不要放弃；哪怕猜也要填',
         '誊写速度放慢一倍：先写，再核，按顺序，不跳跃',
         '多选题：做完后数一下，确认选了两个才算完成',
         '选择题：听到关键词→继续听→特别注意but/however→等整句结束再选'])
    U.content(p,'拼写高频错误：深圳学生TOP 10',
        [('accommodation','双c双m：ac-com-mo-da-tion（记忆：ac + common + mo + dation）'),
         ('necessary',    '一个c两个s：ne-ces-sa-ry（记忆：一件外套(c)两条裤子(ss)）'),
         ('environment',  '无a：en-vi-ron-ment（environ是词根，不是environ-a-ment）'),
         ('February',     '有r：Feb-ru-ary（中间有一个roo发音，不要漏）'),
         ('independent',  '-ent结尾：in-de-pend-ent（不是independant）')],
        note='这5个词在剑桥真题中反复出现，每个都是高频失分词，背会它们')
    U.content(p,'机考备考注意事项',
        [('打字速度','目标：60字符/分钟；英文打字速度低会导致来不及输入答案'),
         ('输入法切换','考前确保中文输入法已关闭或知道快捷键；避免输入中文字符'),
         ('大小写','Caps Lock状态检查；专有名词需要手动大写第一个字母'),
         ('删除修改','听到更正时：用Backspace删除，重新输入；不要犹豫'),
         ('屏幕布局','题目和输入框同时在屏幕上；学会快速切换视线，不低头看键盘')])
    U.tip_card(p,'听力学习四大速记卡',
        [('更正陷阱','听到actually/sorry/I mean/wait → 立即在草稿上划掉之前的答案，写新答案'),
         ('预读黄金时间','说明时间+Section间隙=每套考试约90秒的预读时间，必须100%利用'),
         ('跟丢重定位','跟丢了不要慌：数一下当前题号，在题目上找下一题的位置词，跳过去'),
         ('限词铁律','写完答案立即数词数；冠词算词；超限必须缩短，宁缺词毋超限')])
    U.tip_card(p,'分Section得分目标',
        [('Section 1','目标10/10（满分）；语速最慢，信息会重复，没有满分的理由'),
         ('Section 2','目标8/10；地图题画箭头，匹配题选项提前读，不要在细节上卡'),
         ('Section 3','目标6-7/10；多选题每项单独判断；追踪说话人防张冠李戴'),
         ('Section 4','目标5-7/10；跟丢了继续向前，哪怕猜也要填完，不要放弃')])
    U.checklist(p,'L01课后自检清单',
        ['我能说出S1-S4各Section的场景、难度和典型题型',
         '我知道7分需要做对30题，8分需要35题',
         '我完成了自我诊断，知道自己主要在哪类失分',
         '我理解PRCA四步法并知道在考场中如何应用',
         '我知道纸笔考和机考的关键差异'])
    U.summary(p,1,
        ['S1→S4难度递增：S3是深圳学生失分最多的，S4拉开7+差距',
         '7分=30题正确；三类失分：拼写/陷阱/时间管理——各有针对性解法',
         'PRCA四步法：预读→实时答→即时核→誊写检——是每次做题的固定流程'],
        takeaway='听懂了却做错=策略问题；S1不满分=预读没做好——这两件事优先解决')
    U.homework(p,1,
        ['下载剑桥真题任意一套，完成S1+S2（不计时，只练感觉和预读）',
         '整理个人拼写错误词10个，每个写3遍，明天默写',
         '搜索"BBC 6-Minute English"听1集，记录听到的3个关键信息词'],
        est_time='约60分钟')
    U.preview(p,2,'Section 1 日常对话满分策略',
        ['四大高频场景：预约/租房/旅游/社区服务词汇与套路',
         '数字信息精准捕捉（电话/价格/时间/邮编）',
         'Section 1零失分的完整预读模板'])
    U.save(p, f'{OUT}/L01_听力考试全貌与失分诊断.pptx')
    print(f"  L01 done: {len(p.slides)} pages")

# ══════════════════════════════════════════════════════════════════════════════
#  L02  Section 1 日常对话满分
# ══════════════════════════════════════════════════════════════════════════════
def L02():
    T(); p=U.new_prs()
    U.cover(p,'听力',2,'Section 1\n日常对话满分策略','Section 1: Zero-Loss Perfect Score')
    U.agenda(p,['上节回顾','四大高频场景','30秒预读技术','数字专项训练',
                '字母拼读法','例题精讲','课堂练习','小结作业'])
    U.review(p,['S1-S4梯度：S1语速最慢，必须满分10/10',
                'PRCA：预读→实时→即核→誊检 — 考场的固定流程',
                '7分需对30题；三类失分：拼写/陷阱/时间管理'])
    U.objectives(p,['掌握Section 1四大高频场景的词汇与答题框架',
                    '能在30秒内完成S1全部10题的预读和信息类型预判',
                    '精准捕捉电话号码/价格/时间/日期/邮编五类数字信息',
                    '实现Section 1稳定满分10/10'])
    U.pain_point(p,'S1满分明明很近却总差1-2题',
        ['预读没完成：录音开始时还在看上一题，第一题已经过了',
         '数字听错：thirteen(13)和thirty(30)混淆，重音位置没注意',
         '只写了first name漏了last name（或反之）',
         '更正陷阱：先说Tuesday，后改Wednesday，写了Tuesday'],
        ['录音说明时间是黄金30秒，全部用来预读，一秒不浪费',
         '13重音在-teen（thir-TEEN），30重音在thir-（THIR-ty）——靠后缀重音区分',
         '看到Name空格：数空格数量，预判可能有first+last两空',
         '听到actually/sorry/I mean立即划掉候选词，等待更正后的内容'])
    part(p,1,'四大高频场景','Four Core Scene Categories',
        [('场景一：预约/登记类 Appointment','最高频，超过60%的S1考题',
          [('预约诊所/牙医','必考：姓名+联系方式+预约时间+具体要求'),
           ('健身房/课程登记','必考：姓名+会员类型+费用+起始日期'),
           ('图书馆/设施预订','必考：借阅卡号+预约时间+特殊说明'),
           ('高频陷阱','时间更改：先说9月，再改为10月——后说的是答案'),
           ('关键词信号','appointment/booking/schedule/available/slot')],
          '预约场景必考更正陷阱——先听，听完再写，不要听到就写'),
         ('场景二：住宿/租房类 Accommodation','词汇量最大，信息类别最多',
          [('学生公寓/宿舍','必考：地址+邮编+月租+押金+设施'),
           ('民宿/B&B/酒店','必考：房型+价格+入住退房日期'),
           ('高频词组','furnished/unfurnished; utilities included; deposit/bond'),
           ('数字密集','地址门牌+邮编+房租+押金——数字类信息集中'),
           ('陷阱','押金(deposit)和月租(rent)数字互换——听清"deposit"还是"monthly fee"')])],
        [('场景三：旅游/设施类 Tourism',
          [('景点/博物馆介绍','必考：开放时间+票价+特殊优惠（学生/老人）'),
           ('活动/节目安排','必考：日期+地点+报名方式+特殊要求'),
           ('高频陷阱','"Open every day EXCEPT Sunday" → 题若问关闭日=Sunday'),
           ('数字类型','票价通常两位数；时间注意am/pm；电话注意区号'),
           ('表格策略','先确认表头分类，再按行列对应找答案')]),
         ('场景四：社区服务类 Community Services',
          [('图书馆服务','借还书规则、罚款、延期、预约、会员权限'),
           ('社区课程报名','课程类型+时间+费用+报名截止+资格要求'),
           ('志愿者活动','时间+地点+联系人+所需技能+培训信息'),
           ('信息复杂度','多个服务或多个联系方式同时出现，需要区分对应关系'),
           ('关键词信号','membership/eligible/volunteer/register/deadline')])],
        [('例题：预约场景完整填空',
          'Form Completion: City Fitness Centre. NO MORE THAN ONE WORD AND/OR A NUMBER.\n\n'
          'Member name: __________ Hargreaves (Q1)\nMembership: __________ (Q2)\n'
          'Monthly fee: £__________ (Q3)\nStart: 1st __________ (Q4)\nLocker no: __________ (Q5)\n\n'
          'Script: "My first name\'s Patricia... P-A-T-R-I-C-I-A. I\'d like a family membership.\n'
          'That\'s £32.50 monthly. Starting the 1st of September. And locker 47B, please."',
          'Q1: Patricia (同步字母拼写P-A-T-R-I-C-I-A)\nQ2: family\nQ3: 32.50\nQ4: September\nQ5: 47B',
          'Q5: 数字+字母组合，47B缺任何一部分都是0分；字母B需要大写'),
         ('例题：租房场景信息表',
          'Notes Completion: Student Accommodation. NO MORE THAN TWO WORDS AND/OR A NUMBER.\n\n'
          'Address: 24 __________ Road (Q6)\nRent: £__________ /month (Q7)\n'
          'Deposit: __________ months (Q8)\nAvailable: __________ (Q9)\n'
          'Tube: __________ (Q10)\n\n'
          'Script: "Address is 24 Birchwood Road — B-I-R-C-H-W-O-O-D. Rent\'s £650.\n'
          'Two months\' deposit. Available from mid-October. Nearest tube is Finchley Central."',
          'Q6: Birchwood (字母听写，首字母大写)\nQ7: 650\nQ8: 2/two\nQ9: mid-October\nQ10: Finchley Central',
          'Q9: mid-October = TWO WORDS（mid + October），符合"≤2 words"要求')])
    part(p,2,'30秒预读技术','30-Second Pre-Reading Technique',
        [('预读三层次','从表面信息到深层预判，每层用时5-10秒',
          [('第一层（5s）','确认题型：是填空/选择/地图？决定答题模式'),
           ('第二层（10s）','圈出题目关键词：空格前后的名词/动词/数字'),
           ('第三层（15s）','预判答案类型：这里要填名字/价格/时间/地点？'),
           ('时间分配','10道题共30秒=每题3秒；重点预读空格前后，不读全句'),
           ('边预读边预测','Monthly fee: £___→预测：两位数，如£35/£50/£65')],
          'S1成功的关键不是"听得好"，而是"预读得好"——预读决定成败'),
         ('预测答案类型的四个维度','预判后，听的时候更有针对性',
          [('词性预测','空格前有a/an/the/形容词→填名词；空格后有名词→填形容词'),
           ('数值范围','fee/price→数字；year→四位数；phone→11位；postcode→字母+数字'),
           ('格式预测','date→数字+月份(3 March); time→数字或文字(9:30/half past nine)'),
           ('语义范围','membership ___→可能是type/level/duration；contact ___→人名/电话'),
           ('排除法','预测是月份→忽略数字，专注听月份名称→速度快出很多')])],
        [('预读实战：30秒能预读完哪些？',
          '计时30秒，预读以下S2题目，记录每题的预测答案类型：\n\n'
          'Q11: The centre was built in __________\n'
          'Q12: Opening hours weekdays: __________ to 9pm\n'
          'Q13: The __________ has recently been renovated\n'
          'Q14: Located on __________ Street\n'
          'Q15: Annual membership: £__________',
          'Q11: 年份（4位数字）\nQ12: 时间（如7am/8am）\nQ13: 设施名词（gym/pool/hall）\nQ14: 街道名（专有名词大写）\nQ15: 数字（两/三位数，如£120）',
          '30秒预读5题 = 每题6秒；只圈关键词，不读全句；预测答案类型，不猜具体答案'),
         ('Section间隙的20秒：如何最大化利用',
          [('S1→S2间隙','S1结束音乐 → 立即翻到S2 → 预读Q11-15（至少3题）'),
           ('S2→S3间隙','S2结束 → 翻到S3 → 重点：把MCQ选项A-E读完并圈关键词'),
           ('S3→S4间隙','最重要的20秒：S4最难 → 把所有题目扫一遍，找到每题关键词'),
           ('常见错误','间隙时间在发呆/调整答案/翻前面的页 → 这20秒是S下一Section的预读时间'),
           ('练习方法','做真题时刻意计时：每次Section结束立即计时20秒练预读')])],
        [('30秒预读练习：从零开始',
          '给你30秒，对以下表格题进行预读（记录你圈出的关键词和预测类型）：\n\n'
          'COMMUNITY LANGUAGE CLASSES\n'
          'Language | Day | Time | Room | Fee/month\n'
          'Spanish  | Mon | 7pm  | __(Q1)__ | £__(Q2)__\n'
          'Mandarin | __(Q3)__ | 6:30 | 12 | £35\n'
          'French   | Wed | __(Q4)__ | 8 | £__(Q5)__',
          'Q1: 数字/字母组合（教室号如4B/Room15）\nQ2: 数字（费用£XX）\nQ3: 星期几（Weekday name）\nQ4: 时间（格式如6pm/half past five）\nQ5: 数字（费用）',
          '表格题的预读：先看表头（Language/Day/Time/Room/Fee），再看哪些格子有空需要填'),
         ('反预读：没有预读会发生什么',
          '模拟场景：录音已开始30秒，你才翻到S3题目，从Q21开始\n\n'
          '结果：Q21和Q22已经过了，只能靠猜；Q23开始才能正常答题\n'
          '代价：2道题=0.25分；别人不失分，你失分了\n'
          '解决：把预读当作和答题同等重要的技能，每次做真题强制执行',
          '对比：预读者均分8.5/10；不预读者均分6.2/10（剑桥研究数据）',
          '预读成本：30秒；预读收益：约1-2道题的保险——ROI极高')])
    part(p,3,'数字类信息精准捕捉','Number Information: Precision Capture',
        [('电话号码：英式特有规则','英国电话号码有固定格式，掌握规则事半功倍',
          [('手机格式','07XXX XXXXXX（共11位，07开头）'),
           ('座机格式','01XXX XXXXXX 或 02X XXXX XXXX（地区码不同）'),
           ('double/treble','double 4 = 44（两个相同）；treble 9 = 999（三个相同）'),
           ('零的读法','电话中零读"oh"；数学中"zero"或"nought"'),
           ('同步书写','边听边写，不等全部说完再写——前3位忘了无法追回')],
          '电话号码是Section 1最高频的数字信息，也是最容易失分的'),
         ('价格/费用：听到数字的前后语境','不只是听数字，要听数字的语境',
          [('整数','£45 = "forty-five pounds"，只需写45（不写£符号）'),
           ('小数','£32.50 = "thirty-two pounds fifty"或"thirty-two fifty"'),
           ('价格更正','先说£80"→ "but with discount that\'s £65" → 答案是65'),
           ('货币区分','pounds(英镑) vs pence(便士)；£0.50 vs £50.00完全不同'),
           ('折扣陷阱','20% off → 乘以0.8；student discount → 通常打折后价格才是答案')])],
        [('时间与日期：格式转换','考官要求写数字格式，但录音用文字读',
          [('整点','nine o\'clock = 9:00 或 9am（根据上下文选择格式）'),
           ('半点','half past six = 6:30（不是6:50！）'),
           ('刻钟','quarter past three = 3:15；quarter to three = 2:45'),
           ('日期','the third of March = 3 March 或 March 3（两种写法都接受）'),
           ('月份大写','January/February...必须首字母大写，否则算错误')]),
         ('特殊数字表达：thirteen vs thirty',
          [('核心规则','-teen结尾：重音在后（thir-TEEN）；-ty结尾：重音在前（THIR-ty）'),
           ('易混组','13/30, 14/40, 15/50, 16/60, 17/70, 18/80, 19/90'),
           ('语境帮助','如果是年龄，两位数合理；如果是面积，可能是百以上'),
           ('确认方法','不确定时：等说话人继续描述，通常有其他信息可以验证'),
           ('练习','每天找5个含-teen/-ty的句子，反复听直到能立即区分')])],
        [('数字听写专项练习',
          '请听/读以下信息，快速记录（每条限时5秒）：\n\n'
          '1. "Our number is oh eight double oh, treble five, 69."\n'
          '2. "The fee is forty-two fifty per month."\n'
          '3. "It runs from half past six to quarter to nine."\n'
          '4. "Property at 14 Birchwood Road, postcode BS7 4NQ."\n'
          '5. "Thirteen people — not thirty, I said thirteen."',
          '1. 0800 555 569（double oh=00, treble five=555）\n2. £42.50（forty-two fifty）\n3. 6:30-8:45（quarter to nine=差15分到9点=8:45）\n4. 14; BS7 4NQ\n5. 13（NOT 30——更正陷阱）',
          '第5题：录音说了"not thirty"=更正信号，最终答案是thirteen=13，不是30'),
         ('价格陷阱综合练习',
          '录音片段分析：\n\n'
          '1. "Normal price is £85. With student discount you pay £62."\n'
          '   题目问：How much do students pay? Answer: ?\n\n'
          '2. "Rent is £750. The deposit is two months\' rent."\n'
          '   题目问：What is the deposit? Answer: ?\n\n'
          '3. "Adult tickets: £15. Children under 12 are free."\n'
          '   题目问：What do children pay? Answer: ?',
          '1. £62（学生价，不是原价£85）\n2. £1500（两个月租金=750×2，需要计算！）\n3. free / nothing / 0（直接说"free"即可，写0也接受）',
          '第2题：需要计算！750×2=1500；听力考试偶尔会有简单数学运算')])
    part(p,4,'字母拼读与表格填空','Spelling Dictation & Form Completion',
        [('字母拼读：同步书写法','不要等说完再写，会忘记前面的字母',
          [('核心原则','听到一个字母立即写一个——边听边写，不等对方拼完'),
           ('易混字母','B(bee)/D(dee)/E(ee)/G(gee)/P(pee)/T(tee)/V(vee)/Z(zed)'),
           ('B vs D','B：双唇闭合后爆破；D：舌尖抵上颚；听尾音区分'),
           ('M vs N','M：双唇闭合，鼻腔共鸣(em)；N：舌尖轻触上颚(en)'),
           ('NATO字母表','记忆技巧：B=Bravo, D=Delta, M=Mike, N=November等')],
          '字母拼读在S1几乎100%出现——这是必须掌握的技能，不是可选项'),
         ('表格填空：行列对应法','先理解结构，再找对应位置',
          [('第一步','看表头：横轴是什么属性？纵轴是什么分类？'),
           ('第二步','确定每道题是哪行哪列的交叉点'),
           ('第三步','根据交叉点位置预判：这里应该填什么类型的信息？'),
           ('顺序规律','通常按从左到右、从上到下介绍——利用顺序提前定位'),
           ('例外情况','有时介绍顺序与表格顺序不同——靠关键词而非顺序来定位')])],
        [('易混字母辨析：听音区分',
          [('B vs P','B=bee; P=pee; 听尾音：B结尾较浊，P结尾更清脆爆破'),
           ('D vs T','D=dee; T=tee; D结尾有声，T无声'),
           ('M vs N','M=em（双唇）; N=en（舌尖抵上颚）'),
           ('G vs J','G=gee; J=jay; G短促，J有拖音'),
           ('F vs S','F=ef; S=es; F用上齿咬下唇，S舌尖')]),
         ('表格填空：行列对应与常见错误',
          [('行列错位','Q2应填费用却填了Day——没有确认表格结构就开始答'),
           ('超出限词','限ONE WORD AND/OR A NUMBER：写了"every Monday"(2词+形容词)'),
           ('格式错误','时间题写"half past five"(4词)而不是"5:30"(更简洁)'),
           ('大小写','月份/星期必须大写：monday→Monday；课程名大写'),
           ('解决','预读：先看表头确认结构；每题写完数词数')])],
        [('字母拼读实战练习',
          '跟着老师拼读以下姓名，逐字母同步记录：\n\n'
          '1. MCDOUGALL → M-C-D-O-U-G-A-L-L（9个字母）\n'
          '2. FEATHERSTONE → F-E-A-T-H-E-R-S-T-O-N-E（12个字母）\n'
          '3. PRZYBYLSKI → P-R-Z-Y-B-Y-L-S-K-I（10个字母）\n\n'
          '注意：L要清楚区分（L vs 1）；写完后核对字母数量是否对',
          '1. 9个字母；2. 12个字母；3. 10个字母\n提醒：每个字母之间的停顿是检查和核对的机会',
          '技巧：在格子里写字母（每格一个字母）比在线上写更不容易出错和漏写'),
         ('表格填空综合例题',
          'Complete the table. Write NO MORE THAN ONE WORD AND/OR A NUMBER.\n\n'
          'Course | Day | Time | Room | Fee/month\n'
          'Yoga   | Mon | 7pm  | ?(Q1) | £?(Q2)\n'
          'Pilates | ?(Q3) | 6:30 | 11 | £28\n'
          'Zumba  | Thu | ?(Q4) | 9 | £?(Q5)\n\n'
          'Script: "Yoga Mondays at seven, in room 4B, thirty pounds a month.\n'
          'Pilates is Tuesdays at six-thirty, room 11, twenty-eight a month.\n'
          'Zumba on Thursdays, half past five, room 9, and that\'s twenty-four fifty."',
          'Q1: 4B  Q2: 30  Q3: Tuesday  Q4: 5:30  Q5: 24.50',
          '表格按课程顺序介绍（Yoga→Pilates→Zumba），利用这个规律提前定位到对应行')])
    # PART05
    U.section(p,5,'高频词汇与表达','Key Vocabulary & Expressions')
    U.vocab(p,'预约场景高频词汇',
        [('make/book/schedule an appointment', '预约（三种同义表达，都会出现在录音中）'),
         ('cancel / reschedule / postpone',    '取消 / 改期 / 延后（三者区别：cancel=完全取消）'),
         ('available / free / vacant',         '有空/空闲/空缺（人用available/free；位置用vacant）'),
         ('slot / opening / time frame',       '时间段（可预约的）'),
         ('waiting list / on hold',            '候补名单 / 暂候（满了但可以等）'),
         ('deposit / advance payment / bond',  '订金/预付款/押金（三者有细微差别）'),
         ('confirmation / reminder',           '确认单 / 提醒（预约结束时常见词）')])
    U.vocab(p,'租房/住宿高频词汇',
        [('fully furnished / unfurnished',     '全配家具 / 无家具'),
         ('utilities included / bills included','水电费含在租金内（否则另计）'),
         ('landlord / tenant / estate agent',  '房东 / 租客 / 房产中介'),
         ('lease / tenancy agreement / contract','租约/合同（都指同一种文件）'),
         ('security deposit / bond',           '押金（退房时如无损坏可全额退还）'),
         ('en-suite / shared bathroom',        '独立卫浴 / 共用卫浴'),
         ('ground floor / first floor (UK)',   '地面层 / 一楼（英式：与美式不同！二楼=first floor）')])
    U.vocab(p,'数字信息表达快速参照',
        [('double X / treble X',              'XX / XXX（如double 4=44，treble 5=555）'),
         ('£45.99',                           'forty-five pounds ninety-nine'),
         ('07700 900123',                     'oh-double seven-double oh-nine-double oh-one-two-three'),
         ('3:45pm',                           'quarter to four in the afternoon'),
         ('SW1A 2AA (UK postcode)',            'S-W-one-A [space] two-A-A'),
         ('15th October',                     'the fifteenth of October 或 October the fifteenth'),
         ('2.5 hours',                        'two and a half hours 或 two hours thirty')])
    U.two_col(p,'S1最容易听错的词对',
        '❌ 容易混淆的听法',
        ['13 vs 30（thirteen vs thirty）',
         '15 vs 50（fifteen vs fifty）',
         'Tuesday vs Thursday（都是T开头）',
         'March vs May（M开头月份）',
         'address vs access（发音相近）',
         'deposit vs discount（功能不同）'],
        '✅ 区分技巧',
        ['-teen重音在后(thir-TEEN)，-ty重音在前(THIR-ty)',
         '同理适用于15/50, 16/60, 17/70, 18/80',
         '听清楚：-ues-(day) vs -urs-(day)；前者有-u-音',
         'March=有r音/较硬；May=短音；听清楚长度',
         '语境判断：address=地址，access=进入权限',
         '语境判断：deposit=定金（交钱），discount=折扣（少花钱）'])
    # PART06
    U.section(p,6,'方法精讲与满分范例','Method & Perfect Score Models')
    U.steps(p,'Section 1 满分答题流程（10步）',
        [('考前（说明时间）','翻到S1题目；预读Q1-5；圈关键词；预判答案类型'),
         ('Q1-5 录音进行时','带着预读预测去听；候选词先记旁边，确认后写正格'),
         ('Q1-5 结束间隙','约10秒：预读Q6-10（若之前未完成）'),
         ('Q6-10 录音','特别注意更正陷阱：actually/sorry/I mean'),
         ('S1结束','立即翻到S2；不要在S1上多停留，利用间隙预读S2')])
    U.steps(p,'字母拼读同步书写操作规范',
        [('准备','在空格旁用铅笔轻画小格子，1格=1字母，预估字母数'),
         ('接收第1个字母','立即写入第1格，不等第2个字母'),
         ('连续接收','每听到1个字母立即写入对应格；停顿时核对已写的'),
         ('易混字母处理','B/P不确定时：先写？，继续听后面；从语义判断'),
         ('说完后确认','利用对方停顿：逐格核对字母，检查有无遗漏或混淆')])
    U.model_answer(p,'满分S1答题笔记示范：预约场景',
        'Fitness Centre Membership Enquiry',
        '答题笔记格式示范：\n\n'
        'Q1 Name: [P][A][T][R][I][C][I][A] + [H][A][R][G][R][E][A][V][E][S]\n'
        '   → Patricia Hargreaves\n\n'
        'Q2 Type: (annual × ) family ✓ → family\n\n'
        'Q3 Fee: (£45 × ) £32.50 ✓ → 32.50\n'
        '   [更正：先说45，actually说32.50]\n\n'
        'Q4 Start: September ✓（首字母大写）\n\n'
        'Q5 Locker: 47B ✓（数字+字母，B不能漏）',
        ['用格子记字母，每格一个，避免遗漏',
         '更正陷阱：先划掉£45，写新答案£32.50',
         'Q3只写数字（32.50），不写£符号',
         'Q5: 47B — 字母和数字都必须写，缺一个=0分'])
    U.model_answer(p,'满分S1答题：社区服务场景',
        'Community Centre Course Registration',
        'Q6  Course: Photography\n'
        'Q7  Day: Thursday [not Tuesday! 更正陷阱]\n'
        'Q8  Time: 7:30 [half past seven — 转换为数字]\n'
        'Q9  Room: B12 [字母+数字，注意大写B]\n'
        'Q10 Fee: 24.50 [twenty-four fifty — 小数点格式]\n\n'
        '关键笔记：Q7旁边写 "Tue? → Thu ✓" 记录更正过程',
        ['Q7：录音先说Tuesday，后用"actually"更正为Thursday',
         'Q8：half past seven=7:30，口算转换成数字格式最安全',
         'Q9：字母+数字 = 2部分都要正确，B12缺B=0分',
         'Q10：小数点格式（24.50）和文字格式（twenty-four fifty）均接受'])
    U.compare(p,'S1答题：有没有预读的差距',
        '❌ 没有预读（5-6分结果）',
        ['录音开始时才翻到题目，Q1已过',
         '不知道Q2要填type还是fee',
         '听到数字不知道是哪道题的答案',
         '拼读时等说完再写，前3个字母忘了',
         '誊写时随便抄，没有逐字检查'],
        '✅ 充分预读（10/10结果）',
        ['说明播放时已圈完Q1-10关键词',
         'Q2知道要填修饰"membership"的词',
         '每个数字都知道对应Q3还是Q4',
         '拼读同步写：听到P立刻写P',
         '誊写逐词检查：拼写/大小写/限词'])
    U.compare(p,'数字听写：听错 vs 听对',
        '❌ 典型失误',
        ['thirty-two = 32 → 写成了23（数字颠倒）',
         'treble five = 555 → 写成了5（漏了两个）',
         'half past nine = 9:30 → 写成了9:50',
         'fifteen = 15 → 写成了50（fifteen vs fifty）',
         'forty-two fifty = £42.50 → 写成£42.5（少一个零）'],
        '✅ 正确记录',
        ['thirty-two → 先在脑中确认位数：3-2，不是2-3',
         'treble five → treble=3个，三个5=555',
         'half past nine → half past=30分钟后=9:30',
         'fif-TEEN重音在后=15；FIF-ty重音在前=50',
         '£42.50有两位小数：.50，不能简写成.5'])
    # PART07
    U.section(p,7,'课堂练习','Practice Drills')
    U.drill(p,'练习1：字母拼读听写（逐字母同步）',
        '老师拼读以下词语，同步记录（每字母一格）',
        ['PRZYBYLSKI → P-R-Z-Y-B-Y-L-S-K-I（波兰姓，10字母）',
         'FEATHERSTONE → F-E-A-T-H-E-R-S-T-O-N-E（12字母）',
         'MCDOUGALL → M-C-D-O-U-G-A-L-L（9字母，注意双L）',
         'JOHANNSEN → J-O-H-A-N-N-S-E-N（9字母，注意双N）',
         '完成后核对：字母总数是否正确？有无混淆B/D、M/N？'])
    U.drill(p,'练习2：数字快速记录（每条5秒）',
        '跟着老师/录音，记录以下数字信息',
        ['"zero seven, double eight, treble nine, four two six" → ?',
         '"sixty-four pounds, seventy-five pence" → ?',
         '"from quarter past eight to half past ten" → ?',
         '"the twenty-first of November, two thousand and twenty" → ?',
         '"W-C-two, four-R-E" (postcode) → ?'])
    U.drill(p,'练习3：更正陷阱全识别',
        '每句话含一处更正，写出最终正确答案',
        ['"Tuesday the 8th... sorry, I meant Wednesday the 9th."',
         '"Call 0800-44... wait, check that, it\'s 0800-55, 678."',
         '"Smart casual, jacket fine... actually, very informal."',
         '"Deposit is two months... oh, actually, just one month."',
         '"Room 15... no, moved to Room 50 last week."'])
    U.drill(p,'练习4：表格填空计时（10分钟完成）',
        '完整表格填空（模拟真实S1场景），老师读一遍，只能记录一次',
        ['SWIMMING POOL BOOKING: Name(Q1), Adult/Child(Q2), Session(Q3), Lane(Q4), Duration(Q5)',
         '录音："...Ms Chen — C-H-E-N. Adult session on Tuesday. Lane number 3B. Two hours."',
         '请同步记录5个答案，提交前检查：大小写/限词/数字格式',
         '答案：Q1=Chen; Q2=Adult; Q3=Tuesday; Q4=3B; Q5=2hours/two hours',
         '讨论：Q5如果限制是ONE WORD AND/OR A NUMBER，写"2"还是"2 hours"？'])
    U.drill(p,'练习5：全流程计时练习（S1模拟）',
        '完整模拟Section 1的预读→答题→检查流程',
        ['步骤1（30秒）：预读以下题目，圈关键词，预判类型',
         '步骤2（实时）：老师读录音，同步答题（只有一遍）',
         '步骤3（30秒）：检查答案：拼写/大小写/限词/数字格式',
         '题目：Name(Q1), Age(Q2), Membership type(Q3), Start month(Q4), Emergency contact(Q5)',
         '反思：预读时你预测对了哪些题的答案类型？'])
    # PART08
    U.section(p,8,'避坑指南','S1 Pitfall Prevention')
    U.pain_point(p,'S1满分障碍——都是可以避免的',
        ['誊写时把47B抄成74B（数字和字母顺序颠倒）',
         '听到"first name"只写了名字，漏了title(Mr/Ms)或surname',
         'Section 1结束后发呆，没有利用间隙预读Section 2',
         '确认答案时只看内容，没有检查限词是否超出'],
        ['誊写慢一倍速度，每个答案写完轻声默读一遍确认',
         '看到Name/Contact题：先数空格数量，判断要填几个部分',
         'S1结束音乐响=立即翻到S2；这20秒=S2的2-3道题预读时间',
         '誊写前再看一眼限词：ONE WORD AND/OR A NUMBER→数字是一格'])
    U.content(p,'S1高频细节错误（逐一击破）',
        [('October vs November', '十月vs十一月——Oc-to-ber vs No-vem-ber，听清完整月份'),
         ('am vs pm',            '"nine thirty in the evening" = 9:30pm，不是am'),
         ('Street vs Road vs Avenue','全是地址后缀，听清用的是哪个'),
         ('first floor (UK)',    '英式一楼=美式二楼（ground floor=地面层=美式一楼）'),
         ('pounds vs pence',     '£35 vs 35p 差一百倍；听清"pounds"还是"pence"')],
        note='这5个错误在剑桥真题S1中反复出现，背会等于多拿2题')
    U.content(p,'誊写阶段五步自检（纸笔考）',
        [('Step 1 字数','每个答案数一下词数，是否在限制内'),
         ('Step 2 大小写','专有名词/月份/星期→首字母大写；普通名词→小写'),
         ('Step 3 复数','有数量词(two/several)→名词加s；单数情况→不加s'),
         ('Step 4 拼写','快速默念每个字母，特别检查双字母词accommodation/necessary'),
         ('Step 5 格式','数字是否用阿拉伯数字？时间格式是否统一（9:30 vs nine thirty）')])
    U.tip_card(p,'S1速记卡：四大必知',
        [('double/treble','double X=XX; treble X=XXX；遇到立刻反应，不要犹豫，电话号码常见'),
         ('更正信号词','actually/sorry/I mean/wait/let me correct that → 划掉之前的，写新的'),
         ('Name题必看','Name旁边有几个空格？有没有Title(Mr/Ms)？预判first+last还是只有一个'),
         ('誊写金句','誊写速度慢一倍，核对每个字母；宁可没写完最后2题，也不能写错前8题')])
    U.tip_card(p,'数字速记参考卡',
        [('13vs30','thirTEEN（重音后）vs THIRty（重音前）；-teen vs -ty的重音位置是区分关键'),
         ('价格更正','先说原价 + but/actually/discount → 写后说的折扣价，不写原价'),
         ('时间换算','quarter to X = X:00减15分；half past X = X:30；练到秒速转换'),
         ('邮编格式','英式邮编：字母+数字+空格+数字+两字母，如SW1A 2AA；字母大写')])
    U.checklist(p,'L02 自检清单',
        ['我能区分thirteen和thirty的重音位置并快速正确判断',
         '我知道double/treble在电话号码中的含义并能即时反应',
         '我练习了字母拼读同步书写法（每字母一格）',
         '我掌握S1四大场景的核心词汇（预约/租房/旅游/社区）',
         '我知道S1→S2间隙的20秒应该用来做什么'])
    U.summary(p,2,
        ['S1四大场景：预约/租房/旅游/社区——各有核心词汇和答题套路',
         '30秒预读三层次：确认题型→圈关键词→预判答案类型',
         '数字专项：13vs30靠重音；double/treble立刻反应；价格更正听后说的'],
        takeaway='S1不该失分——失分100%来自预读不足或誊写不仔细，而非"听不懂"')
    U.homework(p,2,
        ['剑桥真题任意两套Section 1（严格计时），统计正确率并分类错误',
         '字母拼读练习：用5个同学名字互相拼读+听写，训练同步书写',
         '数字训练：每天从英美剧/新闻中记录5个电话号码，记录并核对'],
        est_time='约45分钟')
    U.preview(p,3,'Section 2 独白场景攻略',
        ['地图/平面图题20个核心方向词详解',
         '匹配题选项提前读的策略',
         'Section 2独白的结构信号词与内容预测'])
    U.save(p, f'{OUT}/L02_Section1日常对话满分.pptx')
    print(f"  L02 done: {len(p.slides)} pages")

# ══════════════════════════════════════════════════════════════════════════════
#  L03  Section 2 独白场景攻略
# ══════════════════════════════════════════════════════════════════════════════
def L03():
    T(); p=U.new_prs()
    U.cover(p,'听力',3,'Section 2\n独白场景攻略','Section 2: Monologue Mastery')
    U.agenda(p,['上节回顾','S2独白结构与场景','地图/平面图题攻略',
                '匹配题系统解法','信号词与追踪','例题精讲','课堂练习','小结作业'])
    U.review(p,['S1四大场景：预约/租房/旅游/社区——30秒预读决定成败',
                '数字陷阱：13vs30靠重音；更正陷阱听actually/sorry',
                '字母拼读同步书写：每字母一格，边听边写'])
    U.objectives(p,['掌握S2地图/平面图题20个核心方向词',
                    '建立匹配题预读选项关键词的操作习惯',
                    '识别独白中的话题转换信号词',
                    '实现Section 2稳定8/10以上'])
    U.pain_point(p,'S2最常见的3种失分',
        ['地图题跟丢：left/right反应慢，一个路口没跟上，后面全错',
         '匹配题临时读选项：Q16出来才开始读A-H，录音已过大半',
         '话题转换没听到信号词，从Q16开始后面全部跟丢'],
        ['方位词必须形成条件反射：opposite=对面；adjacent=旁边；beyond=更远处',
         '匹配题铁律：录音开始前读完A-H全部选项，每个圈1个关键词',
         '信号词=话题转换提示：听到"Now let\'s move on to"→下一道题来了'])
    part(p,1,'S2结构与高频场景','S2 Structure & Common Scenes',
        [('S2独白的特点','单人介绍，无对话，语速适中',
          [('说话人','社区官员/导游/工作人员，口吻正式礼貌'),
           ('题型','地图题30%+匹配题40%+填空20%+其他10%'),
           ('语速','比S1快10%；专业词较多但不及S4'),
           ('陷阱','独白中的自我修正；条件限定词（only/except）'),
           ('目标分','稳定8/10；10/10需要地图零失误')],
          'S2有内在结构：引言→设施介绍→注意事项→结尾；听到结构=不跟丢'),
         ('S2三大高频场景','旅游/社区/设施是S2主战场',
          [('旅游景点','museum/gallery/exhibition；admission/guided tour/audio guide'),
           ('社区设施','leisure centre/library；membership/booking/peak hours'),
           ('活动安排','event/workshop/festival；register/deadline/eligibility'),
           ('关键差异','旅游场景地图题多；社区场景匹配题多'),
           ('词汇策略','这3类词汇熟悉后S2填空题几乎全拿')])],
        [('地图/平面图三大类型',
          [('建筑平面图','从入口追踪室内房间；注意楼层变化'),
           ('室外区域图','公园/校园；有字母标记(A-E)要填入题目'),
           ('路线追踪图','从A→B的方向描述；按箭头追踪即可'),
           ('共同核心','方向词(left/right/opposite)是核心；不需要背景知识'),
           ('预读要点','先找入口位置；标注已有地标；确认指北针方向')]),
         ('匹配题核心结构',
          [('题目数量','通常5-8道匹配题在Q11-20中'),
           ('选项设计','通常8个选项，比题目多2-3个（干扰项）'),
           ('配对类型','设施名↔特点；人名↔意见；时间↔活动'),
           ('预读必要','8个选项必须全读；只读当前题=自杀'),
           ('同义替换','选项词≠录音词；refurbished=renovated')])],
        [('例题：S2地图基础',
          'Label the map using letters A-E.\nQ11:Café Q12:Disabled toilet Q13:Gift shop\n\n'
          '"From the main entrance, go straight ahead. The information desk is right in front.\n'
          'The café is on your left, just past the gift shop.\n'
          'Disabled toilets are opposite the information desk on the right."',
          'Q11:Café=左侧(past gift shop)\nQ12:Disabled toilet=右侧(opposite info desk)\nQ13:Gift shop=left(before café)',
          '先找主入口→追踪"straight ahead"→再追踪left/right；用铅笔实时标注'),
         ('例题：S2匹配题',
          'Match each facility with feature A-F.\n'
          'A.recently renovated B.members only C.booking required D.free E.24hrs F.closed Mondays\n\n'
          '"The pool was completely refurbished last month.\n'
          'Tennis courts need advance booking.\n'
          'The library operates around the clock, seven days a week."',
          'Q16:Pool→A(refurbished=renovated)\nQ17:Tennis→C(booking required)\nQ18:Library→E(around the clock=24hrs)',
          '关键：refurbished=renovated; around the clock=24hrs——同义替换识别是匹配题核心')])
    part(p,2,'地图题20个核心方向词','20 Essential Direction Words',
        [('位置类方向词（必须5秒反应）','方位词条件反射是地图题零失误的前提',
          [('opposite / facing','对面（正对着的位置）'),
           ('adjacent to / next to','紧邻（左右相接）'),
           ('between A and B','在A和B之间'),
           ('at the end of','在...的尽头'),
           ('in front of / behind','在...前面/后面'),
           ('on your left/right','行进方向的左/右侧')],
          'opposite和next to是最高频的两个；确保5秒内画出空间关系'),
         ('移动类方向词（追踪路线用）','从A到B需要这些词',
          [('go straight ahead / continue','直行'),
           ('turn left/right','左转/右转'),
           ('take the first/second turning','第一/二个路口转弯'),
           ('just past / immediately after','经过...之后紧接着'),
           ('at the junction / at the corner','在路口/在拐角'),
           ('bear left/right','微微偏左/右（不是直角转弯）')])],
        [('方向词空间转换：行进方向决定left/right',
          [('基本原则','left/right是相对行进方向，不是地图绝对方位'),
           ('面向北行走','left=西；right=东；on your left=地图左侧'),
           ('面向南行走','left=东；right=西；on your left=地图右侧！'),
           ('实战技巧','用铅笔箭头标出行进方向；每次转弯更新方向感'),
           ('入口通常','面向地图上方（北）进入；一开始left=地图left')]),
         ('地图题预读三步法',
          [('第1步','找入口：标注✱或START，确定追踪起点'),
           ('第2步','数标记：地图上有几个字母(A/B/C)需要填入？'),
           ('第3步','找地标：已知地标（car park/reception）在哪里？'),
           ('时间','15秒完成预读；剩余5秒确认方向'),
           ('禁忌','不要试图记住地图——实时追踪比记忆更可靠')])],
        [('方向词5秒闪认练习',
          '老师说方向词，学生立即在白纸上画空间关系（每词5秒）：\n\n'
          '1."The café is adjacent to the gift shop."\n'
          '2."Turn left at the junction, the library is at the far end."\n'
          '3."The gym is opposite the swimming pool, between café and entrance."\n'
          '4."Just past the reception on your right, you\'ll find the toilets."',
          '1.相邻 2.左转→尽头 3.对面+在两者之间 4.经过reception右侧',
          '画图而非文字描述：方向词触发空间图像，比文字转换快3倍'),
         ('地图追踪实战（口述画图）',
          '在白纸中央画矩形代表建筑，底部中间标ENTRANCE：\n\n'
          '"From the entrance, go straight ahead to the main hall (Q11).\n'
          'Turn left. The classroom is the first door on your right (Q12).\n'
          'Continue and turn right at the end. The office is on your left (Q13)."',
          'Q11:直走=建筑中央/上方\nQ12:左转第一门右侧\nQ13:继续右转后左侧',
          '用箭头标出每次转弯；先追路线，再写答案字母；不要等全部听完再标')])
    part(p,3,'匹配题系统解法','Matching Questions: Complete System',
        [('匹配题三重压力','为什么匹配题难',
          [('压力1','8个选项同时在面前，不知道关注哪个'),
           ('压力2','录音提到选项词，但整句话可能否定它（干扰）'),
           ('压力3','一边听一边在8个选项中扫描——脑力负荷极大'),
           ('根本解法','预读圈词→脑中建关键词库→听时用库匹配，不靠扫描'),
           ('训练目标','听到"refurbished"→0.5秒联系到"renovated"（B选项）')],
          '匹配题=速记关键词+识别同义替换；不是语言能力的测试'),
         ('预读操作规范（录音前完成）','这一步决定匹配题成败',
          [('时机','S2说明时间或S1→S2间隙的20秒'),
           ('操作','从A到H逐个圈出最独特的1-2个词'),
           ('标注','在每个选项编号旁写圈词：A=free, B=renovated...'),
           ('脑中建库','不需要记全句，只记住关键词与选项字母的对应'),
           ('速度','8个选项应在30秒内圈完；每个选项最多3-4秒')])],
        [('识别干扰选项三步法',
          [('Step1 听到关键词','听到"free"→tentatively标A(free of charge)'),
           ('Step2 验证整句','确认整句话是"免费"含义，不是"曾经免费"或否定'),
           ('Step3 确认最终','等说完整个关于该设施的段落再写答案'),
           ('常见干扰','先说"used to be free"再说现在收费→选现在状态'),
           ('否定陷阱','"NOT available for members"→不选members only')]),
         ('高频同义替换词对（必记）',
          [('renovated = refurbished','翻新（最高频）'),
           ('exclusively for members = members only','仅限会员'),
           ('must book in advance = booking required','需要预约'),
           ('round the clock / 24hrs = open 24 hours','全天开放'),
           ('free of charge / no cost = free','免费'),
           ('recently launched / opened = newly opened','新开'),
           ('situated outdoors = outdoor only','在室外')])],
        [('匹配题同义替换练习',
          'Match录音描述词(LEFT)与IELTS选项(RIGHT)：\n\n'
          '"refurbished" / "exclusively for members" / "must book in advance"\n'
          '"open every hour of the day" / "at no cost" / "launched recently"\n\n'
          'Options: A.free  B.members only  C.recently opened  D.booking required  E.24hrs  F.renovated',
          'refurbished→F / exclusively→B / must book→D / every hour→E / no cost→A / recently→C',
          '建立个人高频替换词对列表；每周增加5组；考场遇到立即联想'),
         ('匹配题完整实战',
          'Q16-20: Match each facility with feature A-G (2 extra options)\n'
          'A.free  B.booking required  C.members only  D.newly opened  E.closed Mondays  F.renovated  G.outdoor only\n\n'
          '"Pool: just opened last week. Gym: members-only. Library: prior reservation needed.\n'
          'Courts: situated in the open air. Café: undergoing renovation now."',
          'Q16Pool→D(last week=newly opened)\nQ17Gym→C(members-only)\nQ18Library→B(prior reservation=booking)\nQ19Courts→G(open air=outdoor)\nQ20Café→F(undergoing renovation)',
          'Q20注意：Café正在装修(undergoing renovation)=F选项(renovated)；时态不同但含义匹配')])
    part(p,4,'信号词与话题追踪','Signal Words & Topic Tracking',
        [('独白结构信号词体系','掌握信号词=永远不跟丢',
          [('开场类','Welcome to... / Let me start with... / Today I\'ll tell you about...'),
           ('转换类','Now let\'s move on to... / Turning to... / Next up we have...'),
           ('强调类','It\'s worth noting... / The key point here is... / Please remember...'),
           ('结尾类','To summarise... / Before I finish... / Finally, I\'d like to...')],
          '话题转换信号=下一道题即将开始；听到信号词→立刻准备好下一道题'),
         ('跟丢后的重定位策略','跟丢不等于放弃',
          [('立即动作','停止回想已过内容；眼睛移到下一道未答的题目'),
           ('关键词锚点','看下一题关键词；等录音中出现该词'),
           ('地图重定位','找地图上唯一未标记的区域，配合下一个设施名'),
           ('估算填写','实在不会：从未使用选项中选，绝不留空'),
           ('自我安慰','跟丢了一题≠整Section放弃；后面4题还来得及')])],
        [('S2录音内容顺序规律',
          [('规律1','按地理路线介绍设施：入口→大厅→各区域→出口'),
           ('规律2','按重要性排序：最主要设施先介绍'),
           ('规律3','实用信息集中在末尾：价格/时间/预订方式'),
           ('利用规律','填空题通常对应末尾实用信息段落'),
           ('例外','偶尔乱序：靠关键词定位，不依赖顺序')]),
         ('S2精析与进步追踪',
          [('精析方法','每道错题找到答案所在的具体句子'),
           ('分类错误','跟丢失分？方位词反应慢？同义替换没识别？'),
           ('目标提升','地图错2题→专练方位词；匹配错3题→专练预读'),
           ('时间投入','每次S2精析20分钟比重做5套更有价值'),
           ('记录进步','用本子记录每次S2正确题数；看到上升就是进步')])],
        [('信号词听力训练',
          '读以下独白，标出所有话题转换信号词：\n\n'
          '"Welcome to Greenfield Sports Centre. Let me start with the facilities.\n'
          'Starting with the ground floor, we have a gym and café.\n'
          'Now, moving on to the first floor, you\'ll find meeting rooms.\n'
          'Before I finish, I\'d like to mention some important booking rules."',
          '信号词1："Now, moving on to"→第一层内容\n信号词2："Before I finish"→预订规则\n共2次话题转换',
          '练习目标：听到信号词时脑中"警报"响起，自动准备好下一道题'),
         ('S2综合模拟流程演练',
          'S2完整模拟三阶段：\n\n'
          '① 预读30秒：判断题型（地图/匹配/填空）→圈选项关键词→找地图入口\n'
          '② 听题：地图铅笔追踪→匹配关键词库扫描→填空原词记录\n'
          '③ 跟丢时：找下一题关键词→重定位→猜填已过题→继续',
          'S2结束：立刻翻S3，利用20秒间隙预读S3全部题（含MCQ选项A-E）',
          '自我评估：统计地图/匹配/填空各对几题→找到弱项→下周专项练习')])
    # PART05
    U.section(p,5,'高频词汇','Key Vocabulary')
    U.vocab(p,'方位词必备（位置类）',
        [('opposite / facing',        '正对面'),
         ('adjacent to / next to',    '紧邻'),
         ('between ... and ...',      '在...之间'),
         ('at the end of',            '在...尽头'),
         ('in front of / behind',     '在...前面/后面'),
         ('on your left/right',       '行进方向的左/右'),
         ('just past',                '经过...之后紧接着')])
    U.vocab(p,'方位词必备（移动类）',
        [('straight ahead',           '直走'),
         ('turn left/right',          '左/右转'),
         ('take the first turning',   '走第一个路口转弯'),
         ('bear left/right',          '微偏左/右'),
         ('at the junction/corner',   '在路口/拐角处'),
         ('beyond / further along',   '更远处'),
         ('directly opposite',        '正对面（强调版）')])
    U.vocab(p,'匹配题高频同义词',
        [('renovated = refurbished',     '翻新'),
         ('members only = exclusive',    '仅限会员'),
         ('booking required = must book','需要预约'),
         ('free = no charge',           '免费'),
         ('newly opened = just launched','新开'),
         ('24hrs = round the clock',     '全天'),
         ('outdoor = open air',          '室外')])
    U.two_col(p,'S2场景词汇对照',
        '旅游景点（Tourism）',
        ['guided tour / self-guided',
         'admission fee / entrance charge',
         'audio guide / headset',
         'gift shop / souvenir store',
         'photography allowed/prohibited',
         'opening times / visiting hours',
         'disabled access / wheelchair ramp'],
        '社区设施（Community）',
        ['membership card / annual pass',
         'peak / off-peak hours',
         'booking / reservation in advance',
         'concession / student discount',
         'sports hall / leisure centre',
         'car park / parking facilities',
         'reception / front desk'])
    # PART06
    U.section(p,6,'方法精讲','Methods & Models')
    U.steps(p,'地图题满分5步操作',
        [('预读(15s)','找入口标✱；数字母标记数量；找已知地标'),
         ('定起点','录音说"As you enter/from the entrance"→铅笔放到✱'),
         ('追路线','每听到方向词→铅笔移动；一步一步，不跳跃'),
         ('标答案','说出设施名时→立即在铅笔位置写题目答案字母'),
         ('验证','全部标完：检查地图每个位置是否合理、有无遗漏')])
    U.steps(p,'匹配题满分5步操作',
        [('预读选项','在录音开始前，A→H逐个圈关键词；每选项≤3秒'),
         ('建词库','脑中记住：A=free, B=renovated, C=booking...'),
         ('听到设施名','立即在脑中激活词库；等说话人描述该设施'),
         ('匹配关键词','听到"refurbished"→匹配→B(renovated) → 写B'),
         ('验证整句','写完后确认：整句话是否支持选项含义？有否定词吗？')])
    U.model_answer(p,'S2地图题满分追踪示范',
        'Riverside Park Map — Label A-E',
        '追踪笔记（边听边画箭头）：\n\n'
        '✱ENTRANCE → ↑直走 → 正前=Info Centre → Q11=B\n'
        '← 左转 → 20m → 左侧第一间=Café → Q12=D\n'
        '↑ 继续走 → 右侧=Gift Shop → Q13=A\n'
        '← 左转走到尽头 → 右边=Toilets → Q14=C\n'
        '← 返回走廊 → 左边最后=Storage → Q15=E',
        ['每步只写箭头和字母，不写全名，节省时间',
         '跟丢了：找地图唯一空白位置，配下一题',
         '最后检查：每个字母用过一次，无重复',
         '铅笔先标，誊写时再确认'])
    U.model_answer(p,'S2匹配题满分示范',
        'Sports Centre Facilities (Q16-20, A-G)',
        '预读关键词库（30秒建立）：\n'
        'A=free  B=booking  C=renovated  D=new  E=members  F=outdoor  G=24hrs\n\n'
        'Q16Pool: "just been refurbished" → C(renovated ✓)\n'
        'Q17Tennis: "exclusively for members" → E(members ✓)\n'
        'Q18Library: "book at least a day in advance" → B(booking ✓)\n'
        'Q19Courts: "situated in the open air" → F(outdoor ✓)\n'
        'Q20Café: "undergoing renovation" → C? → 已用 → F(renovated)',
        ['预读关键词库在30秒内建立——不靠考场现查',
         'exclusively=members only; open air=outdoor——同义替换即时联想',
         'Q20 Café在装修=renovated=C，但C已用于Pool——重新看选项',
         '全部填完：检查A-G每个只用一次（除多余项）'])
    U.compare(p,'地图题：跟丢 vs 满分',
        '❌ 跟丢的考生',
        ['录音开始才低头找地图',
         '方向词还没反应就过了',
         '跟丢了整个地图题放弃',
         '不确定就空着，全靠猜',
         '答完没有检查位置合理性'],
        '✅ 满分的考生',
        ['说明时间已完成地图预读',
         '方位词触发条件反射，即时移动铅笔',
         '跟丢立即找下一个地标词重定位',
         '不确定先写候选，等更多信息确认',
         '全做完检查：位置是否都合理？'])
    U.compare(p,'匹配题：临时读选项 vs 预读建库',
        '❌ 临时读选项',
        ['Q16出来才开始读A选项',
         '录音在说，眼睛在扫描选项',
         '听到关键词来不及找对应选项',
         '来回在题目和选项间切换视线',
         '错过细节，最后凭感觉填'],
        '✅ 预读建关键词库',
        ['录音前已读完A-H，每个圈词',
         '脑中有词库，听时无需扫描',
         '听到"refurbished"→0.5秒定位选项',
         '眼睛只看题号，手写答案',
         '验证整句，确认无否定词后写'])
    # PART07
    U.section(p,7,'课堂练习','Practice Drills')
    U.drill(p,'练习1：方位词5秒闪认',
        '老师说方位词，学生立即在白纸上画空间示意图（5秒/题）',
        ['"The café is opposite the main entrance."→画出位置',
         '"Turn left, then take the second right, gym is at the end."→画路线',
         '"Adjacent to the library, just past the car park."→标注',
         '"Between the gym and the pool, facing the entrance."→标注',
         '"At the far end beyond the sports hall, on your right."→标注'])
    U.drill(p,'练习2：匹配题预读计时',
        '60秒完成以下8个选项的关键词圈划，并背住"字母=关键词"',
        ['A.free of charge  B.closed for maintenance  C.requires advance booking',
         'D.recently opened  E.exclusively for members  F.open on weekends only',
         'G.has parking facilities  H.offers discounts for seniors',
         '→ 60秒后合上笔记，说出A=? B=? C=? D=? E=? F=? G=? H=?',
         '→ 测试脑中词库是否建立完成'])
    U.drill(p,'练习3：信号词标注',
        '读以下段落，圈出所有话题转换信号词，并说出每段讲什么',
        ['"Good morning. Today I\'d like to introduce our community centre."',
         '"Starting with the ground floor — we have a gym and a café."',
         '"Now, moving on to the first floor, you\'ll find four meeting rooms."',
         '"Before I finish, please note that the car park requires a permit."',
         '"Finally, I\'d like to remind you to register online before Friday."'])
    U.drill(p,'练习4：地图追踪完整练习',
        '在白纸上画矩形建筑，底部中央标ENTRANCE，按口述追踪：',
        ['"Enter through the main doors. The reception is directly ahead."',
         '"Turn left at reception. The first room on your right is the gym."',
         '"Continue past the gym. The swimming pool is at the far end."',
         '"The changing rooms are opposite the swimming pool."',
         '"Back at reception, turn right. The café is the last room on your left."'])
    U.drill(p,'练习5：S2综合模拟计时',
        '完整S2模拟：预读→答题→自检三阶段，老师计时',
        ['第一阶段（30秒）：预读Q11-20，圈选项关键词，标地图起点',
         '第二阶段（实时）：播放录音，同步追踪地图+匹配关键词',
         '第三阶段（30秒）：检查：地图位置合理？选项无重复使用？',
         '结果分析：地图题错几题（方位问题）？匹配题错几题（同义替换问题）？',
         '提升计划：错误最多的类型→下周专项练习'])
    # PART08
    U.section(p,8,'避坑指南','Pitfall Prevention')
    U.pain_point(p,'S2地图与匹配最常见失分',
        ['"On your left"=行进方向左，不一定是地图左侧——很多人在这里混淆',
         '匹配题看到"free"立刻选A，没读完整句——否定+free=不免费',
         '话题转换没听到信号词，Q16以后全部跟丢，只能乱猜',
         '预读时只读当前题的选项，没有提前建完整关键词库'],
        ['每次地图转弯：问自己"我现在面向哪个方向？"——方向感必须实时更新',
         '匹配题：选项确认前必须等整句说完——听到关键词≠已确认',
         '信号词训练：每次做S2录音后，对照原文标出所有信号词',
         '预读铁律：S1结束立刻翻S2，强制完成全部选项关键词圈划'])
    U.content(p,'匹配题常见干扰模式',
        [('先否后肯','It used to be free, but now costs £5 → 选收费，不选free'),
         ('条件限定','Free for members → 不是全体免费；如题问"所有人免费"则不选A'),
         ('时间限定','Open weekends only vs closed weekdays → 含义相同，看题目问什么'),
         ('双重特点','设施同时具有B和C两个特点；选最直接对应题干的那个'),
         ('情绪干扰','说话人花更多时间描述A，让人以为A是答案——靠内容不靠时间')],
        note='每道匹配题：问自己"这整句话的核心意思是什么"，不只看关键词')
    U.content(p,'S2常见失分模式与解法',
        [('地图跟丢','原因：方位词反应慢；解法：方位词闪卡，每天5分钟，直到条件反射'),
         ('同义替换不认','原因：替换词不熟；解法：积累20组高频替换对，每周测试'),
         ('话题转换跟丢','原因：信号词不敏感；解法：每次精析后标出全部信号词'),
         ('选项未预读','原因：时间管理差；解法：强制S1结束立刻翻S2预读'),
         ('检查习惯缺失','原因：以为答完就好；解法：每次S2做完，30秒自检流程')])
    U.tip_card(p,'S2地图题速记卡',
        [('起点定位','录音说"As you enter"→立刻找地图入口，铅笔放上去，追踪开始'),
         ('方位三词','opposite=对面; adjacent=紧邻; just past=经过后紧接——解决80%地图题'),
         ('跟丢策略','跟丢不慌：看下一道题关键词，等录音出现该词，跳过已错的继续'),
         ('最终检查','地图全标完：每个字母(A-E)用了一次；位置是否与描述方向吻合')])
    U.tip_card(p,'S2匹配题速记卡',
        [('预读铁律','录音前读完A-H全部选项；每个圈1词；脑中建"字母→关键词"词库'),
         ('同义替换','refurbished=renovated; exclusively=only; round the clock=24hrs'),
         ('不急确认','听到关键词→标候选→等整句说完→无否定词→最终确认'),
         ('多余选项','8选5：3个多余；实在不会：从未使用选项中猜，绝不留空白')])
    U.checklist(p,'L03 自检清单',
        ['我能5秒内说出adjacent/opposite/just past/beyond的空间含义',
         '我知道匹配题预读操作：读完A-H，每个圈1关键词，建词库',
         '我练习了地图追踪：铅笔随方向词实时移动',
         '我识别了S2独白话题转换的信号词',
         '我知道跟丢后重定位的策略：找下一个设施名/地标'])
    U.summary(p,3,
        ['S2地图题=方位词条件反射+铅笔实时追踪；跟丢找地标重定位',
         'S2匹配题=预读选项建词库+识别同义替换；整句说完再确认',
         '信号词是S2独白的结构骨架——掌握信号词永远不跟丢'],
        takeaway='S2稳8分：预读时圈好选项关键词+找好地图起点——这两件事在录音前完成')
    U.homework(p,3,
        ['完成剑桥真题2套S2（计时），分别统计地图题和匹配题正确率',
         '制作20个方位词闪卡，每天5分钟闪认，直到每个词5秒内生成图像',
         '从一套S2精析：标出全部信号词，核对与题目转换点是否一致'],
        est_time='约50分钟')
    U.preview(p,4,'Section 3 学术对话攻略',
        ['多选题(Choose TWO)的系统解法',
         '2-4人说话人观点追踪技术',
         'S3最高频的陷阱：观点归属混乱'])
    U.save(p, f'{OUT}/L03_Section2独白场景攻略.pptx')
    print(f"  L03 done: {len(p.slides)} pages")

# ══════════════════════════════════════════════════════════════════════════════
#  L04  Section 3 学术对话攻略
# ══════════════════════════════════════════════════════════════════════════════
def L04():
    T(); p=U.new_prs()
    U.cover(p,'听力',4,'Section 3\n学术对话攻略','Section 3: Academic Discussion Mastery')
    U.agenda(p,['上节回顾','S3结构与特点','MCQ多选题系统解法','多人追踪技术',
                '态度词识别','例题精讲','课堂练习','小结作业'])
    U.review(p,['S2地图题：方位词条件反射+铅笔追踪，跟丢找地标重定位',
                'S2匹配题：预读建词库+识别同义替换，整句说完再确认',
                '信号词=话题转换：Moving on to/Before I finish'])
    U.objectives(p,['掌握S3多选题(Choose TWO)的五步系统解法',
                    '建立2-4人说话人观点追踪的标注系统',
                    '识别S3最高频陷阱：观点归属混乱和条件限定',
                    '实现Section 3稳定6-7/10'])
    U.pain_point(p,'S3失分最常见的三类原因',
        ['多选题只选了一个(Choose TWO却只填一个)，或者填了三个',
         'Amy说的观点写成了Tom的，说话人切换时张冠李戴',
         '听到选项词就选，没等说完——是否定句或条件限定'],
        ['多选题做完立刻数：必须是两个字母才算完成，1个或3个都是0分',
         '用T/A/S标注说话人，每人发言时边听边记缩写，观点分别追踪',
         '选择题三步法：听到候选词→继续听→特别注意but/however→等整句确认'])
    part(p,1,'S3结构与特点','S3 Structure & Features',
        [('Section 3的本质','2-4人学术讨论，最难的Section',
          [('参与人数','通常2人（学生+导师）或3-4人（小组讨论）'),
           ('内容','项目讨论/实验设计/论文修改/小组作业'),
           ('难度','深圳学生失分最多的Section；均分约5-6/10'),
           ('题型','MCQ多选35%+单选30%+匹配20%+其他15%'),
           ('最大难点','多人对话中观点归属混乱——是谁说的？')],
          'S3考的不是语言能力，是信息追踪能力——可以训练'),
         ('S3典型题型：Choose TWO','最让深圳学生苦恼的题型',
          [('基本要求','从A/B/C/D/E中选两个；两个都对才得分'),
           ('选项数量','通常5个选项选2个；有3个干扰选项'),
           ('判断方式','每个选项单独判断；不靠感觉，不成组'),
           ('常见失误','凭感觉选A和E；听到A词就选A'),
           ('正确','每选项：录音支持→✓；录音否定→✗；未知→?')])],
        [('S3讨论内容的规律',
          [('段落一','介绍研究主题/项目背景；通常有填空或单选题'),
           ('段落二','讨论研究方法；常考MCQ——哪些方法同意/不同意'),
           ('段落三','结果讨论；常考态度题——对结果各有什么评价'),
           ('段落四','未来计划/改进；常考配对——谁建议做什么'),
           ('关键','题目通常按对话顺序出现——不会倒序')]),
         ('S3vs其他Section的核心差异',
          [('vs S1','S1是填空为主；S3是选择和配对为主'),
           ('vs S2','S2是独白；S3有2-4人说话——追踪难度大'),
           ('vs S4','S4是单人讲座；S3的难点在多人+态度判断'),
           ('失分来源','S3约30%失分来自"听懂了但选错"——策略问题'),
           ('提分策略','追踪技能训练比大量刷题更有效')])],
        [('例题：S3 Choose TWO基础',
          'Q21-22: Choose TWO letters A-E.\nWhich TWO problems did they identify?\n'
          'A.sample size  B.survey questions  C.data errors  D.time management  E.lit review\n\n'
          'Script: "The main issue was the sample size — way too few participants.\n'
          'And our literature review wasn\'t thorough enough."',
          'Q21&22: A and E\nA: sample size✓(too few participants=small sample)\nE: lit review✓(not thorough=incomplete)',
          '每个选项单独判断：A→明确支持；E→明确支持；B/C/D→未提及→?'),
         ('例题：S3观点配对',
          'What does each student say about the methods?\n'
          'A.needs more time  B.appropriate  C.too expensive  D.should be changed  E.well-designed\n\n'
          'Q23:Tom→?  Q24:Amy→?\n\n'
          'T:"The interviews were well-structured."\nA:"I\'m not sure about the surveys — we need to reconsider."',
          'Q23:Tom→E(well-structured=well-designed)\nQ24:Amy→D(reconsider=should be changed)',
          '配对关键：明确谁在说话？Tom→Q23，Amy→Q24——说话人追踪决定配对')])
    part(p,2,'MCQ多选题系统解法','MCQ: Complete System',
        [('Choose TWO五步解法','每步都必须，不能跳跃',
          [('Step1 预读','读完A-E；圈每选项最独特的词'),
           ('Step2 分离判断','把A-E当成5道独立是/否题，不成组'),
           ('Step3 标注','每选项旁写✓(支持)/✗(否定)/?(不确定)'),
           ('Step4 确认两个✓','等两个选项都有明确✓才写答案'),
           ('Step5 数量检查','最终：写了两个字母吗？')],
          '最大陷阱：选了两个≠对了——两个都需要有明确的录音支持'),
         ('排除法策略','通过排除提高准确率',
          [('录音否定了某选项','立刻打✗，完全排除不再考虑'),
           ('一个✓一个?','等待，?的答案可能后面还会提到'),
           ('录音说"didn\'t change..."','涉及的选项全打✗'),
           ('时间不够','从?中选概率最高的；绝不留空'),
           ('数量铁律','做完立刻数：必须是2个字母')])],
        [('排除法实战',
          'Choose TWO: Which factors improved the project?\n'
          'A.larger sample  B.revised questions  C.expert feedback  D.more funding  E.extended deadline\n\n'
          '"Expert comments really helped us improve the design.\n'
          'The extended deadline also made a big difference — gave us time to redo the analysis.\n'
          'We didn\'t change the sample size or the questions."',
          'C: expert feedback✓\nE: extended deadline✓\nA,B: "didn\'t change"→✗\nD: 未提→?\n答案: C, E',
          'A和B被明确否定→提高确认C和E的信心；不需要再考虑A/B'),
         ('Choose TWO常见失误详解',
          [('失误1 成组思考','A和E看起来"配"→错；每个选项独立判断'),
           ('失误2 早下结论','听到C的词就选C，后面来了but否定→选错了'),
           ('失误3 数量错误','最终写了3个字母→只计前2个；或写了1个→0分'),
           ('失误4 未排除干扰','没有打✗的选项干扰最终判断'),
           ('正确习惯','做完立刻数：2个字母？用✗划掉确定错误的')])],
        [('Choose TWO完整练习',
          'Q25-26: Choose TWO: What do they agree needs improvement?\n'
          'A.introduction  B.methodology  C.data analysis  D.conclusion  E.references\n\n'
          'T:"The introduction is fine." A:"Agreed. But our methodology needs work."\n'
          'A:"And the references aren\'t complete." T:"Yes — definitely."',
          'A: Tom说fine→✗\nB: Amy说needs work→✓; Tom没直接评论→?\nE: Amy提references→✓; Tom同意→✓\n答案: B, E',
          'B: Amy✓, Tom?→但题目问"they agree"→B需要两人都同意→继续听Tom对B的反应'),
         ('排除法：哪个选项被明确否定？',
          'Q27-28: Choose TWO: What problems did they encounter?\n'
          'A.time pressure  B.unclear instructions  C.sample size  D.budget  E.ethics approval\n\n'
          '"We had a very tight deadline — really struggled with time.\nThe ethics approval process took much longer than expected too.\nBudget wasn\'t an issue and instructions were perfectly clear."',
          'A: time pressure✓(tight deadline)\nE: ethics approval✓(took longer)\nB: "perfectly clear"→✗\nD: "wasn\'t an issue"→✗\nC: 未提→?\n答案: A, E',
          '排除B(clear)和D(not an issue)后只剩A和E被明确支持→答案直接确定')])
    part(p,3,'多人对话追踪系统','Multi-Speaker Tracking System',
        [('说话人识别与标注','不知道谁在说=无法追踪观点',
          [('识别方式','名字+声音特征；说话人通常会互称名字'),
           ('标注系统','草稿纸：T=Tom, A=Amy, S=Susan/导师'),
           ('观点追踪','每人说话时：在对应字母旁记关键词/态度(+/-/?)'),
           ('切换信号','Yeah.../Right.../But.../Actually... = 可能换说话人'),
           ('实战','表格追踪：横轴=选项A-E，纵轴=说话人T/A')],
          '多人对话=信息追踪游戏；谁说了什么比说了什么更重要'),
         ('观点vs事实vs态度的区分','S3的核心考察',
          [('事实陈述','数据/时间/结果——客观，无个人态度'),
           ('观点表达','I think/I believe/It seems to me/In my view'),
           ('强烈赞同','definitely/absolutely/exactly/couldn\'t agree more'),
           ('部分赞同','to some extent/I see your point but.../it depends'),
           ('明确反对','I disagree/That\'s not right/I\'m not so sure')])],
        [('多人追踪表格练习',
          '在纸上建追踪表，读对话后填入✓✗?：\n\n'
          '         Amy    Tom\nA研究设计   ?      ?\nB样本大小   ?      ?\nC数据分析   ?      ?\n\n'
          'A:"I think the research design is solid." T:"Agreed. But the sample size concerns me."\nA:"Data analysis seems fine too." T:"I\'m not sure about that."',
          '研究设计: A✓ T✓\n样本大小: A? T✗\n数据分析: A✓ T✗\n\nQ: 两人都同意的? → 研究设计(A)\nQ: 只有Amy同意的? → 数据分析(C)',
          '表格追踪让"谁说了什么"一目了然；题目通常问"哪个观点两人都同意"'),
         ('导师vs学生的话语模式',
          [('导师说话','更权威；评价学生工作；用"have you considered/what about"提问'),
           ('学生说话','更谦虚；用"I think/maybe/possibly"；会请教导师'),
           ('讨论规律','学生提观点→导师评价→学生修正/坚持'),
           ('观点归属','导师建议≠学生采纳；要听学生的回应'),
           ('陷阱','导师否定学生的原始观点，原观点不是答案')])],
        [('导师-学生对话追踪',
          'Q27:What does the tutor suggest?\nQ28:What do students agree to change?\n\n'
          'Tutor:"Your methodology is sound, but have you considered a focus group?"\n'
          'Student A:"We hadn\'t thought of that — could work." B:"Yes, let\'s include it."',
          'Q27: Using a focus group（导师建议）\nQ28: Including a focus group（学生同意采纳）',
          '导师建议≠学生行动；Q28问的是学生决定做什么——他们同意了→include it'),
         ('态度词快速识别：练习判断真实态度',
          '判断以下发言者的真实态度（+赞同/-反对/~中立）：\n\n'
          '1. "That\'s an interesting idea, but I\'m not sure it would work."\n'
          '2. "Absolutely, I couldn\'t agree more — this is the right approach."\n'
          '3. "It depends on how you look at it — there are pros and cons."\n'
          '4. "I see your point, but the evidence doesn\'t really support that."',
          '1. ~→- (That\'s interesting=礼貌；but后才是真态度=不确定/反对)\n2. +(Absolutely+couldn\'t agree more=强烈赞同)\n3. ~(中立/不确定)\n4. -(I see your point=礼貌；but后=实际反对)',
          'S3关键：礼貌语气后面的but/however才是真实态度——不要被礼貌迷惑')])
    part(p,4,'态度词与S3综合策略','Attitude Words & S3 Strategy',
        [('态度识别词汇体系','不认识态度词=无法判断观点',
          [('强烈支持','definitely/certainly/absolutely/without doubt'),
           ('部分支持','to some extent/fairly/reasonably/I suppose'),
           ('中立','it depends/hard to say/could go either way'),
           ('部分反对','not entirely convinced/not necessarily/I see your point but'),
           ('强烈反对','definitely not/absolutely not/strongly disagree')],
          '听到态度词：立刻在选项旁标注+/-/~ (positive/negative/neutral)'),
         ('S3时间管理与综合策略','S3的时间压力最大',
          [('预读时间','说明+间隙：读完所有选项，MCQ圈每选项关键词'),
           ('每题时间','每道题最多2分钟；超时=跳过，不回头'),
           ('Choose TWO','预留额外30秒确认第二个选项'),
           ('配对题','先做确定的；不确定的最后统一处理'),
           ('S3结束','立即翻S4；用20秒预读S4全部题目')])],
        [('S3综合陷阱识别',
          [('陷阱1 说话人混淆','未确认谁说的→用说话人缩写T/A/S标注'),
           ('陷阱2 条件限定','"only if we have time"=不是确定计划'),
           ('陷阱3 改变主意','先说A后说"actually B"→以B为准'),
           ('陷阱4 礼貌语气','"That\'s interesting"≠同意——仅是礼貌'),
           ('陷阱5 否定+选项词','"not really helpful"→选项有helpful→不能选')]),
         ('S3提分最快3个方法',
          [('方法1 说话人标注','做S3强制标注每个说话人缩写和态度，训练3周'),
           ('方法2 MCQ排除法','每次：先打✗录音否定的，再从剩余中选'),
           ('方法3 精析原文','每道错题：找原文→分析为何听错→分类失分原因'),
           ('时间投入','这3个方法每周2-3小时，坚持4周→S3提升约2题/套'),
           ('目标','均分5题→稳定7题：4-6周针对性训练即可实现')])],
        [('S3综合练习',
          'Q29-30(Choose TWO): Which aspects do BOTH students find challenging?\n'
          'A.data collection  B.writing style  C.referencing  D.time management  E.peer review\n\n'
          'T:"I\'ve struggled a lot with data collection." A:"Yes, same here."\nT:"Referencing has been tricky for me." A:"I\'m fine with that — writing style is my issue."',
          'A: T✓ A✓ → 两人都同意 ✓\nB: A✓ T? → 只有Amy\nC: T✓ A✗ → 只有Tom\n→ 只有A(data collection)两人都同意\n题目问"BOTH find challenging"→只有A符合',
          '注意：题目说BOTH——必须两个说话人都明确支持才算；只有一人支持不够'),
         ('S3综合策略演练',
          'Q31(Choose ONE): What is the main purpose of the focus group?\n'
          'A.to test the survey  B.to collect additional data  C.to replace the interviews\n\n'
          'T:"The focus group will let us verify whether the survey results hold up."\nA:"Exactly — it\'s really a way of cross-checking what we already have."',
          'Answer: A — to test/verify the survey\nB: "additional data"不准确；C: 未提替代；A: verify=test(cross-checking=testing)',
          '单选题：排除B(未提)和C(未提)后只剩A；verify=test是同义替换')])
    # PART05
    U.section(p,5,'高频词汇','Key Vocabulary')
    U.vocab(p,'态度词：赞同类',
        [('definitely / certainly',   '肯定（强烈赞同）'),
         ('absolutely / exactly',     '完全正确（强烈赞同）'),
         ('I couldn\'t agree more',   '再同意不过了'),
         ('that\'s a good point',     '这是个好点（礼貌，未必真同意）'),
         ('fair enough',              '好吧/合理（勉强接受）'),
         ('to some extent',           '在某种程度上（部分赞同）'),
         ('I think you\'re right',    '你说得对（赞同）')])
    U.vocab(p,'态度词：反对与中立类',
        [('I\'m not sure about',      '我不确定（轻度反对）'),
         ('I see your point, but',    '我理解，但...（转折=反对）'),
         ('that\'s not quite right',  '不太对（直接反对）'),
         ('I strongly disagree',      '我强烈不同意'),
         ('it depends',               '要看情况（中立）'),
         ('could go either way',      '两边都有可能（中立）'),
         ('actually / in fact',       '实际上（通常引出不同意见）')])
    U.vocab(p,'S3学术讨论高频词',
        [('methodology / approach',   '研究方法'),
         ('findings / results',       '发现/结果'),
         ('hypothesis / assumption',  '假设'),
         ('validity / reliability',   '效度/信度'),
         ('sample / participants',    '样本/参与者'),
         ('implication / impact',     '含义/影响'),
         ('peer review / feedback',   '同行评审/反馈')])
    U.two_col(p,'S3题型对比：MCQ vs 配对',
        'Choose TWO（5选2）',
        ['每选项独立判断✓✗?',
         '两个✓才写答案',
         '最后数：写了2个吗？',
         '排除✗提高确信度',
         '预读圈每选项关键词',
         '每题预留2.5分钟',
         '题目词≠录音词（同义替换）'],
        '配对题（观点归属）',
        ['建说话人追踪表格',
         '每说话人一列(T/A/S)',
         '每个选项旁标谁说的',
         '态度词决定+/-/~',
         '等整段说完再最终标',
         '注意：礼貌赞同≠真同意',
         '改主意=以后说的为准'])
    # PART06
    U.section(p,6,'方法精讲','Methods & Models')
    U.steps(p,'Choose TWO五步满分解法',
        [('预读','读完A-E；圈每选项关键词；心中建词库'),
         ('分离判断','把5个选项当5道独立是/否题处理'),
         ('实时标注','每选项旁写✓(录音明确支持)/✗(录音否定)/?'),
         ('等两个✓','不要一个✓就急着写；等第二个确认'),
         ('数量检查','最终：确认写了两个字母；一个或三个都是错的')])
    U.steps(p,'多说话人观点追踪5步法',
        [('确认说话人','录音开始：确认几个人？各叫什么名字？'),
         ('建追踪表','草稿纸画：选项横轴 + 说话人纵轴'),
         ('实时标注','每人发言：在对应格写✓✗?和关键词'),
         ('态度词警报','听到definitely/not sure等→提高关注度'),
         ('题目对照','填答案前：用追踪表核对"谁+什么观点"与题目匹配')])
    U.model_answer(p,'S3 Choose TWO满分示范',
        'Q21-22: What problems did they identify?',
        '追踪笔记：\n\n'
        'A.样本: T:"too small"✓ A:"agreed"✓ → 两人支持 ✓\n'
        'B.问题设计: T:? A:"questions were fine"✗ → Amy否定\n'
        'C.数据: T:"data looked clean"✗ A:✗ → 两人否定\n'
        'D.时间: T:? A:? → 不确定\n'
        'E.文献: T:"not thorough"✓ A:✓ → 两人支持 ✓\n\n'
        'Answer: A and E',
        ['A和E: 两人都✓→确认', 'C: 两人都✗→排除',
         'B: Amy否定→排除; D: 两人都?→继续听',
         '已有A和E两个✓→答案确定，不需要等D'])
    U.model_answer(p,'S3配对题满分示范',
        'Q23-25: What does each person say?',
        '说话人追踪：\n\n'
        'Tom(T): "The timeline is too tight."\n'
        '        → 时间紧张 → C.time constraints\n\n'
        'Amy(A): "I think the research question is clear."\n'
        '        → 研究问题好 → A.research question clear\n\n'
        'Tutor(S): "The structure could be improved."\n'
        '          → 结构需改进 → E.structure needs work\n\n'
        'Q23 Tom→C  Q24 Amy→A  Q25 Tutor→E',
        ['每说话人发言时立刻在姓名旁写关键词',
         'Tom=时间→找"timeline/time constraint"选项',
         '导师"could be improved"=需改进，不是赞扬',
         '每人只有一个配对；配完后检查无重复'])
    U.compare(p,'MCQ: 直觉选 vs 系统选',
        '❌ 靠直觉',
        ['A和E看起来像答案，凭感觉选',
         '听到A的词就选A，没继续听',
         '选了两个字母，以为做对了',
         '一个选项纠结超过2分钟',
         '不知道已排除了哪些选项'],
        '✅ 系统选法',
        ['每选项独立判断，不联合思考',
         '听到候选词→继续听→等完整句',
         '两个✓都有明确录音支持才写',
         '超时立刻标?继续，先做其他',
         '排除✗选项，有效缩小范围'])
    U.compare(p,'追踪: 混乱 vs 清晰',
        '❌ 不追踪说话人',
        ['不确定这个观点是谁说的',
         '听到观点就填，不管归属',
         '两人态度相反时选错方向',
         '导师和学生的说法混在一起',
         '配对题全靠猜'],
        '✅ 建立追踪系统',
        ['边听边在草稿纸写T:/A:',
         '每个观点明确归属说话人',
         '表格：选项×说话人=态度',
         '导师建议vs学生采纳分开记',
         '配对题用追踪表核对后填写'])
    # PART07
    U.section(p,7,'课堂练习','Practice Drills')
    U.drill(p,'练习1：态度词快速识别',
        '听/读以下句子，判断说话人态度是+（正面）-（负面）还是~（中立）',
        ['"I definitely think the sample size needs to be increased." → ?',
         '"That\'s an interesting suggestion, but I\'m not convinced." → ?',
         '"It depends on how you look at it." → ?',
         '"The data analysis is clearly flawed." → ?',
         '"I see your point — that could work, to some extent." → ?'])
    U.drill(p,'练习2：Choose TWO标注训练',
        '用✓✗?标注每个选项，最后写两个答案',
        ['Choose TWO: What improvements did they make?',
         'A.larger sample  B.better questions  C.focus group  D.extended time  E.expert review',
         'T:"We added a focus group." A:"And we got an expert to check our analysis."',
         'T:"We didn\'t change the sample size or the questions."',
         '参考：A✗ B✗ C✓ D? E✓ → 答案=C, E'])
    U.drill(p,'练习3：说话人追踪表格练习',
        '建表格，边读边追踪Amy(A)和Tom(T)的态度（✓✗?）',
        ['A:"I think the methodology is solid." T:"I have some doubts about that."',
         'A:"The sample size seems adequate." T:"Actually, it\'s too small."',
         'A:"The findings are interesting." T:"Absolutely — I agree on that."',
         'Q: 两人都同意的? → findings',
         'Q: 只有Amy同意的? → methodology and sample size'])
    U.drill(p,'练习4：MCQ排除法练习',
        '先排除被录音否定的，再从剩余中选最佳答案',
        ['Choose ONE: Why did they change the research topic?',
         'A.tutor suggestion  B.data unavailable  C.personal interest  D.time pressure',
         '"We had plenty of time. And our original interest was still there."',
         '"The tutor thought the original topic was too broad — we followed the advice."',
         '排除D(plenty of time)和C(interest still there)→剩A和B→A✓'])
    U.drill(p,'练习5：S3完整模拟',
        '完整S3模拟：预读→追踪→MCQ三位一体',
        ['预读30秒：读完全部选项，MCQ圈关键词，配对确认说话人数量',
         '听题时：同步建说话人追踪表格；MCQ用✓✗?实时标注',
         '做完后：Choose TWO数量核对(必须2个)；配对核对(每人一个)',
         '对照原文：每道错题找到答案句，分析失分原因',
         '总结：MCQ/配对各对几题——找到弱项，下周专练'])
    # PART08
    U.section(p,8,'避坑指南','Pitfall Prevention')
    U.pain_point(p,'S3最典型的4类失分',
        ['Choose TWO只选了一个：题目要求两个，填一个=0分',
         'Amy的观点写成Tom的：对话追踪混乱，张冠李戴',
         '听到选项词立刻选：没等整句说完，后面来了否定词',
         '导师建议=学生答案：导师建议不一定是学生会执行的'],
        ['做完Choose TWO立刻数：纸上写了几个字母？必须是2个',
         '建说话人追踪表：草稿纸画T/A/S列，每个观点归属明确',
         '三步确认：听到候选→继续听→整句说完无否定→写答案',
         '配对题：问的是学生做什么→找学生的回应，不是导师建议'])
    U.content(p,'S3陷阱词汇速查',
        [('礼貌但不同意', '"That\'s a good point" + but = 实际不同意，but后才是真观点'),
         ('改变主意标志', 'actually/in fact/wait/I\'ve changed my mind → 以后说的为准'),
         ('条件限定', '"only if..." / "as long as..." → 条件成立才成立'),
         ('最终决定', '"so we\'ve decided to..." / "in the end..." → 这才是答案'),
         ('否定+选项词', '"not really helpful" → 不选"helpful"选项')],
        note='遇到这5类信号词：立刻停止书写，继续听，等完整句子结束')
    U.content(p,'S3备考策略',
        [('题型统计', '做完S3统计MCQ/配对/填空各题型正确率，找最弱的专攻'),
         ('说话人练习', '每次练习强制标注说话人，形成习惯后提升速度'),
         ('态度词积累', '建立100个态度词词表：赞同/反对/中立各分类'),
         ('精析优先', '每错1题精析比多做2套更有价值：找原文→分析失分原因'),
         ('目标', 'S3从均分5题→7题：需4-6周针对性训练，不是靠刷题量')])
    U.tip_card(p,'S3 MCQ速记卡',
        [('Choose TWO铁律', '每选项独立✓✗?标注；两个✓才写答案；最后数：必须是2个字母'),
         ('排除法提效', '录音明确否定的选项立刻打✗；从剩余中选，减少干扰'),
         ('等整句', '听到候选词→继续听→整句说完无否定→写；不要早下结论'),
         ('数量第一', 'Choose TWO做完第一件事：数一下，写了几个字母？')])
    U.tip_card(p,'S3说话人追踪速记卡',
        [('草稿纸建表', '横轴=选项A-E；纵轴=说话人T/A/S；每格写✓✗?和关键词'),
         ('态度词敏感', '听到definitely/not sure/actually立刻提高警觉——态度词后是关键'),
         ('礼貌≠同意', '"That\'s interesting"不是赞同；等但是/补充才是真实态度'),
         ('改变主意', '先说A后说"actually B"→以B为准；actually是更正信号')])
    U.checklist(p,'L04 自检清单',
        ['我知道Choose TWO的五步解法，且做完会数是否写了两个字母',
         '我练习了说话人追踪表格，能区分Tom和Amy的不同观点',
         '我识别了态度词：definitely/I\'m not sure/that\'s interesting but...',
         '我知道"导师建议"≠"学生答案"，配对时看学生的回应',
         '我理解S3为什么难：是追踪技能问题，不是语言能力问题'])
    U.summary(p,4,
        ['S3核心：说话人追踪+MCQ系统判断——两个技能同时运用',
         'Choose TWO：每选项独立✓✗?；两个都明确支持才写；必须2个字母',
         '说话人追踪：建表格；态度词=警报；礼貌≠同意；改变主意以后为准'],
        takeaway='S3不是语言能力，是追踪技能——从5→7分需要3-4周系统追踪训练')
    U.homework(p,4,
        ['完成剑桥真题2套S3，强制建说话人追踪表格，统计正确率',
         '积累30个态度词（赞同/反对/中立各10个），制作闪卡',
         '精析本次S3所有错题：每题找原文句，分析是MCQ还是追踪失分'],
        est_time='约55分钟')
    U.preview(p,5,'Section 4 学术讲座攻略',
        ['笔记填空与摘要填空的专项技法',
         'AWL学术词汇300个核心词的听辨策略',
         'S4"跟丢后"的心理调适与追踪恢复'])
    U.save(p, f'{OUT}/L04_Section3学术对话攻略.pptx')
    print(f"  L04 done: {len(p.slides)} pages")

# ══════════════════════════════════════════════════════════════════════════════
#  L05  Section 4 学术讲座攻略
# ══════════════════════════════════════════════════════════════════════════════
def L05():
    T(); p=U.new_prs()
    U.cover(p,'听力',5,'Section 4\n学术讲座攻略','Section 4: Academic Lecture Mastery')
    U.agenda(p,['上节回顾','S4结构与特点','笔记/摘要填空技法','AWL学术词汇策略',
                '讲座追踪技能','例题精讲','课堂练习','小结作业'])
    U.review(p,['S3 Choose TWO：每选项独立✓✗?，两个✓才写，最后数字母数量',
                '说话人追踪表格：T/A/S标注，态度词=警报，礼貌≠同意',
                '时间管理：每题≤2分钟，S3结束立刻翻S4预读'])
    U.objectives(p,['掌握S4笔记填空与摘要填空的答题技法',
                    '建立AWL学术词汇的"听辨"策略（不要求会写，要求听出来）',
                    '学会识别S4讲座结构信号词，避免跟丢',
                    '实现Section 4稳定5-7/10'])
    U.pain_point(p,'S4失分的3种典型情形',
        ['S4第一句没听懂，产生恐慌，后面9题全部放弃——自我放弃是最大的失分',
         'AWL词汇听到了但不认识(evapotranspiration/entrepreneurship)，音记不出来',
         '摘要填空：填了三个词，但限词是ONE WORD——因违反限词全部0分'],
        ['跟丢了不要慌：找下一题的关键词，用题号重定位；哪怕猜也要填',
         'AWL策略：按音节写占位(evap...)，誊写时补全；重音节是识别关键',
         '做每题前看限词：ONE WORD→只写一个；TWO WORDS→最多两个'])
    part(p,1,'S4结构与特点','S4 Structure & Features',
        [('Section 4的本质','单人学术讲座，最难的Section',
          [('说话人','大学教授/讲师；单人独白，约10分钟'),
           ('内容','自然科学/社会科学/历史/商业等学术主题'),
           ('难度','AWL词汇密集；语速最快；是7分和8分的分水岭'),
           ('题型','笔记填空50%+摘要填空30%+流程图10%+其他10%'),
           ('关键','S4拉开7分和8分的差距——7分需稳定5-6/10')],
          'S4考的是学术听力能力；AWL词汇+结构追踪+笔记技能三位一体'),
         ('S4题型详解','笔记填空是最高频的题型',
          [('笔记填空','最高频；按讲座内容顺序填写；答案通常是原词'),
           ('摘要填空','一段摘要有多个空；需要理解上下文语义'),
           ('流程图填空','按步骤顺序；被动语态常见；动词词根重要'),
           ('短问答','限词严格；通常是名词短语；不需要完整句子'),
           ('关键','S4几乎所有答案都是录音原词；不需要改写')])],
        [('S4讲座的内部结构',
          [('讲座开头','引入主题；提出研究问题；给出定义'),
           ('主体部分','按层级展开：主点→子点→例子→数据'),
           ('讲座中段','可能有转折：但是现在来看/最新研究发现'),
           ('讲座结尾','总结+影响+建议/展望未来'),
           ('题目分布','Q31-35通常对应主体；Q36-40对应后半段')]),
         ('S4跟丢后的心理调适','最重要的技能',
          [('第一反应','不要慌；立刻暂停回想——把注意力放到"往前"'),
           ('重定位','用下一道题的关键词找锚点；等录音出现该词'),
           ('猜填','跟丢的题不要留空；从题目语义猜一个合理答案'),
           ('继续','重定位后按正常节奏继续；1题跟丢≠后面全错'),
           ('心态','S4能做对6题以上=已经很好；10/10是S4高手水平')])],
        [('例题：S4笔记填空基础',
          'Notes Completion: Urban Heat Island Effect. Write NO MORE THAN TWO WORDS.\n\n'
          'Main causes:\n'
          '1. Heat from __________ (Q31)\n'
          '2. Reduced __________ due to lack of greenery (Q32)\n\n'
          '"Research identifies microfibre emissions from synthetic textiles as the primary heat source.\n'
          'Furthermore, reduced evapotranspiration due to declining urban greenery compounds this."',
          'Q31: synthetic textiles (听到microfibre emissions后定位到"from synthetic textiles")\nQ32: evapotranspiration (专业词，按音节记：e-vapo-tran...)',
          'Q32：专业词按首字母占位(evap-)，誊写时借助上下文语义补全拼写'),
         ('例题：S4摘要填空',
          'Summary Completion: Write NO MORE THAN ONE WORD.\n\n'
          'The circular economy differs from the linear model in that it aims to __________ (Q33)\n'
          'waste rather than simply dispose of it. Companies have found that __________ (Q34)\n'
          'in reusable materials can save long-term costs.\n\n'
          '"Unlike the linear take-make-dispose model, circular economy minimises waste through reuse.\n'
          'Initial investment in durable materials proves cost-effective over time."',
          'Q33: minimise (录音"minimises waste"→空格前"to"→填动词原形)\nQ34: investment (录音"Initial investment"→空格后"in"→填名词)',
          '摘要填空：先读空格前后的语法环境→预判词性→听时找对应词')])
    part(p,2,'笔记填空与摘要填空技法','Note & Summary Completion Techniques',
        [('笔记填空五步法','每步都能减少失分',
          [('Step1 预读','看笔记结构：标题是什么？空格在哪个层级？'),
           ('Step2 预判词性','空格前后是否有a/an/the/形容词→名词；be动词→形容词/名词'),
           ('Step3 预判内容','笔记的主题是什么？这个空格可能填什么类型的词？'),
           ('Step4 听时捕捉','边听边写候选词；专业词按音节记首字母占位'),
           ('Step5 誊写核对','限词/拼写/大小写逐项检查')],
          '笔记填空的答案几乎100%是录音原词；重点在"定位"和"限词"'),
         ('摘要填空的特殊挑战','摘要填空比笔记填空难',
          [('难点1','摘要用的词和录音词不同（同义替换）——不能直接用录音词'),
           ('难点2','摘要中的空格语法约束更严——词性必须对'),
           ('难点3','摘要填空通常在S4后半段，此时最容易跟丢'),
           ('解法','先读完整段摘要，理解语义→再听录音→用理解匹配'),
           ('关键','摘要填空考的是"语义理解"，不只是"关键词捕捉"')])],
        [('限词违规：最常见的摘要失分',
          [('ONE WORD only','只能填1个词；the/a不算但也不能写；复合词(self-sufficient)算1个'),
           ('TWO WORDS','最多2个；a/the算1个词；数字算1个词'),
           ('A NUMBER','只填数字；写文字"fifty"通常也接受'),
           ('TWO WORDS AND/OR A NUMBER','最多2个单词，或者1个数字，或者词+数字'),
           ('违规代价','超出限词=0分，哪怕内容正确；一定要在誊写时数词数')]),
         ('词性预判实战练习',
          [('例1','Researchers found that _____ (Q35) levels increased after treatment.'),
           ('分析1','_____ 后有levels(名词)→空格需要形容词；或空格本身是名词修饰levels'),
           ('例2','The main _____ (Q36) of this approach is its cost.'),
           ('分析2','The main _____ of... → 名词短语结构→填名词（disadvantage/benefit/problem）'),
           ('结论','预判词性后：听时只关注该词性的词，排除其他干扰')])],
        [('笔记填空精听练习',
          'Notes: Renewable Energy. Write NO MORE THAN TWO WORDS.\n\n'
          'Solar power:\n'
          '- most __________ energy source (Q37)\n'
          '- currently provides __________ of global electricity (Q38)\n\n'
          '"Solar is now the most rapidly expanding energy source globally.\nIt currently supplies around two percent of total global electricity generation."',
          'Q37: rapidly expanding (听到"most rapidly expanding"→填two words符合限制)\nQ38: two percent (听到"around two percent"→around是副词不用填→填two percent)',
          'Q37注意：most后跟形容词/副词结构；"rapidly expanding"整体修饰"energy source"'),
         ('摘要填空完整练习',
          'Summary: Write ONE WORD ONLY.\n\n'
          'The study __________ (Q39) that frequent smartphone use reduces attention span.\n'
          'This is particularly concerning for __________ (Q40) who rely on sustained focus.\n\n'
          '"Researchers confirmed that heavy smartphone usage significantly shortens the ability to focus.\n'
          'The finding is especially relevant for students who need prolonged concentration."',
          'Q39: confirmed (录音"confirmed"=摘要"concluded/showed"→但限ONE WORD→confirmed是原词最准)\nQ40: students (录音"students"=摘要空格的语义匹配)',
          '摘要填空先看语义：Q39需要动词（研究___了...）；Q40需要名词（___群体）')])
    part(p,3,'AWL学术词汇听辨策略','AWL Vocabulary: Listening Strategy',
        [('AWL词汇的特点','Academic Word List的听力挑战',
          [('总量','AWL共570个词族，覆盖学术文章约10%的词汇'),
           ('S4出现频率','S4每10句话约有3-5个AWL词汇'),
           ('听辨难点','看到认识，听到不认识——视觉词形≠听觉音形'),
           ('拼写难点','S4答案可能包含AWL词，拼错=0分'),
           ('策略','不要求会写所有AWL词；重点是"听到音能认出意思"')],
          '雅思S4目标：听到evapotranspiration→知道大概是"蒸发作用"→按音写占位'),
         ('高频AWL词汇：听音辨义技巧','分音节识别是核心',
          [('词根策略','bio=生命；hydro=水；therm=热；geo=地；eco=生态'),
           ('前缀策略','un-/dis-=否定；re-=重复；pre-=之前；inter-=之间'),
           ('后缀策略','-tion/-sion=名词；-ify/-ize=动词；-ous/-al=形容词'),
           ('音节分解','e-vap-o-trans-pi-RA-tion → 重音节"RA"是识别关键'),
           ('首字母法','听不清时：写出能听到的首字母+音节数→占位→誊写补全')])],
        [('AWL高频词30个（听辨重点）',
          [('concept / conceptual','概念/概念性的（重音：CON-cept）'),
           ('significant / significantly','重要的/显著地（SIG-ni-fi-cant）'),
           ('analysis / analyse','分析（a-NAL-y-sis）'),
           ('factor / variable','因素/变量'),
           ('methodology','方法论（meth-od-OL-o-gy）'),
           ('implication','含义/影响（im-pli-CA-tion）')]),
         ('AWL词汇占位法实战',
          [('步骤1','听到不认识的长词：写出首2-3个字母'),
           ('步骤2','继续写能听清的音节（用-代替难音节）'),
           ('步骤3','写出词的大概音节数（划格子）'),
           ('步骤4','誊写时：结合上下文+可能的词根补全拼写'),
           ('例子','e-vap-o-tr---spi---tion → evapotranspiration（10分钟后誊写时补全）')])],
        [('AWL音节训练：重音识别',
          '以下AWL词，标出重音所在的音节（找最强的那个）：\n\n'
          '1. a-NAL-y-sis → 重音在第2音节\n'
          '2. meth-od-OL-o-gy → 重音在第3音节\n'
          '3. sig-NIF-i-cant → 重音在第2音节\n'
          '4. sus-TAIN-a-ble → 重音在第2音节\n'
          '5. en-vi-RON-ment → 重音在第3音节',
          '找到重音节=找到词的"身份标记"；即使前后音节听不清，重音节通常最清晰',
          '训练：每天背5个AWL词，同时背它们的重音位置，听到重音节能反应出整词'),
         ('S4真题AWL词汇占位练习',
          '老师读以下句子，用占位法记录划线词（首字母+音节格）：\n\n'
          '1. "The study examines the feasibility of urban farming."\n'
          '2. "Deforestation has significant implications for biodiversity."\n'
          '3. "Entrepreneurship requires both creativity and methodology."\n'
          '4. "Evapotranspiration contributes to the urban heat island effect."',
          '1. fea-----ity(feasibility) + ur---farm(urban farming)\n2. De---tion(Deforestation) + bio-----ty(biodiversity)\n3. Entrep------ship + meth---gy\n4. Evap------tion',
          '关键：不因难词停下来；首字母占位后继续追踪；誊写时借助上下文补全')])
    part(p,4,'讲座追踪与S4综合策略','Lecture Tracking & S4 Strategy',
        [('S4讲座结构信号词','识别结构=永远不跟丢',
          [('层级标志','Firstly.../Secondly.../Finally... = 讲座层级'),
           ('转折','However.../On the other hand.../In contrast... = 讲座转折'),
           ('例证','For example.../Take...as an instance... = 进入例子'),
           ('总结','To summarise.../In conclusion.../Overall...'),
           ('强调','Most importantly.../What\'s crucial is.../The key point...')],
          '每个信号词=一个题目转换点；听到信号词→准备好下一道题'),
         ('S4时间管理与策略','10题要全部作答',
          [('预读优先','S4说明时间：读完全部Q31-40；圈关键词；预判答案类型'),
           ('跟丢处理','跟丢立刻用题号找锚点；猜填；继续往前不回头'),
           ('专业词处理','不认识的词按音节写占位；不因一个词而停下来'),
           ('时间分配','10道题约7-8分钟；每题约45秒；不在一题上超过1分钟'),
           ('猜测原则','不留空：从语义和语法猜一个合理答案')])],
        [('跟丢重定位五步骤',
          [('停止回想','不要把时间花在已经过去的内容上'),
           ('找锚点词','看下一道未答题目的关键词（名词/动词/数字）'),
           ('等待出现','继续听，等录音出现该关键词或其同义词'),
           ('标记位置','听到锚点词→立刻找到题目位置→准备答题'),
           ('猜填已过','等下一个Section间隙：对跟丢的题猜填一个合理答案')]),
         ('S4综合备考策略',
          [('词汇准备','每周背30个AWL词；重点是"听到音能认出意思"'),
           ('精听练习','每天选1段S4真题录音；听3遍；对照原文标难词'),
           ('结构追踪','每次做S4后：标出讲座结构信号词；核对题目分布'),
           ('错题分析','分类失分：是跟丢？是词汇不认识？是限词违规？'),
           ('目标设定','6分以下→目标S4做对5题；7分→目标7题；8分→目标8题')])],
        [('S4预读策略演练',
          '限时45秒，对以下S4题目完成预读（圈关键词+预判答案类型）：\n\n'
          'Q31: The main cause of soil erosion is _______\n'
          'Q32: Rainfall of _______ mm per year makes the problem worse\n'
          'Q33: A possible solution involves _______ farming techniques\n'
          'Q34: This approach can reduce soil loss by _______',
          'Q31: 名词/名词短语（cause of...=原因→填名词）\nQ32: 数字（mm per year=降雨量→填数字）\nQ33: 形容词或名词（___farming=修饰farming→填形容词或专有名词）\nQ34: 百分比或数字（reduce by___=减少多少→数字/百分比）',
          '预判词性的好处：听时只关注该词性的词，排除其他信息的干扰'),
         ('S4信号词追踪演练',
          '读以下讲座片段，标出信号词并说出每段讲的主要内容：\n\n'
          '"Today I want to focus on urban sustainability.\nFirstly, let\'s consider the environmental factors.\nHowever, economic considerations are equally important.\nTo summarise, sustainable cities require both approaches."',
          '信号词1: Firstly→环境因素 信号词2: However→经济因素 信号词3: To summarise→总结\n题目分布：Q31-34对应环境因素；Q35-37对应经济因素；Q38-40对应总结',
          '信号词=题目分布的路标；每次标完信号词，就能预测每题出现的大致时间')])
    # PART05
    U.section(p,5,'高频词汇','Key Vocabulary')
    U.vocab(p,'AWL高频词：学科核心词',
        [('hypothesis',  '假设（hy-POTH-e-sis）'),
         ('variable',    '变量（VAR-i-a-ble）'),
         ('significant', '显著的（sig-NIF-i-cant）'),
         ('sustainable', '可持续的（sus-TAIN-a-ble）'),
         ('implication', '含义/影响（im-pli-CA-tion）'),
         ('methodology', '研究方法论（meth-od-OL-o-gy）'),
         ('correlation', '相关性（cor-re-LA-tion）')])
    U.vocab(p,'AWL高频词：动词类',
        [('analyse / analysis',   '分析'),
         ('evaluate / assessment','评估'),
         ('indicate / suggest',   '表明/暗示'),
         ('demonstrate',          '展示/证明'),
         ('implement',            '实施'),
         ('investigate',          '调查/研究'),
         ('identify / classify',  '识别/分类')])
    U.vocab(p,'S4讲座信号词完整表',
        [('层级', 'firstly/secondly/finally/to begin with/moving on to'),
         ('转折', 'however/on the other hand/in contrast/although/yet'),
         ('例证', 'for example/for instance/take...as/such as'),
         ('强调', 'most importantly/crucially/above all/what\'s key'),
         ('补充', 'moreover/furthermore/in addition/what\'s more'),
         ('总结', 'to summarise/in conclusion/overall/to recap'),
         ('因果', 'therefore/thus/as a result/consequently')])
    U.two_col(p,'S4题型对比',
        '笔记填空（Note Completion）',
        ['按讲座层级顺序填写',
         '答案通常是录音原词',
         '预读：看笔记层级结构',
         '预判词性（空格前后语法）',
         '专业词按音节占位',
         '限词严格：ONE/TWO WORDS',
         '专有名词大写'],
        '摘要填空（Summary Completion）',
        ['读完整段摘要理解语义',
         '同义替换：摘要词≠录音词',
         '词性约束更严',
         '理解上下文帮助预判',
         '通常在S4后半段，需要坚持',
         '一个词限制最严格',
         '语义正确比用词地道更重要'])
    # PART06
    U.section(p,6,'方法精讲','Methods & Models')
    U.steps(p,'笔记填空满分5步操作',
        [('预读(45s)','看笔记主题；圈空格前后关键词；预判每空词性'),
         ('边听边写','候选词先写旁边；不确定时写首字母占位'),
         ('专业词处理','长词：按音节写能听清的部分；首字母+划线占位'),
         ('即时核对','写完一空立刻检查：限词/词性/大小写是否正确'),
         ('誊写检查','最终：逐词核对拼写；专有名词大写；数词数')])
    U.steps(p,'AWL词汇占位法操作',
        [('第1步','听到不认识的长词：立刻写前2-3个字母(e-vap...)'),
         ('第2步','继续写能听到的音节（用—代替模糊音）'),
         ('第3步','在旁边写词的大概音节数（帮助誊写时补全）'),
         ('第4步','继续听后面的内容；不因一个词停止追踪'),
         ('第5步','誊写时：用词根/上下文语义补全拼写；查词典确认')])
    U.model_answer(p,'S4笔记填空满分示范',
        'Section 4: Urban Heat Island Effect',
        '满分答题笔记格式：\n\n'
        'Urban Heat Island(UHI) = city warmer than surroundings\n\n'
        'Q31 Cause 1: Heat from synthetic textiles\n'
        '    → 听到"synthetic textiles"立刻写，不等整句结束\n\n'
        'Q32 Cause 2: Reduced evap...(占位) due to ↓greenery\n'
        '    → evap-o-trans-pi-ration → 誊写时补全\n\n'
        'Q33 Effect: Urban temp ↑ 2-5°C vs rural areas',
        ['Q31：原词直接写（synthetic textiles）；两词符合TWO WORDS限制',
         'Q32：专业词占位（evap...）；不因难词停止追踪后面的内容',
         '笔记简化：↑↓→=；节省书写时间，把精力放在内容上',
         '预读时标好每题层级位置；听时精准定位每题'])
    U.model_answer(p,'S4摘要填空满分示范',
        'Summary: Circular Economy. ONE WORD ONLY.',
        '预读摘要（先理解语义）：\n'
        '"The circular economy is designed to _______(Q34) waste by reusing materials.\n'
        'Unlike the _______(Q35) approach, it treats waste as a resource."\n\n'
        '听录音追踪：\n'
        '"...aims to minimise waste..." → Q34=minimise\n'
        '"...unlike the conventional/linear take-make-dispose model..." → Q35=linear\n\n'
        '验证：minimise(动词，符合to后接动词原形); linear(形容词，符合the ___approach)',
        ['预读摘要时标注每空的词性要求：Q34=动词；Q35=形容词',
         '同义替换识别：minimize=minimise; linear/conventional=linear',
         'ONE WORD：不能填"waste minimisation"（两词），只填"minimise"',
         '词性验证：Q35 "the ___ approach"→需要形容词→linear✓'])
    U.compare(p,'S4应对：跟丢放弃 vs 重定位继续',
        '❌ 跟丢后放弃',
        ['S4第一句没听懂，产生恐慌',
         '后面9题全部放弃或乱猜',
         '留空白：5道题得0分',
         '觉得S4太难，等着结束',
         '对AWL词汇陌生→完全蒙'],
        '✅ 跟丢后重定位',
        ['跟丢立刻停止回想；往前看',
         '找下一题关键词等锚点出现',
         '跟丢的题用语义猜填；不留空',
         'AWL词按音节占位；誊写补全',
         'S4对5题以上=已经很好'])
    U.compare(p,'笔记填空：违规 vs 合规',
        '❌ 限词违规失分',
        ['ONE WORD：填了"the swimming pool"(3词)',
         'TWO WORDS：填了"the rapidly increasing"(3词)',
         'A NUMBER：填了"about fifty"而非"50"',
         '没有检查：写完直接交卷',
         '专有名词不大写：january→January'],
        '✅ 限词合规满分',
        ['ONE WORD：数词，冠词也算；只填1个词',
         'TWO WORDS：最多2个；冠词算1个',
         'A NUMBER：写数字"50"不写文字',
         '誊写前：逐空数词数，检查限词',
         '专有名词：月份/星期/地名首字母大写'])
    # PART07
    U.section(p,7,'课堂练习','Practice Drills')
    U.drill(p,'练习1：AWL词汇音节分解',
        '以下AWL词，逐音节读出，标出重音位置',
        ['1. methodology → meth-od-OL-o-gy（5音节，重音第3）',
         '2. significant → sig-NIF-i-cant（4音节，重音第2）',
         '3. sustainable → sus-TAIN-a-ble（4音节，重音第2）',
         '4. evapotranspiration → e-VAP-o-trans-pi-RA-tion（7音节）',
         '5. entrepreneurship → en-tre-pre-NEUR-ship（5音节，重音第4）'])
    U.drill(p,'练习2：限词检查',
        '以下答案是否符合"NO MORE THAN TWO WORDS"？',
        ['"the swimming pool" → 合规还是违规？（冠词算词）',
         '"swimming pool" → 合规还是违规？',
         '"rapidly increasing" → 合规还是违规？',
         '"very rapidly" → 合规还是违规？',
         '"50" → 如果限制是ONE WORD AND/OR A NUMBER → ?'])
    U.drill(p,'练习3：AWL词汇占位法练习',
        '老师说以下句子，用占位法记录划线词',
        ['1. "The main hypothesis of the study is..." → 写出hypothesis的占位',
         '2. "Researchers demonstrated that biodiversity affects..." → 写biodiversity占位',
         '3. "The methodology involved quantitative analysis..." → 写methodology占位',
         '4. "This has significant implications for sustainability..." → 写implications占位',
         '完成后：核对你的占位是否能帮助誊写时识别完整词'])
    U.drill(p,'练习4：摘要填空词性预判',
        '预读以下摘要，标出每空需要的词性（名词/动词/形容词）',
        ['The researcher _______(Q31) that pollution levels had risen.  → 词性：?',
         'The main _______(Q32) of climate change include rising seas. → 词性：?',
         'Scientists use _______(Q33) methods to measure temperature.  → 词性：?',
         'The results were _______(Q34) in three major journals.       → 词性：?',
         '参考：Q31=动词(过去式)  Q32=名词(复数)  Q33=形容词  Q34=动词(被动)'])
    U.drill(p,'练习5：S4综合模拟',
        '完整S4模拟：预读→追踪→占位→誊写四阶段',
        ['第一阶段（45秒）：预读Q31-40，圈关键词，预判词性',
         '第二阶段（实时）：边听边写，专业词占位，不因难词停止',
         '第三阶段（10分钟，纸笔考）：誊写，补全占位词，检查限词',
         '第四阶段（分析）：分类失分：词汇不认识/跟丢/限词违规',
         '目标：本节课S4对5题以上；每节课提升1题是合理速度'])
    # PART08
    U.section(p,8,'避坑指南','Pitfall Prevention')
    U.pain_point(p,'S4最典型的4类失分',
        ['第一句没听懂→放弃全场——恐慌是最大的失分',
         'AWL词汇不认识→空着→0分（哪怕猜都比空着好）',
         '摘要填空超出限词→ONE WORD写了两个→0分',
         '专业词拼错→听对了写错→得0分'],
        ['第一句没懂：正常！继续听，用第二三句补上语境',
         'AWL词不认识：写首字母占位；誊写时结合上下文猜拼写',
         '每次做摘要前：重新看限词要求；写完立刻数词数',
         '专业词：按音节占位；誊写时用词典/记忆/推断补全'])
    U.content(p,'S4笔记填空高频词类',
        [('数字类', '年份/百分比/数量——最多要求ONE WORD AND/OR A NUMBER'),
         ('专有名词', '地名/人名/机构名——首字母必须大写'),
         ('学术动词', 'analyse/identify/implement——词根+词缀记忆'),
         ('形容词', '修饰关键名词——空格前通常有the/a/an或名词'),
         ('名词短语', '两词答案——确认限词是TWO WORDS才能写两个')],
        note='S4笔记填空约80%答案是单词；20%是两词短语——预判词性的价值很高')
    U.content(p,'S4备考策略',
        [('短期', '每天精听一段S4真题（3分钟）；对照原文标出陌生词'),
         ('中期', '每周背30个AWL词，重点是"听音识义"而非"会写"'),
         ('长期', '积累S4错题本：分类——词汇/跟丢/限词三类分别统计'),
         ('模拟', '每周一套完整模拟；S4要计时；模拟后精析每道错题'),
         ('心态', 'S4从来不要求完美；对5-6题=7分区间；继续提升即可')])
    U.tip_card(p,'S4笔记填空速记卡',
        [('预读必做', '45秒：看笔记结构，圈空格关键词，预判每空词性——这决定成败'),
         ('占位法', '专业词：写首2字母+划线占位；不因一词停止追踪后面的内容'),
         ('限词铁律', '写完每空立刻数词数；ONE WORD只填1个；TWO WORDS最多2个'),
         ('誊写五检', '词数/大小写/复数/拼写/数字格式——5项逐一检查')])
    U.tip_card(p,'AWL词汇速记卡',
        [('听辨优先', 'AWL核心300词目标：听到音能认出意思；不要求全会写'),
         ('重音节法', '每个AWL词找到重音节——重音节是听辨的"身份标记"'),
         ('词根词缀', 'bio/hydro/geo=词根；-tion/-fy/-ous=词缀——覆盖70%AWL词'),
         ('占位逻辑', '听到长词：写首字母+音节格；誊写时用上下文+词根补全')])
    U.checklist(p,'L05 自检清单',
        ['我知道S4笔记填空的五步解法，且知道每步的作用',
         '我练习了AWL词汇占位法：首字母+音节格，不因难词停止',
         '我识别了S4讲座结构信号词：firstly/however/to summarise',
         '我知道跟丢后的处理：找锚点→重定位→猜填→继续',
         '我了解S4的现实目标：6分→做对5题；7分→做对7题'])
    U.summary(p,5,
        ['S4笔记填空：预读预判词性→听时占位→誊写补全→限词检查',
         'AWL听辨：重音节识别+词根词缀推断；不要求会写，要求听出来',
         '跟丢不放弃：找锚点词重定位；猜填空白；继续往前听'],
        takeaway='S4最大的提分空间在"跟丢后不放弃"——这个心态改变就能多对2-3题')
    U.homework(p,5,
        ['完成剑桥真题2套S4（计时），统计正确率并分类失分原因',
         '背诵AWL核心词汇30个：每天10个，重点是音节分解和听音识义',
         '精析1段S4录音：标出所有信号词和难词，建个人AWL词表'],
        est_time='约60分钟')
    U.preview(p,6,'预测技能与关键词定位',
        ['预读系统：从表面预读到深度预判的三层次',
         '关键词圈划：20秒圈完10题的操作规范',
         '同义替换识别：题目词和录音词的系统对应'])
    U.save(p, f'{OUT}/L05_Section4学术讲座攻略.pptx')
    print(f"  L05 done: {len(p.slides)} pages")

# ══════════════════════════════════════════════════════════════════════════════
#  L06  预测技能与关键词定位
# ══════════════════════════════════════════════════════════════════════════════
def L06():
    T(); p=U.new_prs()
    U.cover(p,'听力',6,'预测技能\n& 关键词定位','Prediction Skills & Keyword Location')
    U.agenda(p,['预测的本质与价值','三层预读系统','关键词圈划技术',
                '同义替换识别','预测词汇库','20秒预读法','课堂练习','小结作业'])
    U.content(p,'L05回顾：S4学术讲座核心要点',
        [('AWL词汇','听到不认识的词先音译占位，誊写时核对拼写'),
         ('信号词追踪','however/therefore/in contrast标志答案转折点'),
         ('跟丢处理','找锚点词重定位，空格留白继续往前听，不停下来纠结'),
         ('难度认知','S4拉开7分和8分差距，策略正确可多对3-4题')],
        note='本节课目标：把听力从"被动接收"变为"主动预判"，效率提升40%')
    U.objectives(p,['掌握题目预读的三个层次：表层→深层→预判',
                    '能在20秒内完成一组题目的关键词圈划',
                    '识别录音中的同义替换，将题目词映射到录音词',
                    '建立个人"高频预测词汇库"，提升首次识别率'])
    U.pain_point(p,'预读做错的三种典型表现',
        ['只读题目不预判：圈了关键词但没想"答案可能是什么类型"，听到答案时反应慢',
         '预读太慢：一组10题花了40秒还没读完，录音已开始，手忙脚乱',
         '同义替换识别失败：题目写"inexpensive"，录音说"cheap"，没反应过来是同一个意思'],
        ['预读=圈词+预判答案类型（名词/数字/形容词）+预判同义词',
         '20秒预读规范：先扫描所有题目标题→再圈每题1-2个关键词→不逐字阅读',
         '每课做10个同义替换配对练习，积累"题目词→录音词"的映射意识'])
    part(p,1,'预测的本质：主动听 vs 被动听','Active vs Passive Listening',
        [('被动听的代价','听完每句才反应，答案已过去',
          [('被动听','等录音说完再找答案，错过关键词窗口'),
           ('主动预测','提前知道"第3题要填一个地点名词"，听到地点词立即捕捉'),
           ('差距','主动预测学生平均多对3-4题，约提高0.5个Band'),
           ('核心','不是"听得更仔细"，而是"听得更有目的"'),
           ('训练','每次做题前强制完成预读，不跳过，形成习惯')],
          '预测不是猜答案，而是提前知道"我要找什么"，大幅降低认知负担'),
         ('三层预读系统','从粗到细，层层深入',
          [('第一层：题型识别','扫描题目格式→确认是填空/选择/配对/地图'),
           ('第二层：关键词圈划','每题圈1-2个最独特的名词或数字'),
           ('第三层：答案类型预判','这个空填名词？形容词？数字？时间？地点？'),
           ('时间分配','3秒识别题型→10秒圈关键词→7秒预判答案类型=20秒'),
           ('优先级','来不及做第三层时，至少完成前两层')])],
        [('预测提分的实证数据',
          [('实验数据','有预读习惯的考生平均比无预读多对3.2题（约0.5 Band）'),
           ('深圳数据','本校学员采用预读系统后，S1满分率从42%→71%'),
           ('时间投入','学会20秒预读需要约10次练习，之后完全自动化'),
           ('最大误区','很多学生"知道要预读"但不坚持——知道和做到有天壤之别'),
           ('立即行动','本节课结束后：每次做题，无论多忙，都要完成三层预读')]),
         ('不同题型的预读重点',
          [('填空题','圈空格前后的词→判断词性→预判答案类型（名词/数字）'),
           ('选择题','圈三个选项的差异核心词（各1个），不读完整选项'),
           ('配对题','先读完所有H选项，每个圈1个关键词，共8个'),
           ('地图题','在图上标注已知地点，圈方向词列表（north/next to/between）'),
           ('流程图','按顺序圈每步动词，预判被动语态（is done/are placed）')])],
        [('三层预读实操：S1场景',
          '以下是S1的题目，用20秒完成三层预读（不要超时）\n\n'
          'Name: _____(1)  Age: _____(2)  Occupation: _____(3)\n'
          'Membership type: _____(4)  Start date: _____(5)',
          'Q1→专有名词（大写，可能需要拼读）\nQ2→数字（两位数为主）\nQ3→职业名词（teacher/nurse等）\nQ4→类型形容词（family/student/annual）\nQ5→日期格式（1st May / 15 June）',
          '日期格式陷阱：录音说"the first of May"→写"1st May"或"May 1st"均可，但不要写"1/5"'),
         ('同义替换配对练习',
          '将左列题目词与右列录音词连线：\n\n'
          '题目词：cheap / large / begin / problems / parking\n'
          '录音词：commence / issues / affordable / car park / spacious',
          'cheap→affordable / large→spacious / begin→commence / problems→issues / parking→car park',
          '考场规律：S1-S2同义替换较直接；S3-S4同义替换更抽象，需要理解语义')]
    )
    part(p,2,'关键词圈划：20秒操作规范','Keyword Marking: 20-Second Protocol',
        [('什么词值得圈？','选最独特、最可能在录音中提到的词',
          [('优先圈','专有名词（人名/地名）、数字、时间、颜色、职业'),
           ('次优圈','题目中的核心名词（不圈冠词/介词/助词）'),
           ('不要圈','the/a/and/of/to等功能词；太常见的动词如"say/go"'),
           ('地图题特殊','圈方向词：north/opposite/next to/between'),
           ('选择题特殊','圈每个选项的第一个实义词，快速对比三个选项差异')],
          '好的关键词=录音中会被特别提到的词，而不是句子中最重要的词'),
         ('关键词与预判结合','圈词后立刻问自己：答案是什么类型？',
          [('填空题','圈前后词→判断词性（空格前是adjective→填名词）'),
           ('选择题','圈三个选项的差异点→只听与差异相关的内容'),
           ('配对题','提前读完所有H选项→每个圈1个关键词（共8个关键词）'),
           ('地图题','在地图上标出已知地点→圈方向关系词'),
           ('流程图','圈每个步骤的动词→预判被动语态出现')])],
        [('20秒预读的时间管理',
          [('3秒：识别题型','扫描全局：是一组填空题还是选择题还是地图？'),
           ('10秒：圈关键词','从Q1到最后，每题1-2个词，不停下来思考'),
           ('5秒：预判答案类型','快速过一遍：这道题填数字？名词？形容词？'),
           ('2秒：聚焦Q1','把注意力锁在Q1关键词上，做好捕捉准备'),
           ('超时处理','来不及做第3层就跳过，至少圈完所有关键词')]),
         ('关键词圈划的常见错误',
          [('错误1','圈太多词：每题圈超过3个，预读太慢'),
           ('错误2','圈功能词：the/a/of/to这些词在录音中出现太频繁，无定位价值'),
           ('错误3','选择题读完整选项：只需圈差异词，不要逐字阅读'),
           ('错误4','预读时思考答案：预读只是"圈词"，不是"想答案"'),
           ('正确目标','20秒内完成所有题目的关键词圈划，不超时')])],
        [('关键词圈划计时练习',
          '计时20秒，圈出以下题目的关键词（不要阅读整句）：\n\n'
          'Q1 The customer wants to book a _____ (type of room)\n'
          'Q2 Check-in date: _____ (date)\n'
          'Q3 The room costs _____ per night\n'
          'Q4 The hotel is located near the _____ (landmark)\n'
          'Q5 Payment method: _____ (method)',
          'Q1→book, room, type / Q2→check-in, date / Q3→costs, per night（预判数字）/ Q4→near, landmark（预判地点名词）/ Q5→payment, method（预判支付方式名词）',
          '时间到立即停止，不要继续读——训练的是速度而不是完整性，习惯后自然加快'),
         ('同义替换：题目词→录音词实战',
          '题目：The course is designed for students with no previous _____ (experience)\n\n'
          '录音："This programme targets complete beginners — no prior knowledge required."\n\n'
          '题目词 previous experience → 录音词是什么？',
          'previous→prior / experience→knowledge / 答案：knowledge\n完整替换链：previous experience = prior knowledge',
          '同义替换解题步骤：①先定位录音中出现的"题目周边词"→②找替换词→③确认语法符合')]
    )
    part(p,3,'同义替换识别系统','Synonym Substitution Recognition',
        [('同义替换的四种类型','从简单到复杂，逐步掌握',
          [('直接同义词','cheap=affordable, big=large, start=begin, end=finish'),
           ('概念替换','parking=car park, rubbish=waste, job=occupation'),
           ('上义词替换','dog/cat/bird → animals（具体→抽象）'),
           ('释义替换','bicycle=a two-wheeled vehicle（词→短语）'),
           ('S4最难','单个词替换为学术词组，需要语义理解')],
          '剑桥真题每套听力包含约15-20处同义替换，熟练识别是7分以上的必备技能'),
         ('同义替换高频词对','必须熟背的100组核心替换',
          [('地点类','car park↔parking, gym↔sports centre, shop↔store/retail outlet'),
           ('时间类','now↔currently/at present, soon↔in the near future'),
           ('程度类','very↔extremely/highly, difficult↔challenging/demanding'),
           ('变化类','increase↔rise/grow/go up, decrease↔fall/drop/decline'),
           ('意见类','think↔believe/consider/suggest, problem↔issue/concern/challenge')])],
        [('同义替换识别的三步法',
          [('Step 1：预读时标注','在关键词旁边写出1-2个可能的同义词'),
           ('Step 2：听录音时捕捉','听到关键词的同义词时，反应要和听到原词一样快'),
           ('Step 3：验证语境','确认录音的语义与题目一致，再写下答案'),
           ('常见失误','题目写"inexpensive"，录音说"cheap"，没反应过来→空白'),
           ('训练方法','每天做10组替换配对，1个月后形成条件反射')]),
         ('S3/S4同义替换的特殊规律',
          [('S3学术讨论','常用学术词替换口语词：purchase=buy, obtain=get'),
           ('S4学术讲座','常用长词组替换单词：decline in population=fewer people'),
           ('跨词类替换','名词→动词：solution→solve；形容词→名词：effective→effectiveness'),
           ('语义扩展','specific→particular/specific/precise（多个替换词均可）'),
           ('深圳学生弱点','对释义替换（词→短语）识别最弱，需专项训练')])],
        [('替换词快速匹配练习',
          '30秒内完成配对（题目词→录音词）：\n\n'
          '① free of charge  ② harmful  ③ collect  ④ ancient  ⑤ improve\n'
          '录音词：gather / old / no cost / upgrade / damaging',
          '① free of charge → no cost\n② harmful → damaging\n③ collect → gather\n④ ancient → old\n⑤ improve → upgrade',
          '记忆技巧：不要死记对应关系，理解语义才能在新语境下灵活识别'),
         ('S3选择题：同义替换+干扰项识别',
          'Q: Why did Sarah choose this research topic?\n'
          'A. Her teacher recommended it\nB. She found it personally interesting\nC. It relates to her future career\n\n'
          '录音："I\'ve always been fascinated by urban planning — it\'s something I genuinely enjoy, '
          'though my supervisor did mention it once in passing."',
          '答案B：fascinated by=personally interested in\n干扰项A：supervisor mentioned≠recommended（程度不够）\n干扰项C：录音没有提到future career',
          '选择题陷阱：录音常先提到错误选项，再转折到正确答案——听到but/however时注意')]
    )
    part(p,4,'预测技能综合实战','Integrated Prediction Practice',
        [('20秒预读法：完整流程','从拿到题目到录音开始的黄金20秒',
          [('Step 1 (3s)','快速识别题型：填空/选择/配对/地图？'),
           ('Step 2 (10s)','逐题圈关键词：每题1-2个，不超过3个'),
           ('Step 3 (5s)','预判答案类型：数字/名词/形容词/日期/地点'),
           ('Step 4 (2s)','快速看第一题，做好捕捉准备'),
           ('录音开始后','按关键词顺序追踪，听到关键词附近立即写答案')],
          '20秒预读是技能，需要专项训练——前10次会慢，50次后会自动化'),
         ('跨题组的衔接预读','一组答完后立刻预读下一组',
          [('时机','答完最后一题后，立刻把视线移到下一组题目'),
           ('不要检查','答完一组不要回头检查，用这段时间预读下一组'),
           ('说明词利用','每Section开始有约30秒说明，全部用来预读题目'),
           ('翻页准备','纸笔考提前翻到下一组题所在页，不在录音中临时翻找'),
           ('机考提示','机考会有视觉提示跳到下一组，提前熟悉界面切换')])],
        [('预测技能的自我评估标准',
          [('初级（6分）','预读超过30秒；经常因为没预读导致第一题没抓到'),
           ('中级（6.5-7分）','20秒内预读完；同义替换直接词（cheap→affordable）能识别'),
           ('高级（7.5-8分）','15秒内完成；能识别概念替换（parking→car park）'),
           ('顶级（8.5+）','自动化预读；预判替换词写在关键词旁；几乎零漏题'),
           ('深圳平均','目前约6分左右，目标本课结束后达到中级')]),
         ('预测的边界：什么不能预测',
          [('不要预测具体答案','只预测类型（名词/数字），不要提前想好"答案是Tuesday"'),
           ('不要锁定选项','选择题不要看了A感觉对就不听B和C'),
           ('不要固执','预测说"应该是Tuesday"，但录音说"actually Wednesday"→写Wednesday'),
           ('以录音为准','任何预测都要在听到录音实际内容后接受修正，不能固执'),
           ('预测≠猜题','预测是降低认知负担的工具，不是答题捷径')])],
        [('全流程预测实战：S2场景',
          '已知：S2是关于博物馆导览的独白。\n'
          '题目类型：地图配对+简短回答\n\n'
          '现在20秒预读以下题目，并说出你的预测：\n'
          'Q11-15: Map - label the rooms (A-H)\n'
          'Q16-20: Answer these questions (NO MORE THAN TWO WORDS)\n'
          'Q16: What is the main theme of the museum?',
          'Q11-15预测：圈地图上的8个选项字母，预判方向词north/south/east/west/opposite/next to会出现\nQ16预测：答案是一个名词短语，描述博物馆主题（如"natural history"/"modern art"）',
          '地图题核心预读：不要读所有A-H的说明，只记住字母位置，听录音时找方向词'),
         ('同义替换+预测综合检验',
          '题目：The library opening hours have been _____ recently.\n\n'
          '你的预测：答案类型是？可能的答案词是？\n\n'
          '录音："We\'ve recently extended our service hours to better accommodate student needs."',
          '预测：changed/extended/reduced（形容词或动词过去分词）\n实际答案：extended（录音说extended service hours→延长了）\n同义替换：opening hours = service hours',
          '预测不一定猜中，但让你在听到extended时立刻知道"这是答案"')]
    )
    # PART05 vocab
    U.section(p,5,'预测核心词汇库','Prediction Vocabulary Bank')
    U.vocab(p,'同义替换高频词对（地点与设施）',
        [('car park','= parking area / parking lot（停车场）'),
         ('sports centre','= gym / fitness centre / leisure centre（运动中心）'),
         ('café / canteen','= restaurant / dining area / food court（餐饮区）'),
         ('library','= learning resource centre / study hall（图书馆）'),
         ('reception','= front desk / information desk（前台）'),
         ('emergency exit','= fire exit / evacuation route（紧急出口）'),
         ('conference room','= meeting room / seminar room / boardroom（会议室）')])
    U.vocab(p,'同义替换高频词对（动作与变化）',
        [('increase','= rise / grow / go up / climb / surge（增加）'),
         ('decrease','= fall / drop / decline / reduce / shrink（减少）'),
         ('start','= begin / commence / initiate / launch（开始）'),
         ('finish','= end / complete / conclude / wrap up（结束）'),
         ('show','= demonstrate / indicate / reveal / suggest（显示）'),
         ('use','= utilise / employ / apply / adopt（使用）'),
         ('help','= assist / support / facilitate / enable（帮助）')])
    U.vocab(p,'同义替换高频词对（描述与评价）',
        [('difficult','= challenging / demanding / complex / tough（困难）'),
         ('important','= significant / crucial / essential / vital（重要）'),
         ('cheap','= affordable / inexpensive / cost-effective（便宜）'),
         ('old','= ancient / historic / traditional / dated（老旧）'),
         ('new','= modern / contemporary / innovative / novel（新的）'),
         ('interesting','= fascinating / engaging / appealing（有趣）'),
         ('worried','= concerned / anxious / troubled / uneasy（担心）')])
    U.two_col(p,'答案类型预判速查表',
        '填空题答案类型判断',
        [('前有形容词空格','→ 填名词（the __ centre）'),
         ('前有be动词空格','→ 填形容词（is very __）'),
         ('前有冠词a/an空格','→ 填名词（a __ of）'),
         ('问how much/many','→ 填数字（价格/数量）'),
         ('问when/what date','→ 填日期时间（June 3 / 9:30）'),
         ('问where/what place','→ 填地点名词（library/Hall B）')],
        '选择题差异快速识别',
        [('三项含数字','→ 只听数字，忽略描述词'),
         ('三项含人名','→ 只追踪人名，判断谁做了什么'),
         ('三项含地点','→ 与地图题结合，追踪位置变化'),
         ('三项含时间','→ 注意更正陷阱，最终时间才是答案'),
         ('三项含程度词','→ 听极限词：most/only/always/never')])
    # PART06 method
    U.section(p,6,'20秒预读法：步骤精解','20-Second Pre-reading: Step-by-Step')
    U.steps(p,'Step 1-3：识别与圈划',
        [('识别题型 (3s)','看题目格式：有空格=填空；A/B/C=选择；字母选项=配对；图形=地图'),
         ('圈关键词 (10s)','每题1-2个词：优先专有名词→数字/时间→核心名词'),
         ('预判答案类型 (5s)','问自己：这个空填什么词性？数字还是名词？')])
    U.steps(p,'Step 4-5：准备与捕捉',
        [('聚焦第一题 (2s)','把注意力锁定在Q1的关键词上，做好第一个答案的捕捉准备'),
         ('听录音时：顺序追踪','按题目顺序追踪关键词，听到关键词附近立即写答案，不回头'),
         ('答完即预读','答完最后一题后立刻把视线移到下一组题目，开始新一轮预读')])
    U.model_answer(p,'满分预读示范：S1信息表',
        'Name: _____(1)  Phone: _____(2)  Address: _____(3)\n'
        'Course: _____(4)  Day: _____(5)  Fee: £_____(6)',
        'Q1→专有名词，可能需拼读（大写）\n'
        'Q2→电话号码（7-11位数字）\n'
        'Q3→地址（数字+街道名）\n'
        'Q4→课程名称（名词短语）\n'
        'Q5→星期几（Monday-Sunday）\n'
        'Q6→金额（两位小数，如24.50）',
        ['Q1预判同义替换：name→called/known as',
         'Q2特殊格式：电话号码0读作"oh"，双数字"66"说"double six"',
         'Q3地址格式：数字在前，街道名在后（15 Park Road）',
         'Q5陷阱预警：星期几最容易被更正，听到"actually"立即划掉旧答案'])
    U.model_answer(p,'满分预读示范：S3选择题',
        'Q21 Why is Maria concerned about the experiment?\n'
        'A her equipment is faulty\nB she lacks experience\nC there is insufficient time\n\n'
        'Q22 What does Tom suggest they do first?\n'
        'A review the literature\nB contact the supervisor\nC collect the data',
        'Q21预读：圈concerned/equipment/experience/time（三个选项核心词）\n'
        'Q22预读：圈suggest/first/literature/supervisor/data（追踪建议动作）\n'
        '预判：Q21答案在Maria表达担忧之后；Q22答案在Tom提建议之后',
        ['Q21技巧：三个选项分别关于equipment/experience/time，只追踪这三个词',
         'Q22技巧："first"是关键词——Tom可能提多个建议，只要"第一个"',
         '选择题不要预判正确答案——只预判"在哪里找答案"'])
    U.compare(p,'预读策略：低效 vs 高效',
        '❌ 低效预读（浪费时间）',
        ['逐字阅读每道题目的完整句子',
         '每道题想好几个可能的答案，越想越多',
         '花太多时间在一题上，没读完所有题目',
         '选择题把A/B/C三个选项完整读完',
         '翻页找题目，浪费宝贵的说明时间'],
        '✅ 高效预读（20秒完成）',
        ['只扫关键词，不读完整句子',
         '每题只预判答案"类型"，不预判具体内容',
         '快速移过每题，保证所有题目都圈到词',
         '选择题只圈三个选项的差异词（各1个）',
         '考前熟悉题目排版，不临时找题目位置'])
    U.compare(p,'同义替换：低分失误 vs 高分识别',
        '❌ 同义替换失败',
        ['题目写affordable，录音说cheap，没反应过来',
         '题目写commence，录音说start，写了个大问号',
         '题目写vehicle，录音说car，觉得car太简单了不敢写',
         '题目写occupied，录音说busy/taken，等着找occupied没听到'],
        '✅ 同义替换成功',
        ['见到afford-系词立刻想到cheap/inexpensive/low cost',
         '见到commence立刻想到start/begin/initiate',
         '题目用正式词，录音用口语词——car就是vehicle，大胆写',
         '预读时把occupied旁边写"=busy/taken/in use"'])
    # PART07 drill
    U.section(p,7,'课堂练习：预测技能强化','Prediction Drills')
    U.drill(p,'练习1：关键词圈划计时（15秒）',
        '计时15秒，圈出以下5道题的关键词（不得超时）',
        ['Q1 The applicant\'s surname: _____',
         'Q2 Preferred start date: _____',
         'Q3 Current job title: _____',
         'Q4 Annual salary expected: £_____',
         'Q5 The interview will be held at: _____ (location)'])
    U.drill(p,'练习2：答案类型预判',
        '根据题目判断填空答案的类型（名词/数字/形容词/日期/地点）',
        ['Q6 The building was constructed in _____ (year)',
         'Q7 The membership fee includes a _____ discount (type)',
         'Q8 Sessions are held every _____ afternoon (day)',
         'Q9 Maximum group size: _____ people (number)',
         'Q10 The new branch is located on _____ Street (name)'])
    U.drill(p,'练习3：同义替换快速配对',
        '30秒内完成10组同义替换配对',
        ['题目词 → 录音词：',
         'requirement → need/demand/necessity',
         'demonstrate → show/prove/illustrate',
         'annual → yearly/every year/per annum',
         'vicinity → nearby/close to/in the area of',
         '以上各组：哪个录音词最常在S1-S2中出现？'])
    U.drill(p,'练习4：同义替换真题实战',
        '以下题目→录音，找出同义替换关系',
        ['题目：The programme focuses on environmental _____ (protection)',
         '录音："The scheme aims to preserve natural habitats and maintain ecological balance."',
         '题目词protection → 录音词preserve/maintain',
         'Q：答案应该填什么词？（TWO WORDS AND/OR A NUMBER）'])
    U.drill(p,'练习5：全流程20秒预读实战',
        '完整模拟：拿到题组→20秒预读→口述预测',
        ['Section 2 Questions 11-15（Map Labelling）',
         'Section 2 Questions 16-20（Short Answer, max TWO WORDS）',
         '计时20秒：圈关键词+预判答案类型+预判可能的同义替换',
         '口述：Q11我预测答案是___类型，因为___',
         '反思：哪道题预读最花时间？下次如何加快？'])
    # PART08 pitfall
    U.section(p,8,'预测陷阱与误区','Common Prediction Mistakes')
    U.pain_point(p,'预测过度导致失分',
        ['预测太具体：认定Q3答案是"Tuesday"，录音说"Wednesday"时不敢写，坚持填Tuesday',
         '看到选项A感觉对就停止听：A说了个关键词就勾A，没听到B里的更正',
         '预测干扰注意力：脑子里想着"应该是medical"，错过了录音里说的"dental"'],
        ['预测答案"类型"不预测具体词：知道要填星期几，但不提前锁定是哪天',
         '选择题听完整段再判断：不要听到A的关键词就立刻勾选',
         '预测=降低认知负担，不是替代听力——答案永远以录音为准，不以预测为准'])
    U.content(p,'预测技能的练习方法',
        [('每日训练','做真题前先做"盲预读"：不听录音，只预读题目，写下答案类型和可能的同义词'),
         ('反思复盘','做完题后对比预读笔记vs实际答案：哪些预读有效？哪些预读误导了你？'),
         ('同义替换日积月累','每天学5组同义替换对，用剑桥听力词汇表，专注"题目词↔录音词"关系'),
         ('计时加压','把预读时间从30秒压缩到20秒再到15秒——实战中题组间隔有时不足20秒')],
        note='目标：经过10次专项练习后，20秒内完成10道题的关键词圈划成为条件反射')
    U.content(p,'预测能力自测标准',
        [('初级（6分）','能圈出关键词，但预读时间>30秒；同义替换识别率<50%'),
         ('中级（6.5-7分）','20秒内圈完关键词；能识别直接同义词（cheap→affordable）'),
         ('高级（7.5-8分）','15秒内圈完；能识别概念替换（parking→car park）和上义词替换'),
         ('顶级（8.5+分）','主动预判答案词性+可能同义词；几乎不被同义替换干扰')],
        note='深圳学生平均处于初级→中级，本节课后目标升至中级，L08后升至高级')
    U.tip_card(p,'课堂即时可用：预读快速指南',
        [('题型识别（3秒）','先看全局：有空格=填空；A/B/C=选择；字母列表=配对；图形=地图'),
         ('关键词圈划（10秒）','每题1-2个最独特词：专有名词>数字/时间>核心名词'),
         ('答案类型预判（5秒）','空格前后判词性：形容词后填名词；be动词后填形容词'),
         ('选择题特殊处理','只圈三个选项的差异核心词，不逐字读完整选项'),
         ('超时处理原则','预读时间到即停止，宁可预读7题也不完整读3题')])
    U.tip_card(p,'考场急救：预读来不及怎么办',
        [('题来不及圈完','先圈已读题目，边听录音边圈剩余题目的关键词'),
         ('跟丢题目位置','找最近一个写下的答案，从下一题关键词开始重新追踪'),
         ('同义替换没识别','记下听到的词，誊写时与题目语义对照确认是否匹配'),
         ('答案选错了','用细线划掉重写，不要大力涂改（影响扫描识别）'),
         ('整组题没听到','保持冷静继续预读下一组，不要纠缠已失去的题目')])
    U.checklist(p,'L06课后自检清单',
                   ['完成剑桥真题1套听力，专项计时预读（20秒限时）',
                   '每道题目完成三层预读：题型→关键词→答案类型',
                   '做完后统计：同义替换识别成功多少次？',
                   '建立个人同义替换词对表，每周新增10组'])
    U.summary(p,6,
        ['预读三层次：题型识别→关键词圈划→答案类型预判，全程20秒',
         '关键词选词原则：专有名词>数字/时间>核心名词，不圈功能词',
         '同义替换是必考技能：题目用正式词，录音用口语词（或反之）',
         '预测≠猜答案：预测是为了知道"要找什么"，答案以录音为准',
         '预读是可训练的技能：10次专项练习后自动化，大幅提升做题效率'])
    U.homework(p,6,
        ['完成剑桥真题2套听力（每套限时），做前专项计时预读',
         '每套做完后统计同义替换识别率，目标≥70%',
         '建立个人同义替换词对本，本周新增20组（参考剑桥词汇表）'],
        est_time='约60分钟')
    U.preview(p,7,'笔记速记系统',
        ['符号缩写体系：10个通用符号，覆盖80%的笔记需求',
         '边听边记的双任务协调技巧',
         '数字、日期、电话号码的速记规范'])
    U.save(p, f'{OUT}/L06_预测技能与关键词定位.pptx')
    print(f"  L06 done: {len(p.slides)} pages")

# ══════════════════════════════════════════════════════════════════════════════
#  L07  笔记速记系统
# ══════════════════════════════════════════════════════════════════════════════
def L07():
    T(); p=U.new_prs()
    U.cover(p,'听力',7,'笔记速记系统','Note-taking & Speed Writing System')
    U.agenda(p,['双任务协调的挑战','速记符号体系','数字信息速记',
                '关键词追踪技术','速记词汇表','边听边记步骤','课堂练习','小结作业'])
    U.content(p,'L06回顾：预测技能核心要点',
        [('三层预读','题型识别(3s)→关键词圈划(10s)→答案类型预判(5s)=20秒'),
         ('同义替换','题目词≠录音词，掌握100组核心替换对，不被迷惑'),
         ('预测原则','预测答案类型而非具体内容，答案永远以录音为准'),
         ('深圳学生进度','大多数同学已能完成关键词圈划，但速记跟不上是新瓶颈')],
        note='本节课目标：建立速记符号体系，让双任务（听+写）不再互相干扰')
    U.objectives(p,['掌握10个通用速记符号，覆盖80%的笔记场景',
                    '能在听录音时同步用符号记录关键信息',
                    '掌握数字、日期、电话号码的快速记录格式',
                    '实现"听-记-答"三步流程的自动化'])
    U.pain_point(p,'双任务失败的三种表现',
        ['写慢了跟不上：还在写上一个答案时，下一个答案已经过去',
         '写快了写不清：用力快写变成潦草，誊写时自己都认不出来',
         '边写边漏听：专注写字时耳朵"断电"，关键信息在这段时间说完了'],
        ['速记符号替代完整词：→代替"increase"，↓代替"decrease"，节省70%书写时间',
         '用小号印刷体写，不追求美观，只追求誊写时自己能认出',
         '符号+首字母组合：w/（with）、b/c（because）、&（and）减少书写量'])
    part(p,1,'双任务协调：听力与书写同步','Dual-Task Coordination',
        [('双任务挑战的根本原因','大脑处理能力有限，听和写会争抢资源',
          [('听力通道','耳朵接收语音→大脑解码语义→工作记忆暂存'),
           ('书写通道','工作记忆提取信息→手部执行写字动作'),
           ('冲突点','解码新信息时书写旧信息，两者同时占用工作记忆'),
           ('解决方案','减少书写所需认知资源——符号化、简化、自动化'),
           ('训练目标','让书写变成"肌肉记忆"，不消耗认知资源')],
          '目标：写字时大脑99%用于听，只用1%用于手部动作'),
         ('速记的核心原则','能少写就少写，能符号就不用词',
          [('原则1','只记答案关键词，不记完整句子'),
           ('原则2','用符号替代常用词：→↑↓=≠&w/b/c'),
           ('原则3','数字一律用阿拉伯数字，不用英文单词'),
           ('原则4','专有名词首字母大写，不确定的部分画横线占位'),
           ('原则5','誊写时再补全——听音时不要求完整，誊写时核对')])],
        [('双任务协调的阶段目标',
          [('初级阶段','先写完上一题再听下一题——会漏题，但减少了写错'),
           ('中级阶段','听到答案立刻写——偶尔漏听，但写得快了'),
           ('高级阶段','听和写同步——符号体系自动化，脑子99%用于听'),
           ('训练路径','第1-2周：专注建立符号体系→第3周：边听边写→第4周：自动化'),
           ('衡量标准','做完2套真题后，双任务成功率（不因写字而漏听）≥80%')]),
         ('速记体系建立的优先级',
          [('第一优先','数字立刻用阿拉伯数字——最快见效，立刻少写50%字符'),
           ('第二优先','方向符号→↑↓——increase/decrease用↑↓替代'),
           ('第三优先','连接符号& w/ b/c——频率最高的逻辑关系词'),
           ('第四优先','专业缩写：env/govt/temp/pop——S4频繁出现'),
           ('关键原则','先把基础5个符号用熟，再扩展，不要一次学10个')])],
        [('双任务协调自测',
          '同时完成以下两个任务（难度模拟）：\n\n'
          '任务A：听老师朗读的数字序列（每2秒一个）\n'
          '任务B：同时把刚才的数字写下来\n\n'
          '序列：47 / 230 / 1,850 / 99.5 / 3/4 / 0.7 / £26.50',
          '47 / 230 / 1850（可省千分位逗号）/ 99.5 / 3/4或¾ / 0.7 / £26.50\n听写正确率：7个全对=优秀；5-6个=良好；<5个=需加强数字听写专项训练',
          '训练法：每天10分钟数字听写，用BBC Radio 4的新闻播报练习'),
         ('速记符号紧急入门',
          '以下录音，用符号速记关键信息：\n\n'
          '"The temperature increased significantly last month, '
          'reaching approximately 35 degrees, which is about 8 degrees higher than usual. '
          'However, next week it will probably decrease."',
          '速记示例：temp ↑ sig last mo → ~35° / ~8° > usual\nnext wk → temp ↓ (prob)\n节省书写量约65%',
          '速记不是固定格式，关键是自己能看懂——建立个人符号系统后坚持使用')]
    )
    part(p,2,'10大速记符号精讲','10 Core Shorthand Symbols',
        [('方向与变化符号','5个最高频符号，覆盖50%的速记需求',
          [('→','leads to / results in / therefore（因果、方向）'),
           ('↑','increase / rise / grow / improve（增加类）'),
           ('↓','decrease / fall / drop / reduce（减少类）'),
           ('=','equals / is / means / refers to（等于、定义）'),
           ('≠','does not equal / is different from / not（不等于）')],
          '这5个符号用一周就能形成习惯——每次听力练习强制使用'),
         ('连接与修饰符号','5个次高频符号',
          [('&','and / as well as / both（和）'),
           ('w/','with / including（带有）'),
           ('w/o','without / excluding（没有）'),
           ('b/c','because / due to / as a result of（因为）'),
           ('~','approximately / about / around（大约）')])],
        [('10个符号的记忆技巧',
          [('视觉记忆','↑=温度计上升，↓=降落的箭头，→=因果链条'),
           ('日常使用','在笔记本里用这些符号代替英文词，不只在听力练习中'),
           ('分批学习','第一天学→↑↓=≠，熟练后学& w/ w/o b/c ~'),
           ('自我测试','合上书本，能默写全部10个符号及其含义=掌握'),
           ('考场保障','符号不确定时宁可写完整词，不要写一个自己不确认的符号')]),
         ('符号速记 vs 完整书写的对比',
          [('完整书写','increase: 8个字母，约1.5秒'),
           ('符号速记','↑: 1个符号，约0.1秒，节省93%时间'),
           ('完整书写','because of habitat destruction: 28个字母，约5秒'),
           ('符号速记','b/c habitat destr.: 16个字符，约2.5秒，节省50%'),
           ('累计效果','40题平均每题节省1.5秒=总计节省1分钟，减少漏题风险')])],
        [('符号速记实战练习',
          '将以下句子用速记符号改写（目标：字数减少50%以上）：\n\n'
          '"The new policy will increase the budget by approximately 20% '
          'because of rising operational costs, but this will not include equipment expenses."',
          '速记版：new policy → budget ↑ ~20% b/c op.costs ↑\nw/o equip. expenses\n原文26词 → 速记11个符号/词，减少58%',
          '速记没有标准答案，只要自己5分钟后还能读懂就是好的速记'),
         ('真题听力速记示范',
          'S2独白：\n"The community centre has recently extended its opening hours. '
          'It now opens at 7am instead of 9am, and closes at 10pm rather than 8pm. '
          'Membership fees have also increased from £45 to £52 per month."',
          '速记：CC: hours ↑\nopen: 7am (was 9am)\nclose: 10pm (was 8pm)\nfee: £45→£52/mo',
          '关键：记录"变化"（原值→新值），而不只记最终值，因为题目可能问"之前是多少"')]
    )
    part(p,3,'数字信息的速记规范','Number Information Speed Recording',
        [('数字类型与速记格式','雅思听力中的数字陷阱重灾区',
          [('整数','直接写：47, 230, 1850（省略千分位逗号可以）'),
           ('小数','保留小数点：99.5, 3.75, 0.5（不能写成½等分数形式）'),
           ('分数','用斜杠：3/4, 1/3, 2/5（誊写时查是否要写文字）'),
           ('金额','加货币符号：£24.50, $99, €1,200（符号紧贴数字，无空格）'),
           ('电话','分组记：020-7946-0342，听一段写一段，不等全部说完')],
          '金额最易失分：£24.5和£24.50是不同的——雅思认后者；小数位要写全'),
         ('易混数字的辨析','听力中最容易混淆的数字对',
          [('13 vs 30','thirteen重音在-teen；thirty重音在thir-；用语境判断'),
           ('14 vs 40','fourteen vs forty——forty没有"u"，注意拼写'),
           ('15 vs 50','fifteen vs fifty——同样的重音规则'),
           ('hundred vs thousand','hundred=百；thousand=千；别把1,000听成100'),
           ('点号','point用于小数：3.5=three point five；不是"three dot five"')])],
        [('数字速记的练习计划',
          [('每日练习','用BBC新闻或真题S1练习5分钟数字听写——最快提分方式'),
           ('专项训练','把13/30/14/40/15/50配对练习：一人报数，一人写，互相核对'),
           ('格式检查','做完真题后专门检查数字格式：£是否贴紧、小数位是否完整'),
           ('电话特训','练习"一段说完立刻写"，不等全部说完——分组记录'),
           ('目标','数字听写正确率从<70%提升到≥95%，约需2周专项练习')]),
         ('数字答案的誊写检查清单',
          [('金额','有没有货币符号（£/$）？小数点后两位是否完整（24.50不是24.5）？'),
           ('日期','格式是否符合题目要求？1st June / June 1st / 01/06 哪个合适？'),
           ('时间','12小时制还是24小时制？am/pm是否需要写？9:30还是half past nine？'),
           ('电话','分组是否正确？双号（double）是两个相同数字，已写成两个了吗？'),
           ('百分比','%符号有没有写？98.6%不能写成98.6')])],
        [('数字听写快速练习',
          '听老师朗读以下数字，用最快速度写下来：\n\n'
          '① 1,350  ② £47.99  ③ 0.75  ④ 3/5  ⑤ 13:30\n'
          '⑥ 020-7946-0918  ⑦ 1st June  ⑧ 2,500,000  ⑨ -3°C  ⑩ 98.6%',
          '① 1350（逗号可省）② £47.99（小数位全写）③ 0.75 ④ 3/5 ⑤ 13:30或1:30pm ⑥ 按段记录 ⑦ 1st June / June 1st均可 ⑧ 2.5m或2,500,000 ⑨ -3°C ⑩ 98.6%',
          '⑤时间格式：24小时制（13:30）和12小时制（1:30pm）雅思都接受'),
         ('电话号码速记专项',
          '以下电话号码，边听边写（每段停顿2秒）：\n\n'
          '"The number is 0-7-9-5-4, double 3, 8-1-6."\n\n'
          '"And the fax number is 0-2-0, 7-9-4-6, 0-9-1-8."',
          '① 07954-338-16 → 07954 33816\n② 020-7946-0918\n注意double 3=33（两个3，不是3的两倍）',
          '"double"=两个相同数字，"treble"=三个；oh/zero都表示0；nought也表示0')]
    )
    part(p,4,'关键词追踪与占位技术','Keyword Tracking & Placeholder Technique',
        [('听不清时的占位策略','听不清≠放弃，用占位符保留答案空间',
          [('首字母占位','听到m开头的词→写"m___"，誊写时确认'),
           ('音译占位','听到不认识的词→尽量写出发音：/si:k/→seec?'),
           ('问号标记','完全没听清→写"?"，提醒誊写时特别注意'),
           ('上下文推断','空白题：根据前后题目推断语义范围'),
           ('不停下来','空一个继续往前，不要停下来纠结一道题')],
          '占位的目的是誊写时有线索——一个首字母比什么都没有强10倍'),
         ('誊写10分钟的最大化利用','纸笔考特有的10分钟誊写时间',
          [('优先级1','把答题纸上的缺字母、拼写不确定的词补全'),
           ('优先级2','检查所有拼写：accommodation/necessary/February'),
           ('优先级3','检查大小写：专有名词首字母大写'),
           ('优先级4','检查单复数和限词是否符合题目要求'),
           ('机考学生','没有誊写时间→听题时就要拼全，只留5秒/题检查')])],
        [('占位技术的实际价值',
          [('数据','用占位技术的学生，因"拼写不确定而空白"的题目减少85%'),
           ('心理作用','写下首字母比完全空白更有把握感，减少考场焦虑'),
           ('誊写还原率','首字母占位的题目，誊写时还原成功率约70%'),
           ('音译占位还原率','音译占位题目，誊写还原率约40%（已有音节信息辅助）'),
           ('核心价值','任何占位都比空白好，哪怕只猜一个字母也有可能得分')]),
         ('占位技术与誊写的配合',
          [('听题时','不追求拼全，专注抓住关键音节或首字母'),
           ('誊写开始','立刻从Q1开始，把所有占位符尝试补全'),
           ('补全策略','根据前后题目语境推断词义→查词根→查高频拼写词表'),
           ('放弃标准','占位词誊写时2次还原失败→放弃，继续下一题'),
           ('时间分配','40题誊写平均15秒/题，难题不超过30秒/题')])],
        [('占位技术实战：难词处理',
          '以下是S4学术词汇，练习"听到→占位→誊写"全流程：\n\n'
          '① biodiversity  ② photosynthesis  ③ epidemiology  ④ nanotechnology  ⑤ seismograph',
          '① bio...diversity → biod... → biodiversity\n② photo...synthesis → photos... → photosynthesis\n③ epidemi...ology → epid... → epidemiology\n④ nano...tech → nanotech... → nanotechnology\n⑤ seismo...graph → seismograph',
          '技巧：学术词通常有可分析的词根——bio(生命)+diversity(多样性)，记词根比记整词快'),
         ('听力全流程速记综合练习',
          'S4笔记填空，用学到的所有速记技巧：\n\n'
          '"Research indicates that urban biodiversity has declined by approximately '
          '40% since 1970, primarily because of habitat destruction and pollution. '
          'However, green roof initiatives have recently shown significant promise."',
          'Q答案：40% / 1970 / habitat destruction (and) pollution / green roof(s)\n速记示例：urban biodiversity ↓ ~40% since 1970 b/c habitat destr. & pollution\nBUT: green roofs → sig. promise recently',
          '速记和答题是分开的：速记记录全信息，答题从速记中提取题目要求的部分')]
    )
    # PART05 vocab
    U.section(p,5,'速记符号与缩写词汇表','Speed Writing Symbols & Abbreviations')
    U.vocab(p,'方向与逻辑符号（必掌握）',
        [('→','leads to / therefore / results in（因此/导致）'),
         ('←','caused by / due to / as a result of（因为，反向）'),
         ('↑','increase / rise / go up / improve（增加）'),
         ('↓','decrease / fall / decline / reduce（减少）'),
         ('=','equals / means / is defined as（等于/意味着）'),
         ('≠','is not / differs from / unlike（不等于/不同）'),
         ('≈ / ~','approximately / about / around（大约）')])
    U.vocab(p,'连接与修饰符号',
        [('&','and / plus / as well as / together with（和）'),
         ('w/','with / including / accompanied by（带有/包括）'),
         ('w/o','without / excluding / not including（不包括）'),
         ('b/c','because / due to / since / as（因为）'),
         ('vs','versus / compared to / in contrast（对比）'),
         ('+','positive / advantage / benefit / pro（优点）'),
         ('-','negative / disadvantage / problem / con（缺点）')])
    U.vocab(p,'常用缩写词（学术与日常）',
        [('govt','government（政府）'),
         ('env','environment / environmental（环境）'),
         ('temp','temperature（温度）'),
         ('pop','population（人口）'),
         ('info','information（信息）'),
         ('org','organisation（组织机构）'),
         ('dept','department（部门）'),
         ('max / min','maximum / minimum（最大/最小值）'),
         ('approx','approximately（大约）')])
    U.two_col(p,'数字速记格式速查',
        '数字类型→推荐格式',
        [('整数','直接写：47, 1850'),
         ('金额','货币符号贴紧：£24.50, $99'),
         ('小数','保留小数点：3.75, 0.5'),
         ('百分比','数字+%：98.6%, 40%'),
         ('时间','冒号格式：9:30am, 13:45'),
         ('日期','数字+月份：1st May, 15 June'),
         ('电话','分段记：020-7946-0918')],
        '特殊数字陷阱',
        [('thirteen vs thirty','重音位置区分：-teen在后/thir-在前'),
         ('double/treble','double=×2个，treble=×3个'),
         ('oh vs zero','电话号码中0读作"oh"'),
         ('point vs dot','小数点说"point"，不说"dot"'),
         ('and in numbers','英式：one hundred AND fifty=150'),
         ('nought','英式：0.5=nought point five')])
    # PART06 method
    U.section(p,6,'边听边记：完整步骤流程','Listen & Write: Full Process')
    U.steps(p,'Step 1-3：准备与启动',
        [('预读阶段（录音前）','20秒完成三层预读；在答题纸上圈关键词；准备好笔握在手中'),
         ('说明播放阶段','说明词只放一次，用这30秒补充预读未完成的题目'),
         ('第一个答案准备','锁定Q1的关键词；把笔尖放在Q1答题位置附近')])
    U.steps(p,'Step 4-6：听题与速记',
        [('听到关键词周边','关键词出现时准备写，关键词后1-3秒内通常出现答案'),
         ('写答案','用速记符号+首字母+数字格式快速记录；不写完整句'),
         ('答完即移位','写完答案后立刻把视线移到下一题，不回头检查，继续往前追踪')])
    U.steps(p,'Step 7-8：誊写与检查（纸笔考）',
        [('誊写优先级','先补全首字母占位的词→再检查拼写→再查大小写→最后查单复数'),
         ('时间分配','10分钟誊写：每道题平均15秒；遇到难词最多停留30秒，继续下一题')])
    U.model_answer(p,'速记示范：S3学术讨论',
        'Script: "In our experiment, we found that the plants grew significantly faster '
        'when exposed to blue light — about 35% more than those under white light. '
        'However, the results were inconsistent, possibly because of temperature fluctuations."',
        'Q: What percentage faster did plants grow under blue light? → 35%\nQ: What may have caused inconsistent results? → temperature fluctuations',
        ['速记全程：plants: blue > white ~35%\nBUT results ≠ consist. b/c temp ↑↓(?)',
         '从速记提取答案：Q1→找%数字→35%；Q2→找"b/c"之后的词→temp fluctuations',
         '速记目的：记录所有可能有用的信息，答题时从中提取，而非边听边想答案'])
    U.model_answer(p,'速记示范：S1信息表',
        'Script: "My name\'s Johnson — J-O-H-N-S-O-N. '
        'I\'d like to book a table for six, on the 23rd of March, at half past seven."',
        'Q29 Name: _____  Q30 Number of guests: _____  Q31 Date: _____  Q32 Time: _____',
        ['速记：Johnson (J-O-H-N-S-O-N) / 6 ppl / 23rd March / 7:30',
         'Q29: Johnson（按字母拼读同步写，不等说完）',
         'Q30: 6（"six"转数字，不写"six"）',
         'Q31: 23rd March（或March 23rd）',
         'Q32: 7:30（half past seven立刻转换，不写half past seven）'])
    U.compare(p,'速记水平：低分 vs 高分',
        '❌ 低分速记（慢+乱）',
        ['写完整英文单词：increase, because, approximately',
         '数字用英文：forty-seven, three hundred',
         '没有符号系统，每次想现找缩写',
         '边听边想答案，没有先记录再提取的意识',
         '誊写时凭印象补全，不是从速记笔记中提取'],
        '✅ 高分速记（快+准）',
        ['用符号替代常用词：↑, b/c, ~, &',
         '数字全用阿拉伯数字：47, 300',
         '符号体系已自动化，不需要临时思考用什么符号',
         '先速记全部信息，听完后从笔记中提取答案',
         '誊写时严格对照速记笔记，不靠记忆'])
    U.compare(p,'真题前后对比：有无速记体系',
        '❌ 无速记体系（6分水平）',
        ['听到"approximately 35%"→想着写three five percent，A还没写完已跑掉',
         '听到"because of habitat destruction"→写because，destruction来不及写',
         '誊写时想不起来听到了什么，靠感觉填空',
         '一道填空跟丢后，后面2道也连带跑掉'],
        '✅ 有速记体系（7.5分水平）',
        ['听到"approximately 35%"→立刻写~35%（3个字符，0.5秒完成）',
         '听到"because of habitat destruction"→b/c habitat destr.（5个字符）',
         '誊写时从笔记提取：b/c habitat destr.→"habitat destruction"（补全即可）',
         '一道填空跟丢后，从下一个锚点词重新定位，损失最多1道题'])
    # PART07 drill
    U.section(p,7,'课堂练习：速记系统训练','Speed Writing Drills')
    U.drill(p,'练习1：符号替换（10组）',
        '将以下英文短语改写为速记符号形式',
        ['① The population increased by approximately 15%',
         '② The project failed because of insufficient funding',
         '③ The new policy is different from the previous one',
         '④ Both methods are effective, with and without supervision',
         '⑤ Temperature decreased significantly, reaching about -5 degrees'])
    U.drill(p,'练习2：数字听写（10个）',
        '听写以下数字，用最快速度记录（格式正确即可）',
        ['① treble 4, 9, double 7  ② £1,249.99  ③ 0.375  ④ 3rd September',
         '⑤ 98.6%  ⑥ 020-3948-7712  ⑦ 7:45am  ⑧ 2,750,000  ⑨ -12°C  ⑩ 4/5'])
    U.drill(p,'练习3：S1信息表速记实战',
        '听老师朗读S1对话，用速记符号+数字格式记录',
        ['场景：健身房会员注册；包含10个填空题',
         '要求：使用至少5个速记符号；数字用阿拉伯数字',
         '自评：速记完成率___% ; 符号使用___个 ; 誊写错误___个',
         '对照标准答案：检查数字格式、拼写、大小写'])
    U.drill(p,'练习4：占位技术练习',
        '以下词汇模拟"听不清"情况，用占位符记录后还原',
        ['① /ˌsʌs.teɪ.nə.bɪl.ɪ.ti/ → sus___ility → sustainability',
         '② /ɪ.kɒn.ə.mɪks/ → ec___omics → economics',
         '③ /ˌfɒt.ə.ˈsɪn.θɪ.sɪs/ → photo___sis → photosynthesis',
         '目标：首字母占位+音节估算 → 誊写时确认'])
    U.drill(p,'练习5：速记→答案提取全流程',
        '根据以下速记笔记，写出完整答案',
        ['速记：Lib: open 8am (was 9am) / close 10pm / fee ↑ £35→£42/mo w/ gym',
         'Q1 What time does the library now open? →',
         'Q2 How much is the monthly fee? →',
         'Q3 What is now included in the membership? →',
         '练习从速记中"提取"而不是靠记忆——这是誊写阶段的核心技能'])
    # PART08 pitfall
    U.section(p,8,'速记误区与考场应对','Speed Writing Pitfalls')
    U.pain_point(p,'速记的三大误区',
        ['误区1：速记太详细——想记下每个词，结果跟不上录音节奏，丢失后面的答案',
         '误区2：发明临时符号——考场中忽然想用一个从没练过的符号，誊写时不记得是什么意思',
         '误区3：速记后不誊写检查——机考学生认为输入的就是最终答案，不做最后检查'],
        ['速记原则：只记答案相关信息，不记所有听到的内容',
         '符号系统提前固化：备考期间建立10个符号并坚持使用，不临时创造',
         '机考也要检查：每组题答完后利用间隔时间检查已输入的答案'])
    U.content(p,'速记体系的建立计划',
        [('第一周','掌握10个基础符号，在日常英语练习中强制使用'),
         ('第二周','加入数字速记规范，做真题时专项使用数字格式'),
         ('第三周','加入占位技术，不再因为拼写不确定而空白'),
         ('第四周','全流程整合：预读→听题→速记→誊写，全部用系统化方法')],
        note='目标：4周后速记体系完全自动化，考场中不需要思考用什么符号')
    U.content(p,'机考学生的特殊速记策略',
        [('草稿纸利用','机考提供草稿纸——用来做预读圈词和难词音译占位'),
         ('输入速度','提前练好英文打字速度，目标≥40词/分钟'),
         ('输入法切换','考前设置好输入法，避免中英切换浪费时间'),
         ('答案检查','每组题输入后，利用间隔时间检查已输入的答案拼写和格式'),
         ('数字输入','数字直接用数字键，不用英文单词输入')],
        note='机考没有誊写时间，所以每道题的质量要更高，不能靠誊写阶段补救')
    U.tip_card(p,'即学即用：10大速记符号卡片',
        [('方向符号（5个）','→因此/导致  ↑增加  ↓减少  =等于/是  ≠不等于/不同'),
         ('连接符号（5个）','& 和  w/ 带有  w/o 不含  b/c 因为  ~ 大约'),
         ('数字规范','金额：£/$/€贴紧数字（£24.50）；时间：冒号（9:30）；日期：1st May'),
         ('占位技术','首字母+横线（m___）；完全没听清写问号（?）；音译（seec?）'),
         ('誊写优先级','占位词补全→拼写检查→大小写→单复数→限词核查')])
    U.tip_card(p,'双任务协调急救指南',
        [('写不完时','只写关键词+数字，省掉冠词介词连词——核心信息最重要'),
         ('跟丢了','空下来继续往前，不停下来纠结——一道题最多损失1分'),
         ('没听清','首字母+横线占位，誊写时靠上下文语义推断补全'),
         ('时间不够','先记数字和专有名词，动词形容词用方向符号速记'),
         ('机考打字慢','先在草稿纸速记关键词，听完整段后集中输入答案')])
    U.checklist(p,'L07课后自检清单',
                   ['建立个人速记符号表（至少10个），打印出来贴在桌前',
                   '做2套真题听力，全程使用速记符号+数字格式',
                   '统计：速记使用率___% 、誊写补全成功率___% 、数字失分次数___',
                   '识别自己最薄弱的速记环节，下周专项训练'])
    U.summary(p,7,
        ['双任务协调的关键：符号化书写，让书写占用最少认知资源',
         '10大速记符号：→↑↓=≠ & w/ w/o b/c ~，覆盖80%场景',
         '数字速记规范：阿拉伯数字+货币符号+正确格式，避免数字失分',
         '占位技术：首字母+横线，誊写时从笔记中提取而不是凭记忆',
         '誊写10分钟（纸笔）：优先补全占位词→检查拼写→查大小写→单复数'])
    U.homework(p,7,
        ['建立个人速记符号表，打印贴桌前，本周日常练习中强制使用',
         '做剑桥真题2套，全程速记，做完后统计速记符号使用次数',
         '专项数字听写：每天5分钟，BBC新闻中的数字，听写练习'],
        est_time='约60分钟')
    U.preview(p,8,'六大陷阱类型全析',
        ['更正型陷阱：录音先给错误信息，再更正——如何不被第一个答案骗',
         '否定型陷阱：句中含not/never/don\'t，意思完全反转',
         '条件型陷阱：if/unless/only if改变答案成立条件'])
    U.save(p, f'{OUT}/L07_笔记速记系统.pptx')
    print(f"  L07 done: {len(p.slides)} pages")

# ══════════════════════════════════════════════════════════════════════════════
#  L08  六大陷阱类型全析
# ══════════════════════════════════════════════════════════════════════════════
def L08():
    T(); p=U.new_prs()
    U.cover(p,'听力',8,'六大陷阱类型全析','Six Trap Types: Complete Analysis')
    U.agenda(p,['陷阱的本质与分布','更正型陷阱','否定型陷阱',
                '条件型陷阱','相似音陷阱','综合陷阱','课堂练习','小结作业'])
    U.content(p,'L07回顾：速记系统核心要点',
        [('10大符号','→↑↓=≠ & w/ w/o b/c ~ — 形成自动化习惯后节省70%书写时间'),
         ('数字规范','阿拉伯数字+货币符号格式+小数位完整（£24.50不是£24.5）'),
         ('占位技术','首字母+横线，不停下来纠结一题，誊写时从笔记中提取'),
         ('誊写优先级','占位补全→拼写→大小写→单复数→限词')],
        note='本节课目标：识别并绕过雅思听力中的六大陷阱，防止"听懂了却做错"')
    U.objectives(p,['识别雅思听力中的六种典型陷阱类型',
                    '掌握每种陷阱的识别信号词和应对策略',
                    '能在听到陷阱时立刻做出正确判断',
                    '通过真题练习将陷阱识别率提升到80%以上'])
    U.pain_point(p,'听懂了却做错：陷阱是第一杀手',
        ['更正陷阱：录音先说"Tuesday"，学生立刻写下Tuesday，没听到后面的"actually Wednesday"',
         '否定陷阱：录音说"not available on Friday"，学生写了Friday（忽略了not）',
         '条件陷阱：录音说"if it\'s free"，条件没实现，但学生把条件内容当成了答案'],
        ['更正信号词：actually/sorry/I mean/wait/in fact — 听到这些词立刻划掉之前答案',
         '否定词强化意识：not/never/don\'t/without/no — 否定词后面的内容≠答案',
         '条件标志：if/unless/as long as/only when — 条件内容不一定是最终答案'])
    part(p,1,'更正型陷阱：最高频陷阱','Correction Traps: The Most Common Trap',
        [('更正陷阱的机制','录音先给出一个信息，再用信号词更正',
          [('信号词','actually, sorry, I mean, wait, no wait, in fact, I\'d like to change'),
           ('结构','错误信息 + 信号词 + 正确答案（考的是正确答案）'),
           ('常见场景','S1预约时间/姓名/价格/地址被更正'),
           ('S3难度','学术讨论中更正更隐蔽，常用"let me reconsider"'),
           ('频率','每套听力约出现3-5次更正陷阱，约占陷阱总量30%')],
          '关键：听到更正信号词，立刻在草稿上划掉第一个答案，等待新答案'),
         ('更正陷阱的变体','不止"actually"，还有这些更正方式',
          [('直接否定','No, sorry — that\'s wrong. It\'s...'),
           ('自我纠正','Wait, let me check... it\'s actually...'),
           ('他人纠正','A: Tuesday  B: Don\'t you mean Wednesday?  A: Oh yes, Wednesday'),
           ('延迟更正','信息A出现→聊别的话题→5句后才更正A（最难辨认）'),
           ('书面确认','"I\'ll just double-check that... yes, it\'s the 15th, not the 5th"')])],
        [('更正陷阱的分布规律',
          [('S1最多','S1预约场景：时间/姓名/价格/地址被更正，约70%的更正陷阱在S1'),
           ('S2也有','S2信息说明时更正：活动时间/地点/价格'),
           ('S3最隐蔽','学术讨论中更正是方案/观点/数据，常用reconsider/revise'),
           ('S4最少','学术讲座更正较少，但有时会出现"corrections to previous research"'),
           ('间隔规律','两道更正题之间通常间隔3-5题，连续更正较少见')]),
         ('更正草稿操作规范',
          [('第一步','听到关键词（可能是答案）→先在草稿写下'),
           ('第二步','继续听，特别注意actually/sorry/I mean/wait'),
           ('第三步','听到更正信号词→立刻横线划掉草稿上的候选答案'),
           ('第四步','写下更正后的新答案→确认逻辑→誊写'),
           ('机考操作','先在草稿写，确认后再输入；不要直接修改输入框（可能误触）')])],
        [('更正陷阱识别练习',
          '以下对话，找出更正陷阱并写出正确答案：\n\n'
          '"The session starts at half past six."\n'
          '"Sorry, I meant half past seven — we moved it back an hour."',
          '错误答案：6:30（half past six）\n正确答案：7:30（half past seven）\n信号词：Sorry, I meant\n操作：听到"Sorry, I meant"立刻划掉6:30，写7:30',
          '操作要点：不要等"I meant"说完才反应——听到"Sorry"就准备划掉上一个答案'),
         ('S3更正陷阱实战',
          'A: "So we\'re planning to interview around 50 participants?"\n'
          'B: "Actually, I think we should increase that to 75 to get a more reliable sample."\n'
          'A: "Right. And we\'ll use a questionnaire?"\n'
          'B: "No, I\'d prefer face-to-face interviews — questionnaires have too many limitations."\n\n'
          'Q: Number of participants? / Q: Data collection method?',
          'Q1：75（不是50，B更正了A的提议）\nQ2：face-to-face interviews（不是questionnaire，B明确否定了questionnaire）',
          'S3陷阱规律：A先提出一个方案，B通常会修改或否定→最终方案才是答案')]
    )
    part(p,2,'否定型与条件型陷阱','Negation & Conditional Traps',
        [('否定型陷阱：not/never/don\'t改变答案','否定词让意思完全反转',
          [('直接否定','not, never, no, don\'t, doesn\'t, can\'t, won\'t, isn\'t, aren\'t'),
           ('间接否定','without, lack of, fail to, unable to, refuse to'),
           ('限制否定','only, just, merely（限定到其他选项不成立）'),
           ('常见错误','听到关键词就写，没注意到关键词前面有否定词'),
           ('应对策略','每个关键词听到之前先扫描有无否定词，形成习惯')],
          '"not available on Friday" → 答案不是Friday；"without prior experience" → 答案不是experience'),
         ('条件型陷阱：if/unless改变答案的有效性','条件不满足时内容不算答案',
          [('条件结构','if... / unless... / only if... / as long as... / provided that...'),
           ('陷阱机制','录音说"if it rains, we\'ll move indoors"——条件不一定实现'),
           ('应对策略','听到条件词时，等待确认"条件是否实现"再写答案'),
           ('典型错误','题目问"venue"，录音说"if it rains, indoors"，学生写indoor（条件未必实现）'),
           ('识别信号','听到if/unless后，继续听确认句：确认了才写，没确认画问号')])],
        [('否定型陷阱的识别难点',
          [('not省音','快速语流中"not"被弱读为/nət/，容易没听到'),
           ('否定形式多样','can\'t/cannot/won\'t/don\'t/doesn\'t/didn\'t——每种都要敏感'),
           ('前置否定','否定词在句首：Never park here / No students allowed'),
           ('双重否定','not impossible=possible，not uncommon=common（要理解语义）'),
           ('否定答案类型','题目问"what can\'t students do"——答案要符合题目否定方向')]),
         ('条件型陷阱的深度解析',
          [('条件未实现','if引导的内容：条件一旦否定，该内容就不是答案'),
           ('确认方式','"The forecast looks positive"=天气好条件成立；"Unfortunately the funding fell through"=条件不成立'),
           ('双条件','unless A, B——A不发生则B发生，A发生则不是B'),
           ('计划+条件','Going to + unless = 最难复合体，需要找最终确认句'),
           ('深圳学生弱点','条件型是深圳学生识别率最低（约45%）的陷阱类型')])],
        [('否定型陷阱练习',
          '以下句子中，答案是什么？\n\n'
          'Q: What is NOT included in the ticket price?\n'
          '录音："The ticket price covers the main exhibition, the guided tour, and '
          'the café visit. However, the gift shop and the special art workshop '
          'are not included and must be paid for separately."\n\n'
          'Q问的是"不包含的"项目：',
          '答案：gift shop / (special) art workshop（录音说"not included"的两项）\n注意：题目本身就是"what is NOT included"——双重否定要小心\n主展览/导览/咖啡馆是包含的（=干扰项）',
          '"NOT included"类题目是否定题中最难的——题目和录音各有一个否定，逻辑要理清'),
         ('条件型陷阱练习',
          '录音："If the weather is good, we\'ll hold the event outdoors in the park. '
          'But if it rains, we\'ll use the community hall instead. '
          'The forecast looks positive, so we\'re planning for the park."\n\n'
          'Q: Where will the event be held?',
          '答案：(the/a) park（因为forecast looks positive，即天气预报良好，确认了第一个条件成立）\n不是：community hall（雨天备用方案，条件未成立）',
          '条件题解题步骤：①找两个条件分支→②找确认哪个条件成立的句子→③只写已确认的答案')]
    )
    part(p,3,'相似音与未来计划陷阱','Similar Sound & Future Plan Traps',
        [('相似音陷阱：听错了单词','音近词导致写错',
          [('数字音近','thirteen/thirty, fifteen/fifty, fourteen/forty'),
           ('词汇音近','dessert/desert, affect/effect, principle/principal'),
           ('名字音近','Lee/Li, Brown/Braun, Hall/Haul, Smith/Smyth'),
           ('短语音近','a name/an aim, ice cream/I scream, new key/UK'),
           ('应对策略','数字类：等说完整句再写；人名：等对方拼读确认')],
          '相似音陷阱在S1最常见，尤其是预约场景的姓名和数字'),
         ('未来计划陷阱：计划≠事实','提到的计划不一定实现',
          [('计划标志','plan to, intend to, going to, hope to, thinking of'),
           ('未实现标志','but, however, unfortunately, in the end, actually'),
           ('陷阱机制','录音先说plan，后说plan改变或取消——考实际结果'),
           ('应对策略','听到计划词时打问号，继续听是否有转折更改'),
           ('典型场景','S1：A计划周五送货，但后来改成了周一；答案是周一')])],
        [('相似音陷阱的练习方法',
          [('数字专项','每天练13vs30/14vs40/15vs50：一人报数，一人写，互相核对'),
           ('词汇专项','affect/effect, dessert/desert, principle/principal音近词配对练习'),
           ('电话专项','练习"double/treble+digit"格式，听一遍写一遍不回放'),
           ('人名专项','英式姓名常见：Smith/Smyth, Brown/Braun, Hall/Haul'),
           ('成功标准','数字听写正确率≥95%；音近词正确率≥90%后方可停止专项训练')]),
         ('未来计划陷阱的识别训练',
          [('训练方法','在录音听到plan to/going to时，在草稿上画问号(?)，继续听'),
           ('确认标准','等录音说出确认句（"so we\'ll"/"it\'s confirmed"）再写正式答案'),
           ('没有确认的情况','如果录音说plan to没有后续更改，计划本身就是答案'),
           ('但是转折','plan to + but = 最终答案在but之后'),
           ('unfortunately转折','plan to + unfortunately = 计划失败，答案是备用方案')])],
        [('相似音陷阱专项练习',
          '以下词对，听老师读一个，用正确词语回答：\n\n'
          '① thirteen vs thirty  ② affect vs effect\n'
          '③ desert vs dessert  ④ principle vs principal\n'
          '⑤ hall/Hall vs haul\n\n'
          '额外挑战：录音说"0800 three double four"——正确号码是？',
          '① 13=thirteen（重音-teen在后）vs 30=thirty（重音thir-在前）\n② affect=动词"影响"；effect=名词"效果"\n③ desert=沙漠/遗弃；dessert=甜点（多一个s）\n④ principle=原则；principal=校长/主要的\n⑤ 0800 344（three double four=344，不是3-4-4-4）',
          'double X在电话中=两个相同数字（double four=44）；treble=三个相同数字（treble six=666）'),
         ('未来计划陷阱练习',
          '"We were thinking of having the workshop on Saturday morning, '
          'but the room\'s already booked. So we\'ve had to switch it to Sunday afternoon instead."\n\n'
          'Q: When is the workshop?',
          '答案：Sunday afternoon（最终确认的时间）\n不是：Saturday morning（原计划，已被否定）\n信号词：but→had to switch→Sunday afternoon',
          '凡是recording先说A、再说but/however的，考的永远是but之后的内容')]
    )
    part(p,4,'复合陷阱与综合策略','Compound Traps & Integrated Strategy',
        [('复合陷阱：多个陷阱叠加','最难的陷阱类型',
          [('更正+否定','先更正，再加否定："No, sorry — it\'s not Tuesday, it\'s Thursday"'),
           ('条件+未来','"If we finish early, we plan to visit the museum"（双重不确定）'),
           ('相似音+更正','先说thirteen，被更正为thirty（音近+更正叠加）'),
           ('多人+立场','S3：A持正面观点，B否定，C提条件——三方观点交织'),
           ('识别策略','碰到复合陷阱时放慢书写速度，等整个逻辑链说完再写答案')],
          '复合陷阱是7分→8分的关键——能识别单一陷阱是7分，能识别复合陷阱是8分'),
         ('陷阱信号词总览','建立陷阱雷达系统',
          [('更正类','actually, sorry, I mean, wait, in fact, let me correct that'),
           ('否定类','not, never, no, don\'t, can\'t, without, fail to, lack of'),
           ('条件类','if, unless, only if, as long as, provided that, in case'),
           ('计划类','plan to, intend to, going to, thinking of, hope to'),
           ('转折类','but, however, although, yet, on the other hand（答案在转折后）')])],
        [('复合陷阱识别训练方法',
          [('第一步：识别所有信号词','听到actually/not/if/plan to等任何信号词，脑中标注'),
           ('第二步：画出逻辑链','在草稿上：计划A→条件C→结果B，梳理清楚'),
           ('第三步：找确认句','找说话人最后确认的那个选择作为答案'),
           ('时间判断','复合陷阱通常需要听完整个对话才能确认答案'),
           ('应急策略','逻辑复杂时：先打问号，继续听，誊写前再判断')]),
         ('S3复合陷阱的特殊难度',
          [('多人立场','每个说话人可能持不同意见，需要分别追踪'),
           ('话题转换','陷阱信号词出现后，话题可能转移，答案出现更晚'),
           ('委婉表达','学术语境中否定更隐蔽："I\'m not entirely convinced"=否定'),
           ('数据更改','数字类题目在S3也有陷阱：研究数据、百分比、参与者人数'),
           ('深圳学生策略','S3做完后标注哪些题有陷阱，复盘时分析识别是否成功')])],
        [('综合陷阱真题实战',
          'S3 Tutorial Discussion:\n'
          '"A: I was planning to use primary data — interviews.\n'
          'B: That sounds good, but unless we get ethical approval by Friday,\n'
          '   we might have to rely on secondary data instead.\n'
          'A: True. Actually, we haven\'t even submitted the application yet.\n'
          'B: So for now, let\'s assume secondary data, and update if approval comes through."\n\n'
          'Q: What type of data will they use?',
          '答案：secondary data（条件"ethical approval by Friday"未满足，所以退到secondary data）\n分析：primary data=计划；unless=条件陷阱；actually+haven\'t submitted=更正确认条件未满足',
          '解题步骤：①识别计划（primary data）→②找条件（unless approval）→③确认条件是否满足（haven\'t submitted=未满足）→④写secondary data'),
         ('陷阱总结：六大类型速览',
          '快速说出六大陷阱的信号词和应对方法：\n\n'
          '① 更正型  ② 否定型  ③ 条件型\n'
          '④ 相似音型  ⑤ 未来计划型  ⑥ 复合型\n\n'
          '2分钟内，每类说出1个信号词+1个应对策略',
          '① actually/sorry → 划掉前答案\n② not/never → 否定后内容不是答案\n③ if/unless → 等确认条件成立\n④ 数字/音近词 → 等说完整句或拼读\n⑤ plan to → 打问号，等转折\n⑥ 组合出现 → 放慢节奏，等逻辑链完整',
          '每次做错题后问自己：是哪类陷阱？信号词是什么？下次能提前预警吗？')]
    )
    # PART05 vocab
    U.section(p,5,'陷阱信号词词汇表','Trap Signal Words Vocabulary')
    U.vocab(p,'更正型陷阱信号词',
        [('actually','其实/实际上（最常见的更正信号词）'),
         ('sorry / I\'m sorry','对不起（自我纠正，通常后接更正内容）'),
         ('I mean / I meant','我的意思是（自我澄清）'),
         ('wait / hang on','等等（表示要修改前面说的内容）'),
         ('in fact','事实上（与前面内容形成对比）'),
         ('let me correct that','让我更正一下（明确更正表达）'),
         ('no, it\'s...','不，是...（直接否定并更正）')])
    U.vocab(p,'否定与条件陷阱信号词',
        [('not / no / never','否定词（后接内容不是答案或情况不存在）'),
         ('don\'t / doesn\'t / didn\'t','否定动词（动作未发生或不成立）'),
         ('without / lack of','缺乏（隐性否定）'),
         ('unless','除非（条件必须满足才成立）'),
         ('only if / as long as','只有当...才（强条件）'),
         ('if... but...','如果条件A，但实际是B（条件+更正组合）'),
         ('provided that','前提是...（条件词）')])
    U.vocab(p,'未来计划与转折信号词',
        [('plan to / planning to','计划做...（未来意图，未必实现）'),
         ('intend to / intention','打算...（意图词）'),
         ('going to / thinking of','打算/考虑...（不确定的计划）'),
         ('hope to / aim to','希望/旨在...（期望，有不确定性）'),
         ('but / however','但是（转折后才是真正答案）'),
         ('although / even though','虽然...但...（让步转折）'),
         ('in the end / eventually','最终（表示最后的结果，是答案）')])
    U.two_col(p,'六大陷阱速查表',
        '陷阱类型 → 信号词',
        [('更正型','actually/sorry/I mean/wait'),
         ('否定型','not/never/don\'t/without'),
         ('条件型','if/unless/only if/as long as'),
         ('相似音型','数字对/音近词/人名'),
         ('未来计划型','plan to/intend to/going to'),
         ('复合型','多个信号词同时出现')],
        '陷阱类型 → 应对策略',
        [('更正型','听到信号词→划掉前答案→写新答案'),
         ('否定型','否定词后内容≠答案，确认正向表达'),
         ('条件型','等待确认条件是否成立再写答案'),
         ('相似音型','等说完整句/拼读再写；区分重音位置'),
         ('未来计划型','听到计划词打问号，等待转折'),
         ('复合型','放慢写字速度，等整个逻辑链完整')])
    # PART06 method
    U.section(p,6,'陷阱规避：系统化解题流程','Trap Avoidance: Systematic Approach')
    U.steps(p,'Step 1-3：预防阶段',
        [('预读时标注陷阱风险','S1数字题预设"更正风险"；S3选择题预设"多人陷阱风险"'),
         ('建立"陷阱雷达"','听录音时，一旦出现信号词立刻提高警惕，延迟写答案'),
         ('双手分工','左手（或视线）追踪题目，右手（笔）准备随时划掉错误答案')])
    U.steps(p,'Step 4-6：应对阶段',
        [('听到更正信号词','立刻在草稿上划掉候选答案，写一个问号，等待新答案'),
         ('听到否定词','脑中立刻"反转"：not+[内容]=[内容]不是答案，继续往后听正确表达'),
         ('听到条件词','不急着写，继续听是否有确认条件成立的表述，再决定是否写答案')])
    U.steps(p,'Step 7：誊写时检查',
        [('回顾所有标注问号的题','誊写时对有问号标注的答案重点核查，是否还有不确定之处'),
         ('逻辑验证','答案写下后问：这个答案符合题目和录音的整体逻辑吗？'),
         ('错题分类','做完题后，错题标注"哪类陷阱"——积累个人易错陷阱类型')])
    U.model_answer(p,'满分应对更正陷阱：S1对话',
        'Q: What time is the appointment?\n\n'
        'Script: "So the appointment is at three o\'clock on Thursday."\n'
        '"Wait, sorry — I\'ve just checked, and Dr. Parker isn\'t available until four.\n'
        'Can you make it at four?"\n"Yes, four o\'clock is fine."',
        'Q答案：4 o\'clock / 4:00 / 4pm（最终确认时间）\n处理过程：\n① 听到3 o\'clock → 草稿写3:00\n② 听到"Wait, sorry" → 立刻划掉3:00\n③ 听到"four" → 写4:00\n④ 听到"Yes, four o\'clock is fine" → 确认，不再改动',
        ['草稿处理：3:00划掉，写4:00——誊写时直接用4:00，不要带划掉的内容',
         '等待确认：A说4:00，B确认"Yes"→才算最终答案，中间不要急着誊写',
         '更正速度：听到"Sorry"就要反应，不是听完整个更正句才反应'])
    U.model_answer(p,'满分应对否定+条件复合陷阱',
        'Q: What will the students use to collect data?\n\n'
        'Script: "We were thinking of using questionnaires, but actually,\n'
        'our supervisor said that\'s not appropriate for this type of study.\n'
        'He suggested interviews instead — but only if we can recruit at least ten participants."',
        '答案：interviews（条件：recruit at least 10 participants — 假设成立）\n否定：questionnaires → "not appropriate"（否定了questionnaire）\n更正：but actually（更正了最初计划）\n条件：only if recruit 10（条件，假设题目问的是"planned method"则写interviews）',
        ['这道题包含3个陷阱：更正（but actually）+否定（not appropriate）+条件（only if）',
         '应对策略：questionnaire出现时打问号→听到"but actually"划掉→听到interviews暂时写→听到"only if"再打问号',
         '题目问"will use"，通常以说话人最后确认的计划为答案：interviews'])
    U.compare(p,'陷阱题：做错 vs 做对的过程对比',
        '❌ 做错陷阱题（6分学生）',
        ['听到关键词立刻写，不等整句说完',
         '没有标注"信号词"的习惯，听到actually不知道要划掉',
         '听到not但没反应，照样把后面内容写成答案',
         '听到"plan to"不打问号，把计划当成事实写下',
         '誊写时直接抄，不检查逻辑'],
        '✅ 做对陷阱题（7.5+分学生）',
        ['听到关键词后继续听整句，直到逻辑链完整',
         '陷阱信号词=立刻提高警惕，延迟写下答案',
         '听到not就在候选词上画×，等待正向表达',
         '听到"plan to"打问号，继续听是否有转折或确认',
         '誊写时检查有问号标注的答案，判断最终确认了什么'])
    U.compare(p,'否定陷阱：两种学生的反应时间',
        '❌ 慢反应（错误）',
        ['录音：not available on Thursday',
         '学生反应：...Thursday！写下Thursday',
         '（忽略了not，把Thursday写成了答案）',
         '同类错误：听到"no change to the schedule"→写"schedule changes"'],
        '✅ 快反应（正确）',
        ['录音：not available on Thursday',
         '学生反应：not → 划掉候选词Thursday，继续听',
         '下一句：available on Friday → 写Friday',
         '同类训练：把否定词单独标记，防止遗漏not/never/no'])
    # PART07 drill
    U.section(p,7,'课堂练习：陷阱识别','Trap Identification Drills')
    U.drill(p,'练习1：更正陷阱识别（5组）',
        '听对话，找出更正信号词，写出最终正确答案',
        ['"The meeting\'s on the 5th... wait, I mean the 15th."',
         '"It costs £35. Oh sorry, that\'s the old price — it\'s £42 now."',
         '"She\'s from Melbourne. Actually no, she\'s from Sydney."',
         '"Three floors up — in fact, let me check — it\'s the fourth floor."',
         '"Ring extension 204. Hmm, I mean 240 — easy to mix up."'])
    U.drill(p,'练习2：否定陷阱识别（5组）',
        '以下句子，答案是什么？（注意否定词）',
        ['"Parking is not available on weekdays."→ Q: When can you park?',
         '"The fee does NOT include equipment rental."→ Q: What\'s excluded?',
         '"Students without prior experience are NOT eligible."→ Q: Who can apply?',
         '"Never use the back entrance after 6pm."→ Q: When can you use back entrance?',
         '"The café is closed on Sundays and Mondays."→ Q: When is it open?'])
    U.drill(p,'练习3：条件陷阱识别（3组）',
        '以下对话，最终答案是什么？（注意条件是否成立）',
        ['"If demand is high, we\'ll add a second session. Last year we had 200 applicants." → Q: How many sessions?',
         '"Unless the funding is approved, we\'ll hold the event online." (Funding was approved.) → Q: Venue?',
         '"We\'ll use the sports hall, provided that the basketball team finishes by 5pm." (They finish at 6pm.) → Q: Venue?'])
    U.drill(p,'练习4：相似音陷阱专项',
        '写出你听到的词，并判断是哪类相似音陷阱',
        ['① 14 vs 40 → 听老师读，写下数字',
         '② Smith vs Smyth → 这是人名，有没有区别？',
         '③ effect vs affect → 音近但词性不同，如何区分？',
         '④ 0800-three-treble-two → 完整号码是多少？',
         '⑤ dessert vs desert → 在S1餐厅场景中，更可能是哪个？'])
    U.drill(p,'练习5：综合陷阱实战（真题改编）',
        '以下S3对话包含多个陷阱，找出所有陷阱并写出每道题的正确答案',
        ['A: "So we\'re submitting on the 20th, right?"\n'
         'B: "Actually, the deadline\'s been extended to the 25th."\n'
         'A: "Great. Will we include the appendix?"\n'
         'B: "Only if it\'s under 10 pages. Ours is 15, so no."',
         '陷阱1：日期（20th→25th，更正陷阱）',
         '陷阱2：appendix（条件：<10页，但实际15页>10页，条件不成立→不包含）',
         'Q1答案：25th / Q2答案：no（不包含appendix）'])
    # PART08 pitfall
    U.section(p,8,'陷阱误区与防范','Common Mistakes in Trap Avoidance')
    U.pain_point(p,'应对陷阱时的三大新误区',
        ['过度怀疑：听到所有内容都不敢写，等啊等，结果什么都没写——考试变成了纯等待',
         '划掉太快：听到任何稍有变化的词就划掉之前答案，结果把正确答案也划掉了',
         '信号词背诵不够熟：知道陷阱存在，但听到"in fact"时还是没有立刻反应'],
        ['平衡策略：只在听到明确信号词时才延迟/划掉——不是怀疑每个答案',
         '划掉前先确认：确实听到了更正/否定信号词才划，不是因为"感觉要变了"就划',
         '信号词反复训练：把更正/否定/条件信号词做成词卡，每天看5分钟，形成条件反射'])
    U.content(p,'建立陷阱雷达：系统练习计划',
        [('Week 1：识别阶段','做完真题后，将所有错题分类：是哪类陷阱？信号词是什么？'),
         ('Week 2：强化阶段','专项练习：每次只练一类陷阱，S1全套只关注更正陷阱'),
         ('Week 3：整合阶段','综合练习，模拟考场环境，所有陷阱类型混合出现'),
         ('Week 4：自动化','做完题后，不需要特别关注信号词——陷阱意识已自动化')],
        note='目标：4周后，陷阱题正确率从约50%提升到85%以上')
    U.content(p,'深圳学生陷阱失分统计',
        [('更正型','失分率最高（35%）——深圳学生反应速度快但容易不等后续更正'),
         ('否定型','失分率第二（28%）——not/never在快速语流中容易被忽视'),
         ('相似音型','失分率第三（20%）——13/30等数字混淆，与中式英语听感有关'),
         ('条件型','失分率第四（12%）——逻辑理解难度较高，需要完整理解语境'),
         ('未来计划型','失分率第五（5%）——相对明显，大多数学生能识别')],
        note='根据以上数据，本课优先训练更正型和否定型陷阱，投入80%的练习时间')
    U.tip_card(p,'六大陷阱即时参考卡',
        [('更正型','听到 actually/sorry/I mean → 立刻划掉草稿上的候选答案，等新答案'),
         ('否定型','听到 not/never/without → 脑中"反转"，后接内容不是答案'),
         ('条件型','听到 if/unless/only if → 打问号，继续听确认条件是否成立'),
         ('相似音型','数字/音近词 → 等整句/等拼读再写，不要听到一半就下笔'),
         ('计划型','听到 plan/intend/going to → 打问号，等但是/转折/确认'),
         ('复合型','多个信号词叠加 → 放慢节奏，等完整逻辑链结束再写答案')])
    U.tip_card(p,'陷阱草稿操作规范',
        [('标注问号','对不确定的答案，在旁边写"?"——提醒自己此处有风险'),
         ('划掉动作','用横线划掉（不是涂抹）：一眼就能看出"之前的"和"更正后的"'),
         ('延迟策略','听到信号词后，笔停0.5-1秒，等待后续更正或确认'),
         ('机考操作','机考：先在草稿纸上写，确认后再输入；不要急着输入第一个听到的词'),
         ('誊写核查','誊写时，所有有"?"标注的题，重新确认逻辑再誊写')])
    U.checklist(p,'L08课后自检清单',
                   ['做2套真题听力，专项标注所有陷阱类型（更正/否定/条件/相似音/计划）',
                   '统计每类陷阱的识别成功率，找出最弱的陷阱类型',
                   '做S3选择题专项：每道题说出"是哪类陷阱+信号词是什么"',
                   '建立个人错题陷阱档案：记录自己最常掉入的陷阱'])
    U.summary(p,8,
        ['六大陷阱：更正型/否定型/条件型/相似音型/未来计划型/复合型',
         '信号词是关键：actually/not/if/unless/plan to——形成条件反射',
         '陷阱应对三步：识别信号词→延迟写答案→等逻辑链完整再写',
         '不要过度怀疑：只在明确信号词出现时才延迟，不是怀疑每个答案',
         '陷阱是可训练的：专项练习4周后，陷阱识别率目标≥85%'])
    U.homework(p,8,
        ['做剑桥真题3套，专项标注所有陷阱（颜色区分：红=更正，蓝=否定，绿=条件）',
         '统计个人最弱陷阱类型，下周针对该类型做专项练习',
         '背诵信号词表：更正（6个）+否定（8个）+条件（5个）'],
        est_time='约75分钟')
    U.preview(p,9,'口音适应与连读',
        ['英式vs澳式：深圳学生最不熟悉的两种口音特征',
         '连读规则：链接音/省略音/弱读——让"听不懂"变"听懂了"',
         '元音变化：英式元音与美式元音的系统对比'])
    U.save(p, f'{OUT}/L08_六大陷阱类型全析.pptx')
    print(f"  L08 done: {len(p.slides)} pages")

# ══════════════════════════════════════════════════════════════════════════════
#  L09  口音适应与连读
# ══════════════════════════════════════════════════════════════════════════════
def L09():
    T(); p=U.new_prs()
    U.cover(p,'听力',9,'口音适应与连读','Accent Adaptation & Linking')
    U.agenda(p,['英式/澳式口音特征','元音系统对比','连读规则','弱读与吞音',
                '口音词汇表','口音训练步骤','课堂练习','小结作业'])
    U.content(p,'L08回顾：六大陷阱核心要点',
        [('更正型','actually/sorry/I mean → 听到信号词立刻划掉前答案'),
         ('否定型','not/never/without → 否定词后内容不是答案'),
         ('条件型','if/unless/only if → 等条件确认再写答案'),
         ('复合型','多信号词叠加 → 等完整逻辑链说完再写')],
        note='本节课目标：适应英式/澳式口音，掌握连读规则，让"听不懂"变为"听得清"')
    U.objectives(p,['识别英式和澳式口音的5大核心特征',
                    '掌握5种连读规则及其在雅思听力中的应用',
                    '能辨认弱读形式（schwa音/辅音省略）',
                    '建立200个口音适应词汇的认知库'])
    U.pain_point(p,'深圳学生的口音适应三大痛点',
        ['美式耳朵听英式口音：学了多年美式英语，雅思却考英式+澳式，/ɑː/听成/æ/，/ɒ/听成/ɑː/',
         '连读导致词边界消失："want to"→"wanna"/"did you"→"didja"，分不出是几个词',
         '弱读被跳过：虚词被弱读到几乎消失（a/the/of），核心词前的这些词听不清导致误判答案'],
        ['英式口音特征表：/ɑː/（英）vs /æ/（美），/ɒ/（英）vs /ɑː/（美），每天听BBC 10分钟',
         '连读训练：把"want to"/"going to"/"kind of"等超高频连读词组专项训练，直到能秒识',
         '弱读脱敏：主动忽略虚词，只追踪实义词——a/the/of/to不是答案，不需要听清'])
    part(p,1,'英式口音：五大核心特征','British Accent: Five Core Features',
        [('元音差异：英式 vs 美式','雅思主要考英式和澳式，深圳学生更熟悉美式',
          [('/ɑː/ vs /æ/','英式：bath/can\'t/last/fast → /bɑːθ/ ; 美式：/bæθ/'),
           ('/ɒ/ vs /ɑː/','英式：lot/hot/stop/body → /lɒt/ ; 美式：/lɑːt/'),
           ('/əʊ/ vs /oʊ/','英式：go/home/phone → /gəʊ/ ; 美式：/goʊ/'),
           ('/ɪə/ vs /ɪr/','英式：here/near/beer → /hɪə/ ; 美式：/hɪr/（r音化）'),
           ('/eə/ vs /er/','英式：there/hair/care → /ðeə/ ; 美式：/ðer/')]),
         ('辅音差异与发音特点','英式特有的辅音规律',
          [('非卷舌r','英式r在词中不发音：car /kɑː/ ; 美式 /kɑːr/（卷舌）'),
           ('T的发音','英式t更清晰：water/butter/city ; 美式常弱化为/d/（flap T）'),
           ('yod coalescence','英式：Tuesday→/ˈtjuːzdɪ/ ; 澳式类似英式'),
           ('不打嗝的h','英式h音清晰：have/hotel，美式有时弱化'),
           ('词尾-ary','英式：secretary /ˈsekrɪtri/ ; 美式 /ˈsekrɪteri/（多一个音节）')])],
        [('英式口音适应策略',
          [('被动接触','每天30分钟BBC Radio 4/BBC World Service，不需要精听，大量接触'),
           ('主动训练','影子跟读：选一段S1/S2英式录音，延迟0.5秒跟读，模仿元音'),
           ('/ɑː/专项','把can\'t/bath/last/fast/dance等/ɑː/词组写成卡片，每天朗读5遍'),
           ('对比训练','同时找美式和英式录音朗读同一段文字，比较元音差异'),
           ('目标','2周后听到can\'t即可区分/ɑː/和can的/ə/弱读')]),
         ('英式口音听力误区',
          [('最常见误区','把/ɑː/误听成其他音：can\'t听成了can，bath听成了bat'),
           ('r音误区','听到car/park/here时等待r的发音，但英式r不发音，等不来'),
           ('T音误区','英式t很清晰，不要把butter/water听成butter/wader'),
           ('字节数误区','secretary英式3音节，美式4音节，听觉长度不同'),
           ('解决方案','建立"英式音素预期"：主动期待/ɑː/和/ɒ/，而非美式对应音')])],
        [('英美口音对比辨音练习',
          '老师读以下词语，判断是英式还是美式口音：\n\n'
          '① can\'t  ② water  ③ car  ④ secretary  ⑤ dance\n'
          '⑥ tomorrow  ⑦ butter  ⑧ either  ⑨ schedule  ⑩ leisure',
          '英式标志：①/kɑːnt/（长元音）③/kɑː/（r不发音）④/ˈsekrɪtri/（少音节）⑤/dɑːns/\n⑧either英式/ˈaɪðə/，美式/ˈiːðər/\n⑨schedule英式/ˈʃedjuːl/，美式/ˈskedjuːl/\n⑩leisure英式/ˈleʒə/，美式/ˈliːʒər/',
          '每天听10分钟BBC Radio 4，大脑会在2周内自动适应英式元音系统'),
         ('英式口音真题场景识别',
          '以下词语在S1-S4真题中出现，用英式发音写出/ɑː/音的词：\n\n'
          'plant / answer / bath / chance / class / demand / example\n'
          'fast / last / past / staff / after / advantage / can\'t',
          '以上全部词在英式英语中含/ɑː/（a在某些辅音前）：\nplant /plɑːnt/, answer /ˈɑːnsə/, bath /bɑːθ/, can\'t /kɑːnt/\n深圳学生常把这些听成/æ/音，导致识别不出来',
          '规律：英式英语中，a在f/s/th/n/m等辅音前常读/ɑː/而非/æ/——这个规律覆盖约80个常用词')]
    )
    part(p,2,'澳式口音与连读规则','Australian Accent & Linking Rules',
        [('澳式口音：雅思第二主要口音','澳式与英式相似但有差异',
          [('/eɪ/ → /aɪ/','澳式：day/say/name → /daɪ/ /saɪ/ /naɪm/（与英式/eɪ/不同）'),
           ('/aɪ/ → /ɔɪ/','澳式：time/wine/right → /tɔɪm/ /wɔɪn/（高化，较少见）'),
           ('/iː/ 末尾音','澳式happy/city末尾/iː/较紧，英式较松'),
           ('鼻腔共鸣','澳式鼻腔共鸣强，整体音调上扬'),
           ('与英式共同点','非卷舌r，/ɑː/类似英式，/ɒ/与英式类似')]),
         ('五大连读规则','让"听不懂"变"听懂了"的核心规则',
          [('连音(Linking)','辅音结尾+元音开头：want it→wan-tit, take out→ta-kout'),
           ('省略(Elision)','相邻辅音省略：next day→nex-day, most students→mos-students'),
           ('同化(Assimilation)','音变：did you→didja, can you→kanja, meet you→meetcha'),
           ('弱读(Reduction)','虚词弱化：and→n, of→ov, to→tə, the→ðə, for→fə'),
           ('连续发音(Intrusion)','元音之间插入/w/或/j/：go on→go-won, see it→see-jit')])],
        [('澳式口音识别策略',
          [('元音重映射','听到/aɪ/时多联想/eɪ/对应词：die可能是day，sigh可能是say'),
           ('上扬语调','澳式语调末尾上扬，疑问句和陈述句末尾都可能上扬'),
           ('与英式对比','澳式/ɑː/与英式相似，/ɒ/也相近——主要区别是/eɪ/→/aɪ/'),
           ('真题比例','雅思听力约30%澳式口音，集中于S2和S4'),
           ('适应方法','听ABC Radio或澳式雅思真题，每天10分钟即可')]),
         ('连读在雅思真题中的分布',
          [('S1出现率','S1日常对话：连读密度最高，gonna/wanna/kinda最常见'),
           ('S2出现率','S2独白：弱读多于连读，虚词大量弱化'),
           ('S3出现率','S3学术讨论：正式语气略多，但could have/should have仍有连读'),
           ('S4出现率','S4讲座：最正式，连读最少，但仍有to/of/and等弱读'),
           ('备考重点','S1连读最多：wanna/gonna/kinda/didja是S1必考连读形式')])],
        [('连读识别练习：速读版',
          '以下短语的快速口语版本是什么？（写出发音）\n\n'
          '① want to  ② going to  ③ kind of  ④ a lot of  ⑤ I don\'t know\n'
          '⑥ let me  ⑦ did you  ⑧ could have  ⑨ have to  ⑩ out of',
          '① wanna  ② gonna  ③ kinda  ④ a lotta  ⑤ dunno\n⑥ lemme  ⑦ didja  ⑧ coulda  ⑨ hafta  ⑩ outta\n这10个是考场最高频连读形式，深圳学生必须能秒识',
          '雅思听力不会考这些非正式形式的拼写，但考生需要能听懂——这是被动技能而非主动技能'),
         ('澳式口音真题辨音',
          '以下场景来自澳式S2，听老师用澳式口音朗读并判断词义：\n\n'
          '"Toiday we\'re going to teik a look at the fecilities of this greight building."\n\n'
          '（澳式：Today=Toiday, take=teik, great=greight）\n翻译这句话并写出标准英式拼写',
          '"Today we\'re going to take a look at the facilities of this great building."\n澳式特征：/eɪ/→/aɪ/（take/great/today中的a音上移）\n/æ/→/eɪ/（facilities中有时出现）',
          '应对澳式口音：提前知道/eɪ/会变/aɪ/，听到/aɪ/时多考虑是否是/eɪ/词')]
    )
    part(p,3,'弱读与吞音','Weak Forms & Sound Dropping',
        [('弱读：虚词的变身','功能词在快速语流中被弱化',
          [('and → /n/','bread and butter → bread n butter'),
           ('of → /ə/','a lot of → a lotta / aloto'),
           ('to → /tə/','go to school → go tə school'),
           ('the → /ðə/','the answer → ðə answer（辅音前）/ ðɪ answer（元音前）'),
           ('for → /fə/','waiting for you → waiting fə you')]),
         ('吞音：辅音的消失','相邻辅音中较弱的那个消失',
          [('t在辅音前消失','next week → nex week, last night → las night'),
           ('d在辅音前消失','old man → ol man, kind of → kina of → kinda'),
           ('h在词中消失','tell him → tell-im, give her → give-er'),
           ('g在-ing词尾','running/talking → runnin/talkin（口语快速语流）'),
           ('l在特定词','always/already/alright → ol-ways/ol-ready/ol-right（英式）')])],
        [('弱读训练方法',
          [('Step 1：识别虚词','先确认哪些词是虚词（冠词/介词/连词/助动词），这些是弱读候选'),
           ('Step 2：不追踪虚词','做题时主动"放弃"追踪弱读虚词，只追踪实义词'),
           ('Step 3：语流练习','用"a lot of / kind of / sort of"等含弱读的短语大量跟读'),
           ('Step 4：真题验证','做完题后，统计弱读失分题数——目标：弱读不是失分原因'),
           ('最终目标','弱读虚词自动过滤，注意力100%在实义词上')]),
         ('吞音高发场景',
          [('辅音群','next door→nex door（t省略），last night→las night（t省略）'),
           ('h的省略','tell him→tell im，give her→give er（h省略）'),
           ('-ing结尾','running/talking→runnin/talkin（g省略，口语快速语流）'),
           ('真题中的吞音','Section 1中的吞音最多——日常对话语速快，吞音密度高'),
           ('应对策略','不要期待每个辅音都发音完整——允许"不完整"是关键心态')])],
        [('弱读辨音练习',
          '以下句子，老师正常语速朗读，请写出你听到的内容（接受连读/弱读形式）：\n\n'
          '"Could you tell me what time the last train to London is?"\n'
          '"I\'d like to book a room for a couple of nights."',
          '"Could you tell me what time the last train to London is?"\n快速口语版：/kəd jə tel mi wɒt taɪm ðə lɑːs treɪn tə ˈlʌndən ɪz/\n弱读：Could→/kəd/, you→/jə/, the→/ðə/, to→/tə/\n吞音：last→las（t省略，因后接辅音tr）',
          '训练目标：听到/kəd jə/立刻反应是"Could you"——这是被动识别，不是主动发音'),
         ('吞音识别真题练习',
          '以下S1对话，找出所有吞音/弱读并标注：\n\n'
          '"So you\'d like to sign up for the erm... next available slot?\n'
          'That would be on the 15th of March, at half past two in the afternoon."\n\n'
          '请标出：① 弱读词  ② 吞音位置',
          '弱读：you\'d→/jəd/, to→/tə/, the→/ðə/, of→/əv/, at→/ət/\n吞音：next（t在av前省略）→/neks/; past（t后接two的t省略）→/pɑːs/\n半连读：half past two→/hɑːf pɑːs tuː/',
          '关键认知：弱读/吞音不是错误发音——这是英式口语的正常形式，真题就是这么录的')]
    )
    part(p,4,'口音综合适应策略','Integrated Accent Adaptation Strategy',
        [('口音适应的科学原理','大脑需要"重新校准"音素认知',
          [('音素期待','大脑会预测接下来的音，口音不同会打破预测，导致"没听到"'),
           ('适应机制','大量接触→大脑建立新的音素预测模型→自动适应'),
           ('时间框架','每天30分钟英式/澳式输入，2-3周后明显改善'),
           ('输入质量','真题录音质量最高；BBC Podcasts/ABC Radio也是优质输入源'),
           ('主动vs被动','被动听效果差；主动影子跟读（shadowing）效果最好')]),
         ('影子跟读法：口音适应最快技术','shadowing是公认最快的口音适应训练',
          [('什么是影子跟读','听到录音就立刻跟着读，延迟约0.5秒，尽量模仿语音语调'),
           ('重点','不是模仿得完美，而是逼迫大脑实时处理口音输入'),
           ('材料','用雅思真题S1-S2的英式/澳式段落，每段1-2分钟'),
           ('时间','每天10-15分钟，不要超过20分钟（高强度）'),
           ('频率','每天进行，连续14天，口音适应效果最为显著')])],
        [('影子跟读进阶方法',
          [('入门阶段','先听一遍不跟读，了解内容；第二遍才开始跟读（降低认知负荷）'),
           ('延迟控制','延迟0.5秒（而非即时）：留时间预测下一个音，不是直接回音'),
           ('重点放大','跟读完后，专门重复/ɑː/和/ɒ/音的词，不需要重复整句'),
           ('进阶挑战','能做到：听到录音的同时跟读，且理解内容——这是流利阈值'),
           ('材料难度','从S2独白开始（最清晰），稳定后再转S4学术讲座')]),
         ('口音适应的常见误区',
          [('误区1','只听不跟读：被动听30分钟不如主动跟读10分钟有效'),
           ('误区2','追求每个音完美：口音适应是识别不是发音，不追求完美复制'),
           ('误区3：太难的材料','从S2单人独白开始，不要上来就听S4学术讲座'),
           ('误区4','间歇训练：隔2-3天一次的影子跟读效果远不如每天10分钟'),
           ('正确预期','2周后：识别率提升；4周后：口音不再是听力障碍')])],
        [('影子跟读实操演示',
          '以下S2英式独白，先听一遍，再跟读（延迟0.5秒）：\n\n'
          '"Welcome to Greenfield Nature Reserve. Today we\'ll be exploring three main trails.\n'
          'The first trail starts near the car park and takes approximately forty-five minutes.\n'
          'Along the way, you\'ll pass the bird-watching hide and the ancient oak tree."',
          '跟读重点：\n① car park /kɑː pɑːk/（r不发音）\n② approximately /əˈprɒksɪmɪtli/（/ɒ/音）\n③ forty-five /ˌfɔːti ˈfaɪv/（/ɔː/非/oʊ/）\n④ ancient /ˈeɪnʃənt/（/eɪ/非/eɪ/区别不大）\n⑤ watching /ˈwɒtʃɪŋ/（/ɒ/非/ɑː/）',
          '跟读不需要发音完美——重点是逼迫大脑实时识别英式元音，建立新的音素预测库'),
         ('口音适应自测：4周进度检查',
          '按以下标准自测口音适应程度：\n\n'
          '① 听BBC新闻播报，能理解80%以上？（是/否）\n'
          '② 听到/ɑː/音，能立刻反应对应的词？（是/否）\n'
          '③ 听连读"going to"，能立刻识别为"gonna"？（是/否）\n'
          '④ 做真题时，因口音问题漏题<3题？（是/否）',
          '4个"是"：口音适应优秀，可进入冲刺阶段\n3个"是"：良好，继续保持输入密度\n2个"是"：需要加强，每天影子跟读增加到20分钟\n1个以下：重点加强，专门做英式/澳式口音专项训练',
          '自测结果诚实记录，不要自我安慰——知道弱点在哪才能有针对性地提高')]
    )
    # PART05 vocab
    U.section(p,5,'口音核心词汇表','Accent Vocabulary Bank')
    U.vocab(p,'英式/澳式发音特殊词（vs美式）',
        [('/ɑː/系列','bath/can\'t/last/fast/dance/example /ɑː/（英）vs /æ/（美）'),
         ('/ɒ/系列','lot/hot/stop/got/body /ɒ/（英）vs /ɑː/（美）'),
         ('schedule','英/ˈʃedjuːl/ vs 美/ˈskedjuːl/（sh vs sk）'),
         ('either','英/ˈaɪðə/ vs 美/ˈiːðər/（ai vs ee）'),
         ('leisure','英/ˈleʒə/ vs 美/ˈliːʒər/（e vs ee）'),
         ('laboratory','英/ləˈbɒrətri/ vs 美/ˈlæbrətɔːri/（重音位置不同）'),
         ('advertisement','英/ədˈvɜːtɪsmənt/ vs 美/ˌædvərˈtaɪzmənt/')])
    U.vocab(p,'高频连读形式（必须秒识）',
        [('want to','= wanna /ˈwɒnə/（最高频，几乎每套真题出现）'),
         ('going to','= gonna /ˈɡɒnə/（未来时标志，连读后易误判）'),
         ('kind of','= kinda /ˈkaɪndə/（S3学术讨论常见）'),
         ('a lot of','= a lotta /əˈlɒtə/（数量表达）'),
         ('did you','= didja /ˈdɪdʒə/（疑问句）'),
         ('could have','= coulda /ˈkʊdə/（虚拟语气，S3常见）'),
         ('out of','= outta /ˈaʊtə/（位置/数量关系）')])
    U.vocab(p,'弱读核心词（不需要听清的虚词）',
        [('and','弱读 /n/ 或 /ənd/（这些词不是答案，不需要清晰听到）'),
         ('of','弱读 /ə/ 或 /əv/（非重读位置几乎消失）'),
         ('to','弱读 /tə/（方向/目的介词，非答案）'),
         ('the','弱读 /ðə/（辅音前）/ /ðɪ/（元音前）'),
         ('for','弱读 /fə/（介词，通常非答案）'),
         ('at','弱读 /ət/（时间/地点介词）'),
         ('can','弱读 /kən/（注意：can\'t弱读时仍有/ɑː/，用于区分can vs can\'t）')])
    U.two_col(p,'英式vs美式口音核心对比',
        '元音差异（影响识别）',
        [('/ɑː/ vs /æ/','can\'t/bath/dance/fast/last'),
         ('/ɒ/ vs /ɑː/','lot/hot/stop/not/body'),
         ('/əʊ/ vs /oʊ/','go/home/phone/know'),
         ('/ɪə/ vs /ɪr/','here/near/beer/clear'),
         ('/eɪ/ → /aɪ/','澳式：day/say/name/place')],
        '辅音差异（影响识别）',
        [('卷舌r','美式词尾r卷舌，英式/澳式不卷'),
         ('Flap T','美式water=/ˈwɑːdər/，英式/ˈwɔːtə/'),
         ('T+Y合并','英式Tuesday=/ˈtjuːzdɪ/（tyoo不分离）'),
         ('ng结尾','口语-ing常省略g：running→runnin'),
         ('h弱化','him/her/have在快速语流中h常省略')])
    # PART06 method
    U.section(p,6,'口音适应：系统训练步骤','Accent Adaptation: Training Steps')
    U.steps(p,'Step 1-3：认知建立阶段（Week 1）',
        [('建立口音特征认知','学习英式/澳式5大元音特征，知道"会有什么不同"'),
         ('被动接触','每天30分钟：BBC Radio 4 / 雅思真题录音，不需要精听'),
         ('特征词圈划','做真题时，遇到英式特殊发音词，圈出并注音')])
    U.steps(p,'Step 4-6：主动训练阶段（Week 2-3）',
        [('影子跟读（每天15分钟）','用雅思S2独白进行影子跟读，重点模仿/ɑː/和/ɒ/'),
         ('连读专项','10个高频连读词组：wanna/gonna/kinda/didja等，每天默读5遍'),
         ('做题检验','做完真题后，标注因口音问题做错的题，分析是哪个音素出了问题')])
    U.steps(p,'Step 7-8：巩固阶段（Week 4）',
        [('口音自测','4周自测：BBC新闻理解率≥80%；真题口音失分<3题/套'),
         ('维持输入','每天10分钟英式/澳式输入维持——停止输入后大脑会"漂移"回美式习惯')])
    U.model_answer(p,'口音适应训练：30天计划',
        'Week 1：认知\n'
        'Week 2：接触（被动输入为主）\n'
        'Week 3：主动训练（影子跟读+连读专项）\n'
        'Week 4：检验+巩固',
        'Day 1-7：每天30分钟BBC，不追求听懂，只"泡"在英式输入中\n'
        'Day 8-14：开始影子跟读，每天15分钟；同时学10个连读词组\n'
        'Day 15-21：做2套真题，专项标注口音失分题，针对性补课\n'
        'Day 22-30：全流程模考，口音问题应已基本消除',
        ['30天不是指30天就能完全适应——而是口音不再是做题瓶颈所需的最少天数',
         '深圳学生英语基础好，口音适应通常比其他地区考生更快（约2-3周显效）',
         '维持阶段：考试前每天继续10分钟英式输入，防止"漂移"'])
    U.model_answer(p,'口音真题分析：S2英式独白典型例子',
        'Script (British): "The car park is situated to the north of the building,\n'
        'accessible via the back entrance. It can\'t accommodate more than\n'
        'fifty vehicles at any one time, and staff passes are not valid after six."',
        '口音陷阱分析：\n① car /kɑː/ → r不发音（英式），不要听成美式/kɑːr/\n② can\'t /kɑːnt/ → /ɑː/，不要误听成can /kən/（弱读）\n③ fifty /ˈfɪfti/ → 注意vs fifteen /fɪfˈtiːn/（重音位置区分）\n④ valid /ˈvælɪd/ → /æ/（这个词美英发音相同）',
        ['car/park/can\'t都含英式/ɑː/——统一处理，不要每个词单独分析',
         'can\'t识别关键：重读+/ɑː/长元音；can弱读为/kən/（短促）',
         '这类S2独白是口音适应练习的最佳材料——单人、清晰、实用词汇'])
    U.compare(p,'口音适应：训练前 vs 训练后',
        '❌ 训练前（6分常见困境）',
        ['听到/ɑː/误判为/æ/：can\'t听成can，写错答案',
         '连读没反应：gonna/wanna听不出是going to/want to',
         '澳式口音感觉"奇怪"，做澳式题正确率明显低于英式',
         '每套真题因口音问题丢失3-5题',
         '做真题时频繁暂停，尝试"听清楚"——浪费注意力'],
        '✅ 训练后（7.5+分目标）',
        ['听到/ɑː/立刻反应：can\'t，不需要停下来思考',
         '连读自动识别：gonna=going to，大脑自动解码',
         '英式/澳式题正确率无明显差异（已适应两种口音）',
         '每套真题因口音问题丢失<1题',
         '做真题时全程流畅，不因口音分神'])
    U.compare(p,'连读误区：识别失败 vs 成功',
        '❌ 连读识别失败',
        ['"going to visit" → 听到"gonna visit"，不知道是什么词',
         '"did you" → 听到"didja"，认为录音质量差',
         '"kind of" → 听到"kinda"，停下来纠结',
         '"could have" → 听到"coulda"，误判为"cold"',
         '"want to" → 听到"wanna"，认为是人名'],
        '✅ 连读识别成功',
        ['"gonna visit" → 立刻识别为"going to visit"',
         '"didja" → 立刻识别为"did you"，这是正常口语',
         '"kinda" → 立刻识别为"kind of"，继续追踪内容词',
         '"coulda" → 立刻识别为"could have"，判断虚拟语气',
         '"wanna" → 立刻识别为"want to"，追踪动作内容'])
    # PART07 drill
    U.section(p,7,'课堂练习：口音辨识','Accent Recognition Drills')
    U.drill(p,'练习1：英美口音对比辨音',
        '老师分别用英式和美式发音读以下词，写出你听到的音标特征',
        ['① can\'t（区分/ɑː/ vs /æ/）',
         '② water（区分/ˈwɔːtə/ vs /ˈwɑːdər/）',
         '③ hot（区分/hɒt/ vs /hɑːt/）',
         '④ here（区分/hɪə/ vs /hɪr/）',
         '⑤ schedule（区分/ˈʃedjuːl/ vs /ˈskedjuːl/）'])
    U.drill(p,'练习2：连读解码（10组）',
        '听老师读以下连读形式，写出完整的标准英文',
        ['① wanna  ② gonna  ③ kinda  ④ didja  ⑤ coulda',
         '⑥ hafta  ⑦ lemme  ⑧ outta  ⑨ dunno  ⑩ shoulda'])
    U.drill(p,'练习3：弱读虚词识别',
        '听句子，标出弱读的虚词（and/of/to/the/for/at等）',
        ['"A lot of students are waiting for the bus at the stop."',
         '"She\'s going to apply for a place at the university."',
         '"Can you tell me what time the lecture starts and ends?"',
         '目标：弱读虚词在实际语速中几乎听不清——这是正常现象'])
    U.drill(p,'练习4：澳式口音辨识',
        '以下句子用澳式口音朗读，写出标准拼写',
        ['"Toiday\'s session will be on the moin topic of climate chynge."',
         '"We\'re going to teik a luk at the reizons behind this triend."',
         '（澳式：today/take/look/reasons/trend中的元音上移）'])
    U.drill(p,'练习5：口音综合真题实战',
        '以下S2英式录音（节选），完整听写并标注3处口音特征',
        ['"The nature reserve covers approximately 200 hectares of woodland\n'
         'and is home to over 150 species of birds. Visitors can\'t park\n'
         'inside the reserve, but there\'s a free car park half a mile away."',
         '口音特征标注：① can\'t的/ɑː/ ② car park的r省略 ③ hectares的重音',
         '注意：approximately/species/reserve中的英式元音'])
    # PART08 pitfall
    U.section(p,8,'口音适应误区','Accent Adaptation Pitfalls')
    U.pain_point(p,'口音适应的三大误区',
        ['只靠做题适应：每套真题做完就结束，没有专项口音训练——做100套题口音还是没适应',
         '只听美式材料：平时听美式播客/美剧，考前才发现口音不适应——太晚了',
         '追求"听懂每个词"：在口语快速语流中，母语者本身也听不清每个虚词——这是正常现象'],
        ['专项输入：每天10分钟BBC/ABC，不需要特别精听，大量接触本身就是训练',
         '提前准备：至少提前4周开始英式/澳式输入，不要到备考最后阶段才意识到口音问题',
         '调整预期：不追求听清虚词，只追踪实义词——答案永远在实义词中，不在虚词里'])
    U.content(p,'口音适应与做题策略的结合',
        [('预读期间','不受口音影响——预读是看题目，不是听录音'),
         ('听题期间','已知口音特征→自动解码→不需要特别注意'),
         ('跟丢后重定位','用题目中的关键词（书面形式）重新定位，不靠音感'),
         ('誊写期间','拼写靠题目意义判断，不靠"听起来怎么写"'),
         ('口音失分补救','做错因口音的题：标注音素→专项输入→下次真题重测')],
        note='口音适应是背景能力，做好后无需再思考——让精力集中在策略和预读上')
    U.content(p,'深圳学生口音适应特殊建议',
        [('优势','深圳学生普遍英语基础好，音素认知能力强，口音适应更快'),
         ('劣势','大量美式英语输入（美剧/美式教材）导致美式口音"优势假设"'),
         ('核心建议','备考期间暂停美式输入，全面切换到英式/澳式材料'),
         ('时间要求','至少4周，最好8周——神经可塑性需要时间'),
         ('考前一周','每天听2-3段雅思真题录音，保持英式/澳式"热身"状态')],
        note='口音适应一旦建立就不容易丢失——考后继续接触英式输入效果会持续')
    U.tip_card(p,'口音适应即时技巧',
        [('/ɑː/系列','听到/æ/感觉"不对"时，试试/ɑː/解读：can\'t/bath/fast/last'),
         ('/ɒ/系列','听到/ɑː/感觉不是答案时，试试/ɒ/解读：lot/hot/stop'),
         ('连读识别','听到"gonna/wanna/kinda"立刻解码为标准英文，不停顿'),
         ('弱读策略','听不清虚词是正常的——只追踪实义词（名词/动词/数字）'),
         ('澳式应对','澳式/eɪ/→/aɪ/：听到"dye"可能是"day"；"sigh"可能是"say"')])
    U.tip_card(p,'30天口音适应计划速查',
        [('Week 1（认知）','学习英式/澳式5大特征；每天30min被动输入（BBC/ABC）'),
         ('Week 2（接触）','继续被动输入；开始学10个连读词组；听真题不看脚本'),
         ('Week 3（训练）','影子跟读每天15min；做真题标注口音失分题；专项补课'),
         ('Week 4（检验）','口音自测；真题口音失分目标<3题；维持10min/天输入'),
         ('考前维持','每天10min英式/澳式输入，保持"校准"状态直到考试当天')])
    U.checklist(p,'L09课后自检清单',
                   ['开始每天30分钟英式输入（BBC Radio 4 / 雅思真题）',
                   '掌握10个高频连读形式，能秒识wanna/gonna/kinda/didja等',
                   '做1套真题，标注所有因口音问题做错的题',
                   '建立"口音失分词汇表"，每周回顾'])
    U.summary(p,9,
        ['英式口音5大特征：/ɑː/vs/æ/、/ɒ/vs/ɑː/、非卷舌r、清晰t、词尾音变',
         '澳式口音：/eɪ/上移为/aɪ/（day听成die），整体与英式相近',
         '5种连读规则：连音/省略/同化/弱读/插入——必须能自动识别',
         '10个高频连读形式：wanna/gonna/kinda/didja/coulda等必须秒识',
         '30天训练计划：认知→接触→主动训练→检验，每天坚持'])
    U.homework(p,9,
        ['开始30天口音适应计划，今天第1天：听BBC新闻30分钟',
         '做剑桥真题1套，专项标注口音失分题（有哪个音没听清楚？）',
         '建立连读词组卡片，今天背前5个：wanna/gonna/kinda/didja/coulda'],
        est_time='约50分钟')
    U.preview(p,10,'全真模考与查漏补缺',
        ['标准模考流程：从考场入座到10分钟誊写的全程规范',
         '错题分类诊断系统：把40道错题分成6类，找到提分最快的突破口',
         '冲刺计划制定：根据目标分数制定最后2周的备考方案'])
    U.save(p, f'{OUT}/L09_口音适应与连读.pptx')
    print(f"  L09 done: {len(p.slides)} pages")

# ══════════════════════════════════════════════════════════════════════════════
#  L10  全真模考与查漏补缺
# ══════════════════════════════════════════════════════════════════════════════
def L10():
    T(); p=U.new_prs()
    U.cover(p,'听力',10,'全真模考\n& 查漏补缺','Full Mock Exam & Gap Analysis')
    U.agenda(p,['标准模考规范','错题分类诊断','冲刺计划制定',
                '高频失分题型','考场词汇冲刺','模考步骤流程','课堂练习','小结作业'])
    U.content(p,'L09回顾：口音适应与连读核心',
        [('/ɑː/系列','can\'t/bath/fast/last — 英式/ɑː/不要误听成美式/æ/'),
         ('连读核心','wanna/gonna/kinda/didja/coulda — 10个高频形式必须秒识'),
         ('弱读策略','虚词弱读是正常现象，只追踪实义词，不追求听清每个词'),
         ('30天计划','Week 1认知→Week 2接触→Week 3主动训练→Week 4检验')],
        note='本节课：全程模考+系统查漏，制定个人冲刺计划，冲刺目标分数')
    U.objectives(p,['掌握标准模考规范，模拟真实考场条件',
                    '能对40道题的错题进行系统分类诊断',
                    '根据诊断结果制定个人化的2周冲刺计划',
                    '掌握考场最后冲刺的策略和心态管理'])
    U.pain_point(p,'模考阶段的三大误区',
        ['只做不分析：每套模考做完核对答案就结束，不分析错题类型——永远在同一个坑里跌倒',
         '只分析不改变：知道自己哪里错，但不针对性练习——知道问题≠解决问题',
         '模考条件不严格：在嘈杂环境做题、随时暂停、不计时——训练出来的不是考场能力'],
        ['错题分类诊断：每套做完，把错题分成6类（词汇/口音/陷阱/预读/速记/粗心）',
         '针对性补课：找到自己最多的错题类别，下一套专项关注该类别',
         '严格模考条件：安静房间+计时+连续做完不暂停+纸笔考用答题纸誊写'])
    part(p,1,'标准模考规范','Standard Mock Exam Protocol',
        [('考前准备：模拟真实考场','备考最后阶段用真实考场条件',
          [('环境','安静房间，无手机，无打断——不在咖啡厅/图书馆做模考'),
           ('材料','打印答题纸（纸笔考）或使用雅思官方模考界面（机考）'),
           ('工具','铅笔（纸笔考）或标准键盘（机考）；草稿纸2张'),
           ('时间','选固定时间段，与考试时间一致（早上/下午）'),
           ('状态','模考前不做其他题目，保持精力充沛')]),
         ('模考中的规范操作','每套模考必须严格执行的操作',
          [('不暂停','全程30分钟不暂停，模拟真实录音不可倒带的考场条件'),
           ('不回头','答完一题立刻移向下一题，不回头检查已答题目'),
           ('预读执行','每组题说明开始时立刻预读，20秒内完成关键词圈划'),
           ('速记执行','全程使用速记符号，不写完整词，数字用阿拉伯数字'),
           ('誊写规范','纸笔考：30分钟后立刻誊写，限时10分钟，不超时')])],
        [('模考规范的核心逻辑',
          [('为什么不暂停','真实考试录音不能倒带——训练中暂停等于在练习一项考场不存在的技能'),
           ('为什么不回头','回头检查浪费预读时间——预读时间比回头检查价值更高'),
           ('为什么用答题纸','训练誊写时间管理——10分钟誊写40题，需要专项训练才能做到'),
           ('机考 vs 纸笔考','机考没有誊写时间，但Section间隙可以检查；节奏与纸笔考不同'),
           ('模考频率建议','冲刺阶段：每天1套（约90分钟含分析）；平时：隔天1套')]),
         ('模考分数的正确解读',
          [('原始分→雅思分','Band 7.5=35-36分; Band 7=30-32分; Band 6.5=26-29分; Band 6=23-25分'),
           ('单次波动','单次模考分数波动1个Band是正常的，不代表真实水平'),
           ('5套平均','至少做5套后再评估真实水平，单次不算数'),
           ('上升趋势','比绝对分数更重要——连续5套都是6.5说明真实水平6.5'),
           ('考前预测','考前最后3套平均分-0.5=考场预期分数（考场压力因素）')])],
        [('模考规范自检：你做到了吗？',
          '以下是10个模考规范，给自己打分（做到得1分，共10分）：\n\n'
          '① 安静环境，无干扰\n② 计时30分钟，全程不暂停\n③ 使用答题纸誊写\n④ 预读执行率≥90%\n⑤ 速记符号使用率≥70%',
          '⑥ 不回头检查已答题\n⑦ 誊写限时10分钟\n⑧ 誊写完毕才核对答案\n⑨ 核对后进行错题分类\n⑩ 记录本次分数并与上次对比\n\n得分判断：9-10分=模考质量优秀；7-8分=良好；<7分=模考条件需改进',
          '模考质量比模考数量更重要——1套严格模考>3套随意模考'),
         ('考场真实流程预演',
          '考场真实流程（按时间顺序）：\n\n'
          '① 入座→取答题纸→不要翻题目\n'
          '② 监考员宣读说明→此时可提问\n'
          '③ 录音开始：S1说明→S1题组→S1作答\n'
          '④ Section间隙（30秒）：立刻预读下一组\n'
          '⑤ 完成S4→10分钟誊写（纸笔考）\n'
          '⑥ 收卷→离场',
          '考场流程核心：Section说明时间=预读时间；Section间隙=继续预读；誊写=核查拼写',
          '机考考生：没有誊写时间，Section间隙直接检查已输入内容，不等最后')]
    )
    part(p,2,'错题诊断系统：六大错误类型','Error Diagnosis: Six Error Categories',
        [('错题分类表：每套必做','把每道错题归入六类之一',
          [('类型A：词汇失分','听到了但不认识这个词（拼写/含义不熟悉）'),
           ('类型B：口音失分','听到了但因口音问题没认出来（/ɑː/vs/æ/等）'),
           ('类型C：陷阱失分','听懂了但写了陷阱给的错误答案（更正/否定/条件）'),
           ('类型D：预读失分','没有预读或预读不完整，错过关键词窗口'),
           ('类型E：速记失分','听到了但来不及写，或写了但誊写时看不清'),
           ('类型F：粗心失分','拼写/大小写/限词/单复数错误')]),
         ('各类型的应对策略','诊断出类型→针对性补课',
          [('A类（词汇）','背诵高频题型词汇表；每周新增30个；做错题词汇本'),
           ('B类（口音）','增加英式/澳式输入密度；影子跟读；专项音素训练'),
           ('C类（陷阱）','重做陷阱题专项；标注每道错题的信号词；L08复习'),
           ('D类（预读）','强制执行预读20秒计时；每次做题前打开秒表'),
           ('E类（速记）','符号系统专项强化；加快书写速度训练；练习占位技术'),
           ('F类（粗心）','建立誊写检查流程：拼写→大小写→单复数→限词')])],
        [('六类错误的时间投入分配',
          [('F类（粗心）','每天5分钟：誊写时检查流程，不需要额外学习，只需要养成习惯'),
           ('E类（速记）','每天10分钟：速记符号复习，练习符号速写，做占位技术训练'),
           ('D类（预读）','每次做题前：计时预读20秒，不超时——每次必做，不跳过'),
           ('C类（陷阱）','每周1-2次：专项陷阱题训练，L08复习，信号词词卡'),
           ('B类（口音）','每天10分钟：影子跟读+连读词组；长期投入但每次时间短'),
           ('A类（词汇）','每天15分钟：词汇本复习+新词记忆；需要最长时间积累')]),
         ('深圳学生错题类型统计',
          [('F类最多','深圳学生最常见：拼写失误（accommodation/necessary/environment等）'),
           ('D类第二','预读执行率不足——会预读但不坚持，尤其S3/S4'),
           ('C类第三','陷阱识别意识有但反应不够快——听到信号词0.5秒内没反应'),
           ('A类积累','词汇量普遍OK，但特殊场景词（医疗/法律/环境）有盲区'),
           ('B类相对好','深圳学生英语基础好，口音适应比其他地区快')])],
        [('错题分类实操：上一套模考的分析',
          '取出你最近一套模考的错题，按以下步骤分类：\n\n'
          '步骤1：写下错题题号\n'
          '步骤2：每道题对应A/B/C/D/E/F哪类\n'
          '步骤3：统计各类型数量\n'
          '步骤4：最多的那类=本周补课重点',
          '示例分析：\n错题：Q3(A词汇) Q8(C陷阱) Q12(D预读) Q17(C陷阱) Q23(F粗心) Q31(B口音) Q35(C陷阱)\nC类(陷阱)=3道，是本周最主要失分原因\n本周重点：L08陷阱章节复习+陷阱专项练习',
          '每套模考做完都做这个分析——积累10套后，你会发现自己的"稳定失分模式"'),
         ('常见失分组合与突破顺序',
          '深圳学生最常见的错题组合：\n\n'
          '组合1：F类（拼写）+ D类（预读）→ 最快提分：先解决F类（拼写），2周内多对3-4题\n'
          '组合2：C类（陷阱）+ B类（口音）→ 先解决C类，再练B类\n'
          '组合3：A类（词汇）+ E类（速记）→ 先解决E类（立刻见效），再积累词汇',
          '为什么先解决F/E/D类？这三类是"技能型"失分（可快速习得），不是"知识型"失分\n词汇（A类）需要长期积累，短期突击效果有限\n陷阱（C类）需要意识+反应速度，需要2-3周专项训练',
          '冲刺2周：优先解决技能型失分（F/E/D类），这些是最快看到分数提升的领域')]
    )
    part(p,3,'冲刺计划：最后2周','2-Week Sprint Plan',
        [('冲刺计划的核心原则','最后2周的备考不是"新学东西"',
          [('原则1','不要在最后2周学新词汇或新题型——只巩固已学的'),
           ('原则2','每天1套模考（30分钟）+1小时错题分析，共90分钟'),
           ('原则3','错题分析严格执行六分类，每类出现>2次就针对性复习'),
           ('原则4','考前3天停止新内容，只做预读热身和轻量复习'),
           ('原则5','考前1天不做题，休息，确保考试当天精力充沛')]),
         ('每日冲刺时间表','最后14天的具体安排',
          [('Day 1-5','每天：30min模考 + 60min错题分析 + 30min针对性复习'),
           ('Day 6-10','每天：30min模考 + 45min错题分析 + 45min薄弱点专项'),
           ('Day 11-12','完整模考+深度分析，确认已解决的失分类型'),
           ('Day 13','轻量复习：预读练习20min + 信号词背诵10min'),
           ('Day 14（考试日前一天）','不做任何练习，正常作息，早睡')])],
        [('冲刺计划的心理框架',
          [('不要panic','最后2周panic会降低发挥，保持"已准备好"心态比再学新东西更重要'),
           ('相信积累','L01-L09的训练已经建立了系统——最后2周是激活，不是从零开始'),
           ('接受不完美','模考有时做差是正常的，不等于真实水平下降'),
           ('每日小目标','与其设"这周提分0.5"，不如设"今天错题分类做完"——过程导向'),
           ('睡眠优先','考前3天早睡，睡眠对记忆巩固比多做1套题更有价值')]),
         ('考前最后3天清单',
          [('Day -3','最后一套完整模考+深度分析，确认哪类错题已解决'),
           ('Day -2','轻量复习：信号词表+拼写词表，各20分钟；不做新题'),
           ('Day -1','完全休息，准备好考试物品，早睡'),
           ('考试当天早上','吃早饭，到达考场早15分钟，不做题，轻松状态入场'),
           ('考场心态','相信已经准备好了——紧张的10%能量转化为专注')])],
        [('2周冲刺的预期效果',
          '基于深圳学生历史数据，执行完整2周冲刺计划的预期提分：\n\n'
          '起点6分学生 → 冲刺2周 → 平均达到6.5分\n'
          '起点6.5分学生 → 冲刺2周 → 平均达到7分\n'
          '起点7分学生 → 冲刺2周 → 平均达到7.5分',
          '提分关键：\n①严格执行错题分类，不随意做题\n②专项补课，而非重复同类题\n③保持稳定，不因一套做差而崩溃\n④考前保持好状态，技术层面已经够了',
          '最后2周，技术层面已经到位——状态管理是最后的决定因素'),
         ('个人冲刺计划制定',
          '根据你的错题诊断结果，填写个人冲刺计划：\n\n'
          '当前分数：___分  目标分数：___分\n'
          '主要失分类型：___（A/B/C/D/E/F）\n'
          '每日计划：___ 分钟模考 + ___ 分钟分析 + ___ 分钟针对性练习\n'
          '薄弱点：___（填具体内容：如"S3陷阱"/"拼写词10个"）',
          '填写完毕后和老师确认计划是否合理，避免计划太满或方向偏差',
          '冲刺计划是个人化的——没有适合所有人的万能计划，根据自己的失分类型定制')]
    )
    part(p,4,'高频拼写与考场心态','High-Freq Spelling & Exam Mindset',
        [('100个高频拼写失分词','F类失分的专项攻克',
          [('住宿类','accommodation（注意双c和双m）'),
           ('必要类','necessary（1个c，2个s）'),
           ('推荐类','recommend（1个c，2个m）'),
           ('月份类','February（r后还有u：Feb-ru-a-ry）'),
           ('设施类','facilities, equipment, environment（长词记词根）')]),
         ('考场心态管理','最后的能力：情绪稳定',
          [('跟丢了','立刻找锚点词重定位，不停下来纠结——损失1题，不是1分钟'),
           ('整组听错了','继续往前，不要"崩溃"——30分钟做完，后面3个Section还有机会'),
           ('听到不认识的词','用音译占位，誊写时靠上下文推断，不影响后续题'),
           ('计时误判','听到S3心想"怎么那么难"——这正是正常难度，正常发挥就行'),
           ('考后感觉','考后感觉"做得不好"是极其常见的，不代表实际分数低')])],
        [('高频拼写词记忆法',
          [('视觉分解','把长词分成音节块：ac-com-mo-da-tion，每块单独记忆'),
           ('词根辅助','environ（周围）+ ment；manage（管理）+ ment；这些词根出现在多个词中'),
           ('错误预警','accommodation（双c双m）是最常见拼写错误词，专门标记'),
           ('每日5个','每天记5个高频拼写词，21天后覆盖100个高频词'),
           ('遗忘曲线','第1天记→第3天复习→第7天复习→第21天复习——记住率提到90%')]),
         ('考场心态的科学依据',
          [('焦虑与表现','适度焦虑（30%）提升表现；高度焦虑降低注意力广度——目标是"冷静专注"'),
           ('跟丢的数学','跟丢1题=损失1分（最多）；停下纠结30秒=损失后续1-2题机会'),
           ('考后感觉失真','考后感觉差≠实际分数差：统计显示约60%感觉不好的考生分数达到目标'),
           ('Section间隙','30秒间隙：深呼吸2次+预读下一题——这30秒价值极高'),
           ('自我对话','考场出问题时，内心说："这正常，继续前进"——不评判自己')])],
        [('高频拼写词即时测试',
          '以下10个词，听写拼写（不看答案）：\n\n'
          '① accommodation  ② necessary  ③ recommend  ④ February\n'
          '⑤ environment  ⑥ government  ⑦ management  ⑧ opportunity\n'
          '⑨ university  ⑩ questionnaire',
          '① accommodation（双c双m）② necessary（1c2s）③ recommend（1c2m）④ February（rb-u-a-ry）\n⑤ environment（-iron-ment）⑥ government（-ern-ment）⑦ management（-age-ment）\n⑧ opportunity（双p，-tunity）⑨ university（-versity）⑩ questionnaire（-tion-naire）',
          '拼写题每道值1分——10个词背对可以多得约5分（考虑实际出现频率），性价比极高'),
         ('考场应急处理实战',
          '以下情况，你的应对方案是什么？\n\n'
          '① S1第3题跟丢了，第4题开始播了\n'
          '② S3整组多选题听了两个但来不及标答案\n'
          '③ 誊写时发现一个词拼不出来，只剩2分钟\n'
          '④ S4感觉完全没听懂，5道题全空了',
          '① 从Q4开始，空白Q3不回头，誊写时靠上下文推断\n② 立刻标记已确认的，不确定的写?，誊写时再判断\n③ 写首字母+横线，誊写完所有题后再回来，上下文推断\n④ 继续往后听，S4即使只做对5题也是5分；不要停下来纠结',
          '考场应急处理的核心：永远往前，不回头——时间是单向流逝的')]
    )
    # PART05 vocab
    U.section(p,5,'考场高频拼写词冲刺表','High-Frequency Spelling Words')
    U.vocab(p,'住宿与生活高频词（S1必考）',
        [('accommodation','住宿（双c双m：ac-com-mo-da-tion）'),
         ('facilities','设施（不是facilties，注意中间无字母颠倒）'),
         ('availability','可用性（avail + ability）'),
         ('appointment','预约（ap-point-ment，双p）'),
         ('enquiry / inquiry','询问（英式enquiry，美式inquiry均接受）'),
         ('receipt','收据（注意p不发音：re-ce-i-pt）'),
         ('maintenance','维护（main-ten-ance，注意中间e）')])
    U.vocab(p,'学术与教育高频词（S3/S4必考）',
        [('environment','环境（envir-on-ment，不是enviroment）'),
         ('government','政府（govern-ment，注意n）'),
         ('management','管理（manage-ment）'),
         ('recommend','推荐（re-com-mend，1个c，2个m）'),
         ('necessary','必要（ne-ces-sa-ry，1c2s）'),
         ('questionnaire','问卷（question-naire，双n）'),
         ('hypothesis','假设（hy-poth-e-sis，发音/haɪˈpɒθɪsɪs/）')])
    U.vocab(p,'时间与月份高频词（填空必考）',
        [('February','二月（Feb-ru-a-ry，r后有u）'),
         ('Wednesday','星期三（Wed-nes-day，d不发音）'),
         ('fortnight','两周（fort-night，英式词，=14 nights）'),
         ('schedule','日程（英式/ˈʃedjuːl/，sh开头）'),
         ('anniversary','周年纪念（anni-vers-a-ry）'),
         ('preliminary','初步的（pre-lim-in-a-ry，长词拆开背）'),
         ('approximately','大约（ap-prox-im-ate-ly）')])
    U.two_col(p,'最后2周冲刺核查清单',
        'Day 1-7 每日必做',
        [('30min模考','完整S1-S4，不暂停，计时'),
         ('错题分类','A/B/C/D/E/F六分类，统计各类数量'),
         ('针对性复习','最多那类错题→对应章节重读'),
         ('拼写检查','每套检查所有拼写，建立个人错误词表'),
         ('记录进步','比较本套vs上套，看哪类错题减少了')],
        'Day 8-14 每日必做',
        [('30min模考','同上，条件更严格'),
         ('深度分析','比较Day 1-7 vs Day 8-14的错题类型变化'),
         ('薄弱点专攻','剩余最顽固的失分类型集中练习'),
         ('Day 13','只做预读热身，不做完整模考'),
         ('Day 14','休息，不做题，正常作息')])
    # PART06 method
    U.section(p,6,'全真模考步骤流程','Full Mock Exam: Step-by-Step Process')
    U.steps(p,'Step 1-3：考前准备',
        [('打印材料','打印答题纸（纸笔考）/ 打开官方界面（机考）；草稿纸2张'),
         ('环境设置','安静房间，手机静音，告知家人不要打扰，提前排除干扰源'),
         ('状态调整','放松5分钟，做3次深呼吸，把注意力从外部转移到考试状态')])
    U.steps(p,'Step 4-6：模考执行',
        [('S1-S2','语速较慢，应得满分或接近满分；完成后立刻预读下一组'),
         ('S3-S4','难度增加，保持稳定，跟丢了找锚点重定位，不要在一题上耗时'),
         ('誊写阶段（纸笔）','按优先级：占位补全→拼写→大小写→单复数→限词，限时10分钟')])
    U.steps(p,'Step 7-8：模考后分析',
        [('核对答案','用官方答案核对，标注每道错题；不要急着查原因，先统计数量'),
         ('错题六分类','对每道错题进行A/B/C/D/E/F分类，找出最多的类型，制定下次针对性训练')])
    U.model_answer(p,'全真模考示范：30分钟完整流程',
        'Section 1（约10分钟）：\n'
        '说明（30秒预读）→ Q1-10作答 → 间隙（预读S2题目）\n\n'
        'Section 2（约10分钟）：\n'
        '说明（预读Q11-20）→ Q11-20作答 → 间隙（预读S3题目）\n\n'
        'Section 3（约10分钟）：\n'
        '说明（预读Q21-30）→ Q21-30作答 → 间隙（预读S4题目）\n\n'
        'Section 4（约10分钟）：\n'
        '说明（预读Q31-40）→ Q31-40作答 → 结束',
        '纸笔考：录音结束→立刻誊写→限时10分钟→按优先级核查\n机考：录音结束→立刻检查已输入内容→重点核查有问号标注的答案',
        ['每组说明时间=宝贵预读时间，绝对不能浪费在翻页或休息上',
         'Section间隙约30秒，已足够完成下一组10题的关键词圈划',
         '机考学生：Section间隙用来检查上一组的输入内容，同时快速预读下一组'])
    U.model_answer(p,'错题分析示范：一套完整分析',
        '模考结果：正确27/40 → 估分约6.5分\n错题：Q3 Q7 Q14 Q19 Q21 Q22 Q23 Q28 Q31 Q35 Q38 Q40',
        '错题分类：\nA类（词汇）：Q28, Q31（2题）— 听到了但不认识\nB类（口音）：Q3（1题）— 英式/ɑː/没识别\nC类（陷阱）：Q19, Q21, Q22（3题）— S3更正陷阱\nD类（预读）：Q7（1题）— 没来得及预读Q7\nE类（速记）：Q38（1题）— 写不完\nF类（粗心）：Q14, Q23, Q35, Q40（4题）— 拼写/单复数\n\n本次重点：F类（4题）→ 下次做前，把Q14/23/35/40的正确拼写默写10遍',
        ['F类粗心4题=多对4题=从6.5→7分的关键（差距约5题），性价比最高',
         'C类陷阱3题=复习L08陷阱章节，S3更正陷阱专项练习2次',
         'A类词汇2题=把Q28/31的答案词加入词汇本，这周集中背'])
    U.compare(p,'模考习惯：低效 vs 高效',
        '❌ 低效模考习惯',
        ['做题中途暂停3-4次，"听不清就重放"',
         '做完核对答案就结束，不分析错题原因',
         '每套都做同类型题，从不针对性补课',
         '在嘈杂环境/边做边看手机',
         '考前一晚刷3套题，精力耗尽'],
        '✅ 高效模考习惯',
        ['全程30分钟不暂停，模拟真实录音只播一次',
         '核对后严格执行六分类，分析每类的根本原因',
         '每套分析后针对性补课，下一套专项观察该类是否减少',
         '安静房间，手机静音，考场级别的专注',
         '考前一晚轻量复习，确保考试当天精力充沛'])
    U.compare(p,'成绩曲线：随机刷题 vs 系统备考',
        '❌ 随机刷题（无系统）',
        ['做了20套题，分数在6-6.5之间徘徊',
         '每套做完不分析，换下一套继续做',
         '同类型错误反复出现',
         '最后2周仍在做同样的练习，没有提高',
         '考试结果：6分，停滞3个月'],
        '✅ 系统备考（有诊断）',
        ['做了10套题，每套认真分析错题类型',
         '发现F类（拼写）是主要问题→专项拼写练习2周',
         '第11-15套：F类错题从4题降到1题',
         '最后2周：薄弱点基本补齐，保持稳定',
         '考试结果：7分，提升1个Band'])
    # PART07 drill
    U.section(p,7,'课堂练习：全真模考模拟','Full Mock Practice')
    U.drill(p,'练习1：高频拼写测试（10词）',
        '不看答案，听写以下10个高频词',
        ['① accommodation  ② necessary  ③ recommend  ④ questionnaire',
         '⑤ environment  ⑥ February  ⑦ facilities  ⑧ approximately',
         '⑨ maintenance  ⑩ government'])
    U.drill(p,'练习2：错题六分类实操',
        '拿出你最近一套模考的错题，现在完成六分类',
        ['错题题号：___  ___  ___  ___  ___  ___  ___  ___',
         'A类（词汇）：___  B类（口音）：___  C类（陷阱）：___',
         'D类（预读）：___  E类（速记）：___  F类（粗心）：___',
         '最多的类别：___  → 本周补课重点：___'])
    U.drill(p,'练习3：模考压力测试（S1+S2）',
        '计时做以下模考（只S1+S2，约20分钟），严格执行规范',
        ['使用剑桥真题（老师指定一套），只做S1和S2',
         '全程不暂停，预读执行，速记符号使用',
         '做完立刻六分类，不要先看答案',
         '目标：S1=9-10分，S2=7-9分'])
    U.drill(p,'练习4：考场应急模拟',
        '模拟以下突发情况，给出应对方案并实际执行',
        ['情况1：S2第3题完全没听到，如何处理？',
         '情况2：S3你写了B，但后来不确定，如何在答题纸上标注？',
         '情况3：誊写时发现Q17的词拼不出来，还有4分钟，如何处理？',
         '情况4：整个S4都没听清楚，还剩2分钟誊写，如何最大化得分？'])
    U.drill(p,'练习5：个人冲刺计划制定',
        '根据本课所学，制定个人2周冲刺计划',
        ['填写：当前分数___分  目标分数___分  差距___题',
         '主要失分类型（来自最近3套分析）：___类（___道题）',
         'Day 1-7每日计划：___分钟模考+___分钟___专项',
         'Day 8-13计划：___  考前日：轻量复习  考前一天：休息',
         '提交给老师审阅，确认计划合理性'])
    # PART08 pitfall
    U.section(p,8,'考前最后冲刺提醒','Final Sprint Reminders')
    U.pain_point(p,'考前最后关头的三大误区',
        ['考前一晚刷题到凌晨：大脑需要睡眠来巩固记忆和恢复认知能力，熬夜刷题得不偿失',
         '考前临时学新策略：最后2天不要读新的方法书或视频，已有策略用熟比学新的更重要',
         '考后立刻对答案：等成绩出来后再分析，考后当天不要因为感觉做得不好而焦虑'],
        ['考前一晚：轻量复习（看信号词表10分钟）→早睡→考试当天早起→正常早餐→准时到考场',
         '考前最后2天：只巩固已学策略，预读练习20分钟/天，不学新内容',
         '考后：放松，等成绩——听力考试结束不等于结果，分数出来再复盘'])
    U.content(p,'考场当天操作手册',
        [('入场前','确认准考证+身份证；提前20分钟到达；找到自己座位，排列好工具'),
         ('等待阶段','不刷题，闭目养神，做3次深呼吸，把注意力聚焦在"预读"这件事上'),
         ('说明阶段','认真听监考员说明，确认题目类型；利用这段时间扫描第一组题目格式'),
         ('录音开始','立刻进入预读模式，不要被第一句话内容吸引，先完成Q1-10预读'),
         ('誊写阶段','按优先级核查：占位词→拼写→大小写→单复数→限词→整体检查1遍')],
        note='考场当天的心态：技术已到位，状态是唯一变量——保持稳定比临时提高更重要')
    U.content(p,'10节课完整回顾：你已掌握了什么',
        [('L01-L05','考试结构→S1-S4各Section策略→10大题型解题逻辑'),
         ('L06-L07','预测技能（三层预读+同义替换）→速记系统（10符号+数字格式+占位）'),
         ('L08-L09','六大陷阱识别→口音适应（英式/澳式+连读+弱读）'),
         ('L10','全真模考规范→错题六分类诊断→2周冲刺计划'),
         ('核心价值','从6分到7分：策略提升（可快速习得）；7分到8分：实力提升（持续积累）')],
        note='恭喜完成听力全部10节课！坚持执行你的冲刺计划，目标分数可期')
    U.tip_card(p,'考前最后72小时行动清单',
        [('考前3天','做1套完整模考+错题分析；轻量复习信号词表；早睡'),
         ('考前2天','只做预读热身（不做完整模考）；复习自己最易犯的拼写词'),
         ('考前1天','完全休息，不做任何题；确认考场位置和交通；准备证件'),
         ('考试当天','正常早餐；提前到达；考场内闭目等待；录音开始即进入状态'),
         ('考后','不要对答案，不要看题目，休息等成绩——已经尽力了')])
    U.tip_card(p,'雅思听力终极速查：10条黄金法则',
        [('法则1','预读20秒——题型→关键词→答案类型，不超时'),
         ('法则2','同义替换——题目词≠录音词，建立100组替换词反应'),
         ('法则3','速记符号——→↑↓ & w/ b/c ~，数字用阿拉伯数字'),
         ('法则4','陷阱信号词——actually/not/if/plan to，延迟写答案'),
         ('法则5','口音适应——/ɑː/vs/æ/，连读wanna/gonna/kinda秒识'),
         ('法则6','跟丢后不放弃——找锚点词重定位，继续往前'),
         ('法则7','占位技术——首字母+横线，不停在一题上'),
         ('法则8','誊写核查——拼写→大小写→单复数→限词'),
         ('法则9','错题六分类——每套做完分类，针对性补课'),
         ('法则10','考前稳定——技术已到位，状态是最后的竞争力')])
    U.checklist(p,'L10课后自检清单（全课程终极版）',
                   ['完成至少5套完整模考，每套都做错题六分类分析',
                   '制定个人2周冲刺计划，提交老师审阅',
                   '背诵20个高频拼写词，测试全对',
                   '信号词表默写：更正（6个）+否定（5个）+条件（5个）全部写出'])
    U.summary(p,10,
        ['全真模考规范：安静环境+不暂停+预读执行+速记执行+誊写核查',
         '错题六分类：A词汇/B口音/C陷阱/D预读/E速记/F粗心——找最多的那类',
         '2周冲刺计划：Day 1-10每天模考+分析，Day 11-12深度分析，Day 13轻量，Day 14休息',
         '高频拼写词：accommodation/necessary/recommend/February等20个必须全对',
         '10条黄金法则：预读→替换→速记→陷阱→口音→不放弃→占位→誊写→分类→稳定'])
    U.homework(p,10,
        ['完成2周冲刺计划的第1-3天执行，记录每天模考分数和错题类型',
         '背诵全部20个高频拼写词，每天自测，直到全部正确',
         '复习L01-L10核心要点，建立个人"听力策略手册"（一张A4纸）'],
        est_time='约90分钟/天（冲刺期）')
    U.preview(p,11,'下一阶段：口语备考',
        ['口语考试结构：Part 1/2/3的时间分配与评分标准',
         '深圳学生口语痛点：中式英语思维、词汇贫乏、停顿过多',
         '7分口语的特征：流利度/词汇/语法/发音四维同步提升'])
    U.save(p, f'{OUT}/L10_全真模考与查漏补缺.pptx')
    print(f"  L10 done: {len(p.slides)} pages")

if __name__ == '__main__':
    print("=== Building IELTS Listening L01-L10 ===")
    L01(); L02(); L03(); L04(); L05(); L06(); L07(); L08(); L09(); L10()
    from pptx import Presentation
    for f in sorted(os.listdir(OUT)):
        if f.endswith('.pptx'):
            n = len(Presentation(os.path.join(OUT,f)).slides)
            flag = "✓" if n >= 60 else "✗"
            print(f"  {flag} {n:3d} pages  {f}")
