# -*- coding: utf-8 -*-
"""
IELTS 授课 PPT 模板库 (v2)
为「每节课一个 60+ 页详细课件」设计，提供丰富的版式构件。
每个单项可调用 set_theme() 设置专属配色。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── 全局尺寸 16:9 ──────────────────────────────────────────────
W, H = Inches(13.333), Inches(7.5)

# ── 通用色板 ───────────────────────────────────────────────────
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_INK   = RGBColor(0x1E, 0x29, 0x3B)   # 正文深色
C_GRAY  = RGBColor(0x55, 0x60, 0x75)   # 次要灰
C_LGRAY = RGBColor(0x8A, 0x93, 0xA6)
C_RED   = RGBColor(0xC0, 0x2A, 0x2A)
C_GREEN = RGBColor(0x1B, 0x7A, 0x3D)
C_AMBER = RGBColor(0xE8, 0x8C, 0x00)

# ── 主题（可被 set_theme 覆盖）─────────────────────────────────
THEME = {
    "primary":  RGBColor(0x1B, 0x3A, 0x6B),  # 主色（深）
    "accent":   RGBColor(0xE8, 0x8C, 0x00),  # 强调色
    "light":    RGBColor(0xEE, 0xF3, 0xFB),  # 浅底
    "softdark": RGBColor(0x0E, 0x24, 0x49),  # 更深的面板
    "name_cn":  "雅思",
    "name_en":  "IELTS",
}

def set_theme(primary, accent, light, softdark, name_cn, name_en):
    THEME["primary"]  = primary
    THEME["accent"]   = accent
    THEME["light"]    = light
    THEME["softdark"] = softdark
    THEME["name_cn"]  = name_cn
    THEME["name_en"]  = name_en


# ── 基础工具 ───────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def slide(prs, bg=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if bg is not None:
        rect(s, 0, 0, 13.333, 7.5, bg)
    return s

def rect(s, l, t, w, h, color, line=None):
    sp = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def rrect(s, l, t, w, h, color, line=None):
    sp = s.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))  # rounded
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    return sp

def text(s, txt, l, t, w, h, size=18, bold=False, color=C_INK,
         align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb

def _bullets(s, items, l, t, w, h, size=16, gap=8, color=C_INK,
             marker="•", head_color=None, head_size=None):
    """items: 字符串 或 (标题, 说明) 元组。"""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    head_color = head_color or THEME["primary"]
    head_size  = head_size or (size + 2)
    for it in items:
        if isinstance(it, tuple):
            head, body = it
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            if not first: p.space_before = Pt(gap)
            first = False
            r = p.add_run(); r.text = f"{marker} {head}"
            r.font.size = Pt(head_size); r.font.bold = True; r.font.color.rgb = head_color
            p2 = tf.add_paragraph(); p2.space_before = Pt(1)
            r2 = p2.add_run(); r2.text = f"   {body}"
            r2.font.size = Pt(size - 1); r2.font.color.rgb = C_GRAY
        else:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            if not first: p.space_before = Pt(gap)
            first = False
            r = p.add_run(); r.text = f"{marker} {it}"
            r.font.size = Pt(size); r.font.color.rgb = color
    return tb


# ── 页眉（内容页统一样式）─────────────────────────────────────
def _header(s, title, kicker=None):
    rect(s, 0, 0, 13.333, 1.0, THEME["primary"])
    rect(s, 0, 1.0, 13.333, 0.08, THEME["accent"])
    if kicker:
        text(s, kicker, 0.45, 0.12, 11.0, 0.3, size=11, bold=True,
             color=THEME["accent"])
        text(s, title, 0.45, 0.38, 12.4, 0.58, size=21, bold=True, color=C_WHITE)
    else:
        text(s, title, 0.45, 0.18, 12.4, 0.66, size=22, bold=True,
             color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # 页脚
    text(s, f"{THEME['name_en']} · 深圳提分课程", 0.45, 7.12, 6.0, 0.3,
         size=9, color=C_LGRAY)


# ════════════════════════════════════════════════════════════════
#  封面 / 分隔 / 结构页
# ════════════════════════════════════════════════════════════════
def cover(prs, skill_cn, lesson_no, lesson_title, subtitle, total=10):
    s = slide(prs, THEME["primary"])
    rect(s, 0, 0, 0.4, 7.5, THEME["accent"])
    # 角标
    rect(s, 9.6, 0, 3.733, 7.5, THEME["softdark"])
    text(s, f"L{lesson_no:02d}", 9.6, 1.6, 3.733, 2.0, size=120, bold=True,
         color=THEME["accent"], align=PP_ALIGN.CENTER)
    text(s, f"/ {total:02d}", 9.6, 3.9, 3.733, 0.8, size=30, bold=True,
         color=C_WHITE, align=PP_ALIGN.CENTER)
    # 主标题
    text(s, f"{THEME['name_en']} {skill_cn}", 0.85, 1.4, 8.4, 0.7, size=20,
         bold=True, color=THEME["accent"])
    text(s, lesson_title, 0.85, 2.1, 8.4, 2.4, size=42, bold=True, color=C_WHITE)
    text(s, subtitle, 0.85, 4.6, 8.2, 1.2, size=18,
         color=RGBColor(0xCF, 0xDC, 0xF2))
    # 底部标签
    rrect(s, 0.85, 6.2, 4.2, 0.6, THEME["accent"])
    text(s, f"第 {lesson_no} 节课 · 课时 90 分钟", 0.95, 6.28, 4.0, 0.45,
         size=13, bold=True, color=C_WHITE)
    return s

def section(prs, idx, title_cn, title_en):
    s = slide(prs, THEME["softdark"])
    rect(s, 0, 3.0, 13.333, 0.05, THEME["accent"])
    text(s, f"PART {idx:02d}", 1.0, 2.0, 11.0, 0.6, size=22, bold=True,
         color=THEME["accent"])
    text(s, title_cn, 1.0, 3.2, 11.3, 1.4, size=40, bold=True, color=C_WHITE)
    text(s, title_en, 1.0, 4.7, 11.3, 0.8, size=20, italic=True,
         color=RGBColor(0xB9, 0xC8, 0xE0))
    return s

def agenda(prs, items):
    s = slide(prs, C_WHITE); _header(s, "本课流程  Lesson Roadmap")
    half = (len(items) + 1) // 2
    cols = [items[:half], items[half:]]
    for ci, col in enumerate(cols):
        x = 0.6 + ci * 6.4
        for i, it in enumerate(col):
            gi = ci * half + i
            y = 1.45 + i * 0.92
            rrect(s, x, y, 0.62, 0.62, THEME["primary"])
            text(s, str(gi + 1), x, y + 0.05, 0.62, 0.5, size=20, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)
            text(s, it, x + 0.8, y + 0.02, 5.4, 0.8, size=14.5, bold=True,
                 color=C_INK, anchor=MSO_ANCHOR.MIDDLE)
    return s

def objectives(prs, items):
    s = slide(prs, C_WHITE); _header(s, "学习目标  Learning Objectives")
    text(s, "完成本课后，你将能够：", 0.6, 1.3, 12.0, 0.5, size=16, bold=True,
         color=THEME["accent"])
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12.0), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i: p.space_before = Pt(14)
        r = p.add_run(); r.text = f"✦  {it}"
        r.font.size = Pt(18); r.font.color.rgb = C_INK
    return s

def review(prs, points):
    s = slide(prs, C_WHITE); _header(s, "上节回顾  Quick Review", kicker="WARM-UP")
    rrect(s, 0.6, 1.3, 12.1, 5.4, THEME["light"])
    _bullets(s, points, 1.0, 1.7, 11.3, 4.8, size=16.5, gap=12, marker="↺")
    return s


# ════════════════════════════════════════════════════════════════
#  内容版式
# ════════════════════════════════════════════════════════════════
def concept(prs, title, lead, points, note=None, kicker="核心讲解"):
    """带导语的讲解页。"""
    s = slide(prs, C_WHITE); _header(s, title, kicker=kicker)
    rrect(s, 0.6, 1.25, 12.1, 0.95, THEME["light"])
    text(s, lead, 0.85, 1.34, 11.6, 0.8, size=15, italic=True, color=THEME["primary"],
         anchor=MSO_ANCHOR.MIDDLE)
    body_h = 4.3 if note else 4.7
    _bullets(s, points, 0.7, 2.45, 12.0, body_h, size=16, gap=9)
    if note:
        rect(s, 0.6, 6.55, 12.1, 0.6, RGBColor(0xFD, 0xF0, 0xE0))
        text(s, f"⚠ 深圳学生注意：{note}", 0.8, 6.6, 11.7, 0.5, size=12.5,
             color=C_RED, italic=True, anchor=MSO_ANCHOR.MIDDLE)
    return s

def content(prs, title, points, note=None, kicker=None):
    s = slide(prs, C_WHITE); _header(s, title, kicker=kicker)
    body_h = 5.0 if note else 5.6
    _bullets(s, points, 0.7, 1.35, 12.0, body_h, size=16.5, gap=10)
    if note:
        rect(s, 0.6, 6.55, 12.1, 0.6, RGBColor(0xFD, 0xF0, 0xE0))
        text(s, f"⚠ 深圳学生注意：{note}", 0.8, 6.6, 11.7, 0.5, size=12.5,
             color=C_RED, italic=True, anchor=MSO_ANCHOR.MIDDLE)
    return s

def two_col(prs, title, lt, litems, rt, ritems, lcolor=None, rcolor=None):
    s = slide(prs, C_WHITE); _header(s, title)
    lcolor = lcolor or THEME["primary"]; rcolor = rcolor or THEME["accent"]
    for ci, (ct, items, cc) in enumerate([(lt, litems, lcolor), (rt, ritems, rcolor)]):
        x = 0.6 + ci * 6.35
        rect(s, x, 1.3, 6.0, 0.5, cc)
        text(s, ct, x + 0.15, 1.34, 5.7, 0.42, size=14.5, bold=True,
             color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        _bullets(s, items, x + 0.1, 1.95, 5.85, 4.9, size=14.5, gap=7)
    return s

def vocab(prs, title, rows, headers=("表达 / Expression", "用法 / 中文")):
    """词汇/句型表：rows = [(term, meaning), ...]，最多约9行。"""
    s = slide(prs, C_WHITE); _header(s, title, kicker="词汇 · 句型库")
    top = 1.35; lh = (6.6 - top) / max(len(rows) + 1, 6)
    lh = min(lh, 0.62)
    # 表头
    rect(s, 0.6, top, 5.4, lh, THEME["primary"])
    rect(s, 6.0, top, 6.7, lh, THEME["primary"])
    text(s, headers[0], 0.75, top, 5.2, lh, size=13, bold=True, color=C_WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, headers[1], 6.15, top, 6.4, lh, size=13, bold=True, color=C_WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
    for i, (term, mean) in enumerate(rows):
        y = top + lh * (i + 1)
        bg = THEME["light"] if i % 2 == 0 else C_WHITE
        rect(s, 0.6, y, 5.4, lh, bg); rect(s, 6.0, y, 6.7, lh, bg)
        text(s, term, 0.75, y, 5.2, lh, size=12.5, bold=True, color=THEME["primary"],
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, mean, 6.15, y, 6.4, lh, size=12.5, color=C_INK,
             anchor=MSO_ANCHOR.MIDDLE)
    return s

def steps(prs, title, step_list, kicker="方法步骤"):
    """step_list = [(标题, 说明), ...] 纵向编号流程。"""
    s = slide(prs, C_WHITE); _header(s, title, kicker=kicker)
    top = 1.4; n = len(step_list); gap = min((6.6 - top) / n, 1.15)
    for i, (h, b) in enumerate(step_list):
        y = top + i * gap
        rrect(s, 0.6, y, 0.7, 0.7, THEME["accent"])
        text(s, f"{i+1}", 0.6, y + 0.06, 0.7, 0.55, size=22, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < n - 1:
            rect(s, 0.92, y + 0.72, 0.05, gap - 0.72, THEME["light"])
        text(s, h, 1.55, y, 11.0, 0.45, size=16, bold=True, color=THEME["primary"])
        text(s, b, 1.55, y + 0.4, 11.0, gap - 0.4, size=13.5, color=C_GRAY)
    return s

def example(prs, title, prompt, answer=None, tip=None, kicker="例题精讲"):
    s = slide(prs, C_WHITE); _header(s, title, kicker=kicker)
    rect(s, 0.6, 1.3, 12.1, 0.4, THEME["primary"])
    text(s, "【题目 / Question】", 0.75, 1.32, 6.0, 0.36, size=12.5, bold=True,
         color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    qh = 2.3 if (answer or tip) else 5.0
    tb = s.shapes.add_textbox(Inches(0.65), Inches(1.8), Inches(12.0), Inches(qh))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = prompt
    r.font.size = Pt(14.5); r.font.color.rgb = C_INK
    y = 1.8 + qh
    if answer:
        rect(s, 0.6, y, 12.1, 0.4, C_GREEN)
        text(s, "【参考答案 / Answer】", 0.75, y + 0.02, 6.0, 0.36, size=12.5,
             bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        ah = 1.6 if tip else 2.4
        tb2 = s.shapes.add_textbox(Inches(0.65), Inches(y + 0.5), Inches(12.0), Inches(ah))
        tf2 = tb2.text_frame; tf2.word_wrap = True
        p2 = tf2.paragraphs[0]; r2 = p2.add_run(); r2.text = answer
        r2.font.size = Pt(14); r2.font.color.rgb = C_GREEN
        y = y + 0.5 + ah
    if tip:
        rect(s, 0.6, y, 12.1, 0.4, THEME["accent"])
        text(s, "💡 技巧点拨", 0.75, y + 0.02, 6.0, 0.36, size=12.5, bold=True,
             color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        tb3 = s.shapes.add_textbox(Inches(0.65), Inches(y + 0.5), Inches(12.0), Inches(1.4))
        tf3 = tb3.text_frame; tf3.word_wrap = True
        p3 = tf3.paragraphs[0]; r3 = p3.add_run(); r3.text = tip
        r3.font.size = Pt(13.5); r3.font.color.rgb = C_GRAY
    return s

def model_answer(prs, title, prompt, body, annotations, kicker="高分范例"):
    """范文 + 右侧批注。"""
    s = slide(prs, C_WHITE); _header(s, title, kicker=kicker)
    if prompt:
        rect(s, 0.6, 1.3, 12.1, 0.42, THEME["softdark"])
        text(s, f"题目：{prompt}", 0.78, 1.32, 11.8, 0.38, size=12, bold=True,
             color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
        ytop = 1.95
    else:
        ytop = 1.4
    rrect(s, 0.6, ytop, 7.7, 6.5 - ytop, THEME["light"])
    tb = s.shapes.add_textbox(Inches(0.85), Inches(ytop + 0.15),
                              Inches(7.25), Inches(6.3 - ytop))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = body
    r.font.size = Pt(13); r.font.color.rgb = C_INK
    text(s, "✍ 得分点批注", 8.5, ytop, 4.3, 0.4, size=13, bold=True,
         color=THEME["accent"])
    _bullets(s, annotations, 8.5, ytop + 0.45, 4.3, 6.0 - ytop, size=12, gap=7,
             marker="▸")
    return s

def compare(prs, title, bad_title, bad, good_title, good):
    s = slide(prs, C_WHITE); _header(s, title, kicker="对比示范")
    rect(s, 0.6, 1.3, 6.0, 0.5, C_RED)
    text(s, bad_title, 0.75, 1.33, 5.7, 0.44, size=14, bold=True, color=C_WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 6.75, 1.3, 6.0, 0.5, C_GREEN)
    text(s, good_title, 6.9, 1.33, 5.7, 0.44, size=14, bold=True, color=C_WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
    rrect(s, 0.6, 1.9, 6.0, 4.9, RGBColor(0xFB, 0xEC, 0xEC))
    rrect(s, 6.75, 1.9, 6.0, 4.9, RGBColor(0xEA, 0xF6, 0xEE))
    _bullets(s, bad, 0.85, 2.1, 5.55, 4.5, size=13.5, gap=7, marker="✗",
             color=C_RED, head_color=C_RED)
    _bullets(s, good, 7.0, 2.1, 5.55, 4.5, size=13.5, gap=7, marker="✓",
             color=C_GREEN, head_color=C_GREEN)
    return s

def pain_point(prs, title, problems, solutions):
    s = slide(prs, C_WHITE)
    rect(s, 0, 0, 13.333, 1.0, RGBColor(0x7A, 0x16, 0x16))
    rect(s, 0, 1.0, 13.333, 0.08, THEME["accent"])
    text(s, f"🎯 深圳学生高频痛点 · {title}", 0.45, 0.18, 12.4, 0.66, size=20,
         bold=True, color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    text(s, "❌ 常见误区", 0.6, 1.25, 6.0, 0.45, size=15, bold=True, color=C_RED)
    _bullets(s, problems, 0.6, 1.8, 6.0, 5.0, size=14, gap=8, marker="✗",
             color=C_RED, head_color=C_RED)
    rect(s, 6.68, 1.2, 0.04, 5.7, THEME["light"])
    text(s, "✅ 正确做法", 6.9, 1.25, 6.0, 0.45, size=15, bold=True, color=C_GREEN)
    _bullets(s, solutions, 6.9, 1.8, 5.9, 5.0, size=14, gap=8, marker="✓",
             color=C_GREEN, head_color=C_GREEN)
    return s

def drill(prs, title, instruction, items, kicker="课堂练习"):
    s = slide(prs, C_WHITE); _header(s, title, kicker=kicker)
    rrect(s, 0.6, 1.3, 12.1, 0.7, THEME["primary"])
    text(s, f"📝 {instruction}", 0.85, 1.38, 11.6, 0.55, size=14, bold=True,
         color=C_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _bullets(s, items, 0.7, 2.25, 12.0, 4.6, size=15, gap=11, marker="▢")
    return s

def tip_card(prs, title, cards, kicker="技巧速记"):
    """cards = [(标题, 内容), ...] 卡片网格 (2 列)。"""
    s = slide(prs, C_WHITE); _header(s, title, kicker=kicker)
    n = len(cards); cols = 2; rows = (n + 1) // 2
    cw, ch = 6.0, min((6.5 - 1.3) / rows - 0.2, 1.6)
    for i, (h, b) in enumerate(cards):
        r_i, c_i = divmod(i, cols)
        x = 0.6 + c_i * 6.35; y = 1.35 + r_i * (ch + 0.25)
        rrect(s, x, y, cw, ch, THEME["light"])
        rect(s, x, y, 0.12, ch, THEME["accent"])
        text(s, h, x + 0.3, y + 0.1, cw - 0.5, 0.45, size=14, bold=True,
             color=THEME["primary"])
        text(s, b, x + 0.3, y + 0.55, cw - 0.5, ch - 0.6, size=12.5, color=C_GRAY)
    return s

def checklist(prs, title, items, kicker="自检清单"):
    s = slide(prs, C_WHITE); _header(s, title, kicker=kicker)
    _bullets(s, items, 0.8, 1.5, 11.8, 5.3, size=16, gap=12, marker="☐")
    return s

def summary(prs, lesson_no, key_points, takeaway=None):
    s = slide(prs, THEME["primary"])
    rect(s, 0, 0, 0.4, 7.5, THEME["accent"])
    text(s, f"L{lesson_no:02d} 课堂小结", 0.8, 0.4, 11.0, 0.8, size=30, bold=True,
         color=C_WHITE)
    text(s, "📌 本课核心要点", 0.8, 1.45, 11.0, 0.5, size=16, bold=True,
         color=THEME["accent"])
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.05), Inches(11.7), Inches(3.8))
    tf = tb.text_frame; tf.word_wrap = True
    for i, kp in enumerate(key_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i: p.space_before = Pt(11)
        r = p.add_run(); r.text = f"✓  {kp}"
        r.font.size = Pt(15.5); r.font.color.rgb = C_WHITE
    if takeaway:
        rrect(s, 0.8, 6.0, 11.7, 1.0, THEME["softdark"])
        text(s, f"💎 一句话记住：{takeaway}", 1.0, 6.12, 11.3, 0.8, size=15,
             bold=True, color=THEME["accent"], anchor=MSO_ANCHOR.MIDDLE)
    return s

def homework(prs, lesson_no, tasks, est_time=None):
    s = slide(prs, C_WHITE); _header(s, f"L{lesson_no:02d} 课后作业  Homework",
                                     kicker="练习巩固")
    if est_time:
        text(s, f"⏱ 建议用时：{est_time}", 0.6, 1.25, 11.0, 0.4, size=13,
             bold=True, color=THEME["accent"])
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12.0), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, t in enumerate(tasks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i: p.space_before = Pt(13)
        r = p.add_run(); r.text = f"作业 {i+1}　{t}"
        r.font.size = Pt(15.5); r.font.color.rgb = C_INK; r.font.bold = True
    return s

def preview(prs, next_no, next_title, points):
    s = slide(prs, THEME["softdark"])
    text(s, "下节预告  Next Lesson", 1.0, 1.4, 11.0, 0.6, size=18, bold=True,
         color=THEME["accent"])
    text(s, f"L{next_no:02d} · {next_title}", 1.0, 2.1, 11.3, 1.0, size=30,
         bold=True, color=C_WHITE)
    _bullets(s, points, 1.0, 3.5, 11.3, 3.0, size=16, gap=10, marker="→",
             color=RGBColor(0xCF, 0xDC, 0xF2))
    return s


# ── 便捷：统计并保存 ───────────────────────────────────────────
def save(prs, path, min_pages=60):
    n = len(prs.slides)
    prs.save(path)
    flag = "OK " if n >= min_pages else "!! "
    print(f"[{flag}] {n:>3} pages -> {path}")
    return n
